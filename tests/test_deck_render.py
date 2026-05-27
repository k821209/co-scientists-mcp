"""Deck rendering: placeholder resolution + slide rendering modes.

PPTX export is covered by a smoke test only — requires python-pptx
optional dep.
"""
from __future__ import annotations

import pathlib

import pytest

from co_scientist_local.backends.base import NotFound
from co_scientist_local.tools import deck_render, decks, figures, papers


def _setup(state):
    papers.create_paper(state, title="My paper")
    decks.create_deck(state, "my-paper", title="d", deck_id="d")
    return "my-paper"


# ─── placeholder resolution ──────────────────────────────────────────────


def test_parse_concept_extracts_pairs():
    concept = """
    Palette:
      bg: #fafaf7  surface: #ffffff  text: #1a1a1a  accent: #b58900
    Typography:
      display: Inter Bold
      body: Inter Regular
    """
    kv = deck_render._parse_concept(concept)
    assert kv["accent"] == "#b58900"
    assert kv["bg"] == "#fafaf7"
    assert kv["display"] == "Inter Bold"
    assert kv["body"] == "Inter Regular"


def test_resolve_replaces_known_placeholders():
    concept = "accent: #b58900"
    text = "Use {accent} for the trend line"
    assert deck_render.resolve_placeholders(text, concept) == "Use #b58900 for the trend line"


def test_resolve_leaves_unknown_placeholders_visible():
    """Unknown placeholders stay literal so the bug is obvious."""
    out = deck_render.resolve_placeholders("Hello {unknown}", "accent: red")
    assert "{unknown}" in out


def test_resolve_none_concept_returns_text_unchanged():
    assert deck_render.resolve_placeholders("hi {x}", None) == "hi {x}"


def test_resolve_no_braces_returns_unchanged():
    assert deck_render.resolve_placeholders("no placeholders", "x: y") == "no placeholders"


# ─── render_slide modes ──────────────────────────────────────────────────


def test_render_paper_figure_copies_blob(state, tmp_path):
    slug = _setup(state)
    # Stage a figure with a real PNG blob
    png_bytes = (b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
    p = tmp_path / "fig1.png"
    p.write_bytes(png_bytes)
    figures.add_figure(
        state, slug, figure_number=1, title="Figure 1",
        caption="x", local_path=str(p),
    )
    s = decks.add_slide(
        state, slug, "d", slide_number=1, role="result", title="Fig 1 slide",
        notes="my notes", render_mode="paper-figure", figure_number=1,
    )
    r = deck_render.render_slide(state, slug, "d", s["id"])
    assert r["size_bytes"] == len(png_bytes)
    assert r["mode"] == "paper-figure"
    # Verify slide doc was updated
    updated = decks.list_slides(state, slug, "d")[0]
    assert updated["image_blob_path"]
    assert updated["status"] == "rendered"


def test_render_ai_image_uses_image_gen(state, monkeypatch):
    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #f00")
    s = decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="cover", notes="hi",
        prompt="Use {accent} on the cover",
        render_mode="ai-image",
    )

    class FakeGen:
        captured: dict = {}
        def generate(self, *, prompt, aspect_ratio="1:1", model="gpt-image-2"):
            FakeGen.captured["prompt"] = prompt
            FakeGen.captured["aspect_ratio"] = aspect_ratio
            return b"fake-png-bytes"

    state.image_gen = FakeGen()
    r = deck_render.render_slide(state, slug, "d", s["id"])
    assert r["mode"] == "ai-image"
    assert FakeGen.captured["prompt"] == "Use #f00 on the cover"
    assert FakeGen.captured["aspect_ratio"] == "16:9"


def test_render_code_shape_requires_local_path(state, tmp_path):
    slug = _setup(state)
    s = decks.add_slide(
        state, slug, "d", slide_number=1, role="result", title="x",
        notes="hi", code="plt.savefig('out.png')",
        render_mode="code-shape",
    )
    with pytest.raises(ValueError):
        deck_render.render_slide(state, slug, "d", s["id"])


def test_render_code_shape_with_local_path(state, tmp_path):
    slug = _setup(state)
    s = decks.add_slide(
        state, slug, "d", slide_number=1, role="result", title="x",
        notes="hi", code="plt.savefig('out.png')",
        render_mode="code-shape",
    )
    p = tmp_path / "code-shape.png"
    p.write_bytes(b"local-png-bytes")
    r = deck_render.render_slide(state, slug, "d", s["id"], local_path=str(p))
    assert r["size_bytes"] == len(b"local-png-bytes")


def test_render_deck_marks_status_rendered_when_all_done(state, tmp_path):
    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: red")
    # Two ai-image slides
    decks.add_slide(state, slug, "d", slide_number=1, role="title",
                    title="t1", notes="n", prompt="Hi", render_mode="ai-image")
    decks.add_slide(state, slug, "d", slide_number=2, role="conclusion",
                    title="t2", notes="n", prompt="Bye", render_mode="ai-image")

    class FakeGen:
        def generate(self, *, prompt, aspect_ratio="1:1", model="gpt-image-2"):
            return b"fake-png"
    state.image_gen = FakeGen()

    result = deck_render.render_deck(state, slug, "d")
    assert len(result["rendered"]) == 2
    assert result["errors"] == []
    deck = decks.get_deck(state, slug, "d")
    assert deck["status"] == "rendered"


def test_render_deck_skips_code_shape(state):
    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="title",
                    title="t1", notes="n", code="x",
                    render_mode="code-shape")
    result = deck_render.render_deck(state, slug, "d")
    assert len(result["skipped"]) == 1
    assert result["skipped"][0]["reason"] == "needs local PNG from agent"


def test_render_missing_slide_raises(state):
    slug = _setup(state)
    with pytest.raises(NotFound):
        deck_render.render_slide(state, slug, "d", "ghost-id")


# ─── text slides + theme ─────────────────────────────────────────────────


def test_theme_colors_from_concept():
    """Palette exposes all 7 keys (todo 007 axis 4). Specified keys
    pass through unchanged; unspecified keys get computed defaults
    (muted = blend of fg/bg; secondary / highlight = accent shifts)."""
    concept = "Palette:\n  bg: #101010  text: #fafafa  accent: #ff8800"
    c = deck_render._theme_colors(concept)
    # Specified keys
    assert c["accent"] == "#ff8800"
    assert c["background"] == "#101010"
    assert c["foreground"] == "#fafafa"
    # Computed defaults
    assert c["surface"] == "#101010"   # falls back to background
    assert "muted" in c and c["muted"] != "#fafafa"
    assert "secondary" in c
    assert "highlight" in c


def test_theme_colors_explicit_muted_secondary_highlight():
    """When the concept declares the optional keys, they pass through
    instead of being computed."""
    concept = (
        "Palette:\n  bg: #fff  text: #111  accent: #b58900\n"
        "  muted: #888  secondary: #226699  highlight: #ee5500\n"
        "  surface: #fafafa"
    )
    c = deck_render._theme_colors(concept)
    assert c["muted"] == "#888"
    assert c["secondary"] == "#226699"
    assert c["highlight"] == "#ee5500"
    assert c["surface"] == "#fafafa"


def test_theme_colors_defaults_when_absent():
    c = deck_render._theme_colors(None)
    assert c["accent"] == "#2E7D32"
    assert c["background"] == "#FFFFFF"
    assert c["foreground"] == "#1A1A1A"
    # The 4 new keys all defined too (with computed defaults)
    for k in ("surface", "muted", "secondary", "highlight"):
        assert k in c, f"missing palette key: {k}"


def test_render_slide_text_mode_raises(state):
    slug = _setup(state)
    s = decks.add_slide(state, slug, "d", slide_number=1, role="outline",
                        title="Agenda", body="- a\n- b", render_mode="text")
    with pytest.raises(ValueError, match="text slides"):
        deck_render.render_slide(state, slug, "d", s["id"])


def test_render_deck_skips_text_slides(state):
    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="outline",
                    title="Agenda", body="x", render_mode="text")
    result = deck_render.render_deck(state, slug, "d")
    assert len(result["skipped"]) == 1
    assert result["skipped"][0]["reason"] == \
        "text slide — native PPTX text on export"


def test_render_deck_all_text_marks_rendered(state):
    """A deck of only text slides is fully rendered — text slides never
    get an image_blob_path, so they must not block the status flip."""
    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="title",
                    title="t", body="x", render_mode="text")
    deck_render.render_deck(state, slug, "d")
    assert decks.get_deck(state, slug, "d")["status"] == "rendered"


def test_pdf_via_soffice_missing_returns_none(state, tmp_path, monkeypatch):
    def _no_soffice(*a, **k):
        raise FileNotFoundError("soffice")
    monkeypatch.setattr(deck_render.subprocess, "run", _no_soffice)
    assert deck_render._pdf_via_soffice(tmp_path / "deck.pptx") is None


# ─── PPTX export smoke test (needs the optional python-pptx dep) ─────────

# A valid 1×1 transparent PNG — python-pptx parses the header for sizing.
import base64 as _b64  # noqa: E402
_PNG_1x1 = _b64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR42mNk"
    "YAAAAAYAAjCB0C8AAAAASUVORK5CYII="
)


def test_export_deck_to_pptx_smoke(state, tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900",
                      aspect_ratio="4:3")
    # one image slide (paper-figure, rendered)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    img = decks.add_slide(state, slug, "d", slide_number=1, role="result",
                          title="Result", notes="n",
                          render_mode="paper-figure", figure_number=1)
    deck_render.render_slide(state, slug, "d", img["id"])
    # one native text slide
    decks.add_slide(state, slug, "d", slide_number=2, role="outline",
                    title="Agenda", body="- point one\n- point two",
                    notes="n2", render_mode="text")
    # keep the test offline + deterministic — skip the LibreOffice step
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)

    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert out.is_file()
    assert res["slide_count"] == 2
    assert res["image_slides"] == 1
    assert res["text_slides"] == 1
    assert res["missing_renders"] == []     # the text slide is not "missing"
    assert res["aspect_ratio"] == "4:3"
    assert res["pdf_skipped"] is True
    # PDF was skipped → per-slide PNGs also skipped (their source).
    assert res["slide_pngs"] == []
    assert res["slide_pngs_skipped"] is True


# ─── regions: hybrid multi-image slides ──────────────────────


def _hybrid_setup(state, tmp_path):
    """paper + deck "d" + figure 1 + a hybrid slide with two regions
    (one paper-figure, one ai-image)."""
    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    s = decks.add_slide(state, slug, "d", slide_number=1, role="result", title="R")
    decks.set_slide_regions(state, slug, "d", s["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1,
         "x": 0.05, "y": 0.25, "w": 0.42, "h": 0.6},
        {"render_mode": "ai-image", "prompt": "{accent} schematic",
         "x": 0.52, "y": 0.25, "w": 0.42, "h": 0.6},
    ])
    return slug, s["id"]


def test_nearest_aspect():
    assert deck_render._nearest_aspect(1920, 1080) == "16:9"
    assert deck_render._nearest_aspect(100, 100) == "1:1"
    assert deck_render._nearest_aspect(1080, 1920) == "9:16"
    assert deck_render._nearest_aspect(400, 300) == "4:3"


def test_render_region_paper_figure(state, tmp_path):
    slug, sid = _hybrid_setup(state, tmp_path)
    r = deck_render.render_region(state, slug, "d", sid, "r1")
    assert r["region_id"] == "r1" and r["mode"] == "paper-figure"
    slide = decks.list_slides(state, slug, "d")[0]
    assert slide["regions"][0]["image_blob_path"]


def test_render_region_ai_image_matches_box_aspect(state, tmp_path):
    slug, sid = _hybrid_setup(state, tmp_path)
    captured: dict = {}

    class FakeGen:
        def generate(self, *, prompt, aspect_ratio="1:1", model="gpt-image-2"):
            captured["prompt"] = prompt
            captured["aspect"] = aspect_ratio
            return b"ai-png"

    state.image_gen = FakeGen()
    decks.update_deck(state, slug, "d", concept="accent: #abc")
    r = deck_render.render_region(state, slug, "d", sid, "r2")
    assert r["mode"] == "ai-image"
    assert captured["prompt"] == "#abc schematic"   # placeholder resolved
    # region box 0.42w × 0.6h on a 16:9 deck is portrait-ish → "4:3",
    # NOT the deck's 16:9 (the original's todo 052).
    assert captured["aspect"] == "4:3"


def test_render_region_code_shape_needs_local_path(state, tmp_path):
    slug = _setup(state)
    s = decks.add_slide(state, slug, "d", slide_number=1, role="result", title="R")
    decks.set_slide_regions(state, slug, "d", s["id"], regions=[
        {"render_mode": "code-shape", "code": "plt...",
         "x": 0.1, "y": 0.2, "w": 0.8, "h": 0.6},
    ])
    with pytest.raises(ValueError, match="local_path"):
        deck_render.render_region(state, slug, "d", s["id"], "r1")
    p = tmp_path / "r.png"
    p.write_bytes(b"region-png")
    out = deck_render.render_region(state, slug, "d", s["id"], "r1",
                                   local_path=str(p))
    assert out["size_bytes"] == len(b"region-png")


def test_render_region_missing_raises(state, tmp_path):
    slug, sid = _hybrid_setup(state, tmp_path)
    with pytest.raises(NotFound):
        deck_render.render_region(state, slug, "d", sid, "r99")


def test_render_slide_hybrid_renders_all_regions(state, tmp_path):
    slug, sid = _hybrid_setup(state, tmp_path)

    class FakeGen:
        def generate(self, *, prompt, aspect_ratio="1:1", model="gpt-image-2"):
            return b"ai-png"

    state.image_gen = FakeGen()
    res = deck_render.render_slide(state, slug, "d", sid)
    assert res["mode"] == "hybrid"
    assert len(res["rendered"]) == 2
    assert res["skipped"] == [] and res["errors"] == []
    slide = decks.list_slides(state, slug, "d")[0]
    assert all(r["image_blob_path"] for r in slide["regions"])


def test_render_deck_hybrid_code_shape_region_skipped(state, tmp_path):
    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    s = decks.add_slide(state, slug, "d", slide_number=1, role="result", title="R")
    decks.set_slide_regions(state, slug, "d", s["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1,
         "x": 0.05, "y": 0.2, "w": 0.4, "h": 0.6},
        {"render_mode": "code-shape", "code": "x",
         "x": 0.5, "y": 0.2, "w": 0.4, "h": 0.6},
    ])
    result = deck_render.render_deck(state, slug, "d")
    assert len(result["rendered"]) == 1                 # paper-figure region
    assert len(result["skipped"]) == 1                  # code-shape region
    assert "code-shape region" in result["skipped"][0]["reason"]
    # one region still unrendered → deck not flipped to 'rendered'
    assert decks.get_deck(state, slug, "d")["status"] != "rendered"


def test_export_hybrid_slide_smoke(state, tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug, sid = _hybrid_setup(state, tmp_path)

    class FakeGen:
        def generate(self, *, prompt, aspect_ratio="1:1", model="gpt-image-2"):
            return _PNG_1x1

    state.image_gen = FakeGen()
    deck_render.render_slide(state, slug, "d", sid)      # render both regions
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)

    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert out.is_file()
    assert res["hybrid_slides"] == 1
    assert res["missing_renders"] == []
    prs = Presentation(str(out))
    pics = [sh for sh in prs.slides[0].shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 2   # two region images placed as separate shapes


def test_export_hybrid_slide_with_body_renders_native_bullets(
        state, tmp_path, monkeypatch):
    """A hybrid slide with `body` renders the bullets as NATIVE editable
    text in the LEFT half, alongside an image region on the right — the
    'title + bullets + figure' layout."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    s = decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Pipeline vs MCP",
        body="- two-week bespoke pipeline\n- 30-second MCP query",
        notes="n",
    )
    decks.set_slide_regions(state, slug, "d", s["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1,
         "x": 0.54, "y": 0.22, "w": 0.42, "h": 0.65},
    ])
    deck_render.render_slide(state, slug, "d", s["id"])
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)

    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    prs = Presentation(str(out))
    slide = prs.slides[0]

    # one image region on the right
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1

    # native body text present and on the LEFT half of the slide. Each
    # bullet is its own textbox (block-based renderer — todo 002), so
    # collect them all and assert (a) both bullets exist as text frames
    # and (b) every body textbox sits in the left half.
    body_tbs = [
        sh for sh in slide.shapes
        if sh.has_text_frame and (
            "two-week bespoke pipeline" in sh.text_frame.text
            or "30-second MCP query" in sh.text_frame.text
        )
        and sh.left is not None
        and sh.left + sh.width <= prs.slide_width / 2 + 1
    ]
    body_text = "".join(b.text_frame.text for b in body_tbs)
    assert "two-week bespoke pipeline" in body_text
    assert "30-second MCP query" in body_text


# ─── placeholder metadata: fit mode + captured image size ────


def test_png_size_parses_dimensions():
    assert deck_render._png_size(_PNG_1x1) == (1, 1)
    assert deck_render._png_size(b"not a png at all") is None
    assert deck_render._image_dims(_PNG_1x1) == {
        "image_width": 1, "image_height": 1,
    }
    assert deck_render._image_dims(b"junk") == {}


def test_render_region_records_image_dimensions(state, tmp_path):
    slug, sid = _hybrid_setup(state, tmp_path)
    deck_render.render_region(state, slug, "d", sid, "r1")
    region = decks.list_slides(state, slug, "d")[0]["regions"][0]
    assert region["image_width"] == 1 and region["image_height"] == 1


def test_render_slide_records_image_dimensions(state, tmp_path):
    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    s = decks.add_slide(state, slug, "d", slide_number=1, role="result",
                        title="R", render_mode="paper-figure", figure_number=1)
    deck_render.render_slide(state, slug, "d", s["id"])
    slide = decks.list_slides(state, slug, "d")[0]
    assert slide["image_width"] == 1 and slide["image_height"] == 1


def test_export_cover_region_crops(state, tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    s = decks.add_slide(state, slug, "d", slide_number=1, role="result", title="R")
    # square image (1×1) into a portrait box with fit=cover → crop sides
    decks.set_slide_regions(state, slug, "d", s["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1, "fit": "cover",
         "x": 0.1, "y": 0.25, "w": 0.3, "h": 0.6},
    ])
    deck_render.render_region(state, slug, "d", s["id"], "r1")
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    prs = Presentation(str(out))
    pic = next(sh for sh in prs.slides[0].shapes
               if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert pic.crop_left > 0   # cover cropped the overflowing edges


# ─── export typography: title layout, markdown, fonts ────────


def test_font_family_strips_weight_words():
    assert deck_render._font_family("Inter Bold") == "Inter"
    assert deck_render._font_family("Source Serif Pro Regular") == "Source Serif Pro"
    assert deck_render._font_family("Inter") == "Inter"
    assert deck_render._font_family(None) is None
    assert deck_render._font_family("   ") is None


def test_theme_fonts_from_concept():
    concept = ("Typography:\n  display: Inter Bold   "
               "body: Source Serif Regular   mono: JetBrains Mono")
    f = deck_render._theme_fonts(concept)
    assert f["display"] == "Inter"
    assert f["body"] == "Source Serif"
    assert f["mono"] == "JetBrains Mono"


def test_theme_fonts_default_none():
    assert deck_render._theme_fonts(None) == {
        "display": None, "body": None, "mono": None,
    }


def test_export_title_slide_is_centered(state, tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN

    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="title",
                    title="My Big Talk", body="A subtitle line",
                    notes="n", render_mode="text")
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    prs = Presentation(str(out))
    tb = next(sh for sh in prs.slides[0].shapes if sh.has_text_frame)
    # cover layout: the title paragraph is centered, not top-anchored
    assert tb.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert "My Big Talk" in tb.text_frame.text


def test_export_body_is_plain_text_no_markdown_parsing(state, tmp_path, monkeypatch):
    """body markdown markers (** _ `) survive into the slide as literal
    text — the renderer is intentionally markdown-agnostic. Rich slide
    design is the job of the slide's `code` field, not the body parser
    (docs/todo/002)."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title="Background", body="plain and **strong** words",
                    notes="n", render_mode="text")
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    prs = Presentation(str(out))
    flat = " ".join(
        sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
    )
    assert "plain and **strong** words" in flat
    runs = [r for sh in prs.slides[0].shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs for r in para.runs]
    body_runs = [r for r in runs if "strong" in (r.text or "")]
    assert body_runs and body_runs[0].font.bold is not True


def test_export_applies_concept_fonts(state, tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.update_deck(
        state, slug, "d",
        concept="Typography:\n  display: Inter Bold   body: Lora Regular",
    )
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title="Heading", body="some body text",
                    notes="n", render_mode="text")
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    prs = Presentation(str(out))
    names = {r.font.name for sh in prs.slides[0].shapes if sh.has_text_frame
             for para in sh.text_frame.paragraphs for r in para.runs}
    assert "Inter" in names   # display font on the title
    assert "Lora" in names    # body font on the body text


# ─── Keynote-safe image normalization (todo 032) ────────────────────────


def test_normalize_image_rgba_becomes_rgb_jpeg(tmp_path):
    pytest.importorskip("PIL")
    from PIL import Image
    src = tmp_path / "rgba.png"
    Image.new("RGBA", (32, 32), (255, 0, 0, 128)).save(src)
    out = deck_render._normalize_image_for_pptx(str(src))
    # Returned as BytesIO; first bytes are the JPEG SOI marker (FFD8).
    assert hasattr(out, "read")
    head = out.read(3)
    assert head[:2] == b"\xff\xd8"
    out.seek(0)
    img = Image.open(out)
    assert img.mode == "RGB"   # alpha flattened against white background


def test_normalize_image_downsamples_oversized(tmp_path):
    pytest.importorskip("PIL")
    from PIL import Image
    src = tmp_path / "huge.png"
    Image.new("RGB", (4000, 2000), (200, 100, 50)).save(src)
    out = deck_render._normalize_image_for_pptx(str(src), target_width=1920)
    img = Image.open(out)
    assert img.width == 1920
    assert img.height == 960   # proportional


def test_normalize_image_keeps_smaller_rgb_at_size(tmp_path):
    pytest.importorskip("PIL")
    from PIL import Image
    src = tmp_path / "small.png"
    Image.new("RGB", (800, 600), (10, 20, 30)).save(src)
    out = deck_render._normalize_image_for_pptx(str(src), target_width=1920)
    img = Image.open(out)
    assert (img.width, img.height) == (800, 600)


def test_normalize_image_falls_back_on_unreadable_input(tmp_path):
    src = tmp_path / "garbage.png"
    src.write_bytes(b"not actually an image")
    # Must not raise — degraded gracefully so the export still produces a PPTX.
    out = deck_render._normalize_image_for_pptx(str(src))
    assert out == str(src)   # original path returned as last-ditch fallback


# ─── Theme-driven type scale (todo 035) ─────────────────────────────────


def test_type_scale_defaults_when_concept_absent():
    ts = deck_render._theme_type_scale(None)
    assert ts["title"] == 32
    assert ts["body"] == 20
    assert ts["head"] == 26
    assert ts["cover_title"] == 40
    assert ts["caption"] == 12
    assert ts["line_spacing"] == pytest.approx(1.22)


def test_type_scale_overrides_from_concept():
    concept = (
        "Type scale:\n"
        "  title: 30  head: 24  body: 18  line_spacing: 1.15\n"
        "  cover_title: 36"
    )
    ts = deck_render._theme_type_scale(concept)
    assert ts["title"] == 30
    assert ts["head"] == 24
    assert ts["body"] == 18
    assert ts["cover_title"] == 36
    assert ts["line_spacing"] == pytest.approx(1.15)
    # Unspecified keys keep defaults.
    assert ts["caption"] == 12


def test_type_scale_ignores_malformed_values():
    concept = "Type scale:\n  title: huge  body: 22"
    ts = deck_render._theme_type_scale(concept)
    assert ts["title"] == 32   # malformed → default
    assert ts["body"] == 22    # well-formed override applied


def test_export_applies_type_scale_to_title(state, tmp_path, monkeypatch):
    """A concept with `Type scale: title: 30` produces a 30pt title run
    instead of the default 32pt."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.update_deck(
        state, slug, "d",
        concept="Type scale:\n  title: 30  body: 18",
    )
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title="Heading", body="some body text",
                    notes="n", render_mode="text")
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    prs = Presentation(str(out))
    title_runs = [
        r for sh in prs.slides[0].shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs for r in para.runs
        if r.text == "Heading"
    ]
    assert title_runs
    # python-pptx font.size is an Emu/Pt object; check via .pt
    assert title_runs[0].font.size.pt == 30


# ─── code slides: agent-authored python-pptx snippets (todo 002) ────────


def test_render_deck_skips_code_slides(state):
    """A code slide's `code` runs at PPTX export time; render_deck has
    nothing to do for it."""
    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title="t", code="pass", render_mode="code")
    result = deck_render.render_deck(state, slug, "d")
    assert len(result["skipped"]) == 1
    assert result["skipped"][0]["reason"] == \
        "code slide — runs at PPTX export time"


def test_render_slide_code_mode_without_regions_raises(state):
    """A code slide with no regions has nothing to pre-render — render_
    slide raises so the caller knows the `code` runs at export time."""
    slug = _setup(state)
    s = decks.add_slide(state, slug, "d", slide_number=1, role="discussion",
                        title="t", code="pass", render_mode="code")
    with pytest.raises(ValueError, match="code slide has no regions"):
        deck_render.render_slide(state, slug, "d", s["id"])


def test_set_slide_regions_preserves_code_mode(state, tmp_path):
    """`set_slide_regions` on a code-mode slide keeps it in code mode
    (regions act as image placeholders for the snippet); on every other
    mode it still snaps to hybrid (legacy behavior). (todo 004 follow-
    up — image placeholder strategy.)"""
    slug = _setup(state)
    # Stage a figure for paper-figure regions
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))

    # Code-mode slide → set_slide_regions keeps it in code
    code_slide = decks.add_slide(
        state, slug, "d", slide_number=1, role="result", title="Code",
        notes="n", render_mode="code",
        code="h.image_region('r1', left=Inches(1), top=Inches(2), "
             "width=Inches(6), height=Inches(4))",
    )
    decks.set_slide_regions(state, slug, "d", code_slide["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1,
         "x": 0.1, "y": 0.25, "w": 0.5, "h": 0.5},
    ])
    refreshed = decks.list_slides(state, slug, "d")[0]
    assert refreshed["render_mode"] == "code"
    assert len(refreshed["regions"]) == 1

    # Text-mode slide → set_slide_regions snaps to hybrid (back-compat)
    text_slide = decks.add_slide(
        state, slug, "d", slide_number=2, role="result", title="Text",
        notes="n", render_mode="text",
    )
    decks.set_slide_regions(state, slug, "d", text_slide["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1,
         "x": 0.1, "y": 0.25, "w": 0.5, "h": 0.5},
    ])
    refreshed = decks.list_slides(state, slug, "d")[1]
    assert refreshed["render_mode"] == "hybrid"


def test_render_slide_code_mode_with_regions_renders_them(state, tmp_path):
    """render_slide on a code slide with regions = same path as a
    hybrid slide: each region's image is materialized so h.image_region
    can pick it up at export."""
    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    s = decks.add_slide(
        state, slug, "d", slide_number=1, role="result", title="Code+regions",
        notes="n", render_mode="code", code="pass",
    )
    decks.set_slide_regions(state, slug, "d", s["id"], regions=[
        {"render_mode": "paper-figure", "figure_number": 1,
         "x": 0.1, "y": 0.25, "w": 0.5, "h": 0.5},
    ])
    result = deck_render.render_slide(state, slug, "d", s["id"])
    assert result["mode"] == "hybrid"   # routed through hybrid renderer
    assert len(result["rendered"]) == 1
    # Region got its image
    refreshed = decks.list_slides(state, slug, "d")[0]
    assert refreshed["regions"][0]["image_blob_path"]


def test_render_deck_all_code_marks_rendered(state):
    """A deck of only code slides is fully 'rendered' — they materialize
    natively at export time, so they must not block the status flip."""
    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="discussion",
                    title="t", code="pass", render_mode="code")
    deck_render.render_deck(state, slug, "d")
    assert decks.get_deck(state, slug, "d")["status"] == "rendered"


def test_export_code_slide_executes_snippet(state, tmp_path, monkeypatch):
    """A `code` slide's snippet runs at export, populating the slide
    with whatever shapes the agent's code adds."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="discussion",
        title="My Title", notes="n", render_mode="code",
        code=(
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            "h.bullet_list(slide, ['alpha', 'beta', 'gamma'],\n"
            "              palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "              left=Inches(0.7), top=Inches(2.0),\n"
            "              width=sw - Inches(1.4), height=Inches(4.0))\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_slides"] == 1
    assert res["code_errors"] == []
    assert res["text_slides"] == 0
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    assert "My Title" in flat
    for token in ("alpha", "beta", "gamma"):
        assert token in flat


def test_export_code_slide_card_grid(state, tmp_path, monkeypatch):
    """The card_grid helper produces N cards arranged in `cols` columns."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Spec", notes="n", render_mode="code",
        code=(
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            "h.card_grid(slide, [\n"
            "    {'title': 'Memory', 'body': 'persists across sessions'},\n"
            "    {'title': 'Hooks', 'body': 'fire on events'},\n"
            "    {'title': 'Slash', 'body': 'reusable commands'},\n"
            "    {'title': 'Context', 'body': 'attention manager'},\n"
            "], left=Inches(0.7), top=Inches(2.0),\n"
            "   width=sw - Inches(1.4), height=Inches(4.0),\n"
            "   palette=palette, fonts=fonts, type_scale=type_scale, cols=2)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    slide = Presentation(str(out)).slides[0]
    # Each card adds: a card body rect (RECTANGLE), an accent stripe
    # (RECTANGLE), a title textbox, a body textbox. We expect 4 cards.
    flat = " ".join(
        sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
    )
    for tag in ("Memory", "Hooks", "Slash", "Context"):
        assert tag in flat
    for body in ("persists across sessions", "fire on events",
                 "reusable commands", "attention manager"):
        assert body in flat
    # 4 cards × (card rect + accent stripe) = 8 RECTANGLE auto-shapes,
    # plus the top accent stripe and the title accent rule = at least 10.
    auto_shapes = [
        sh for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]
    assert len(auto_shapes) >= 10


def test_export_code_slide_error_degrades_to_text(state, tmp_path, monkeypatch):
    """A broken `code` slide's exception is captured in `code_errors`
    and the slide degrades to plain text — the whole export still
    succeeds."""
    pytest.importorskip("pptx")

    slug = _setup(state)
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Broken", body="fallback text", notes="n",
        render_mode="code",
        code="raise RuntimeError('intentional test failure')",
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert out.is_file()
    assert res["code_slides"] == 0
    assert len(res["code_errors"]) == 1
    err = res["code_errors"][0]
    assert err["slide_number"] == 1
    assert "RuntimeError" in err["error"]
    assert "intentional test failure" in err["error"]
    # Slide degraded to text — body content survives.
    from pptx import Presentation
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    assert "fallback text" in flat
    assert res["text_slides"] == 1


# ─── Design grid (todo 004 §D) ──────────────────────────────────────────


def test_grid_cell_basic_geometry():
    """First cell sits at (margin_x, margin_top); column widths divide
    the usable width evenly; horizontal gutters separate columns."""
    from pptx.util import Inches, Pt
    from co_scientist_local.tools import slide_render_helpers as srh

    sw = Inches(13.333)  # 16:9 deck slide width
    sh = Inches(7.5)
    g = srh.grid(sw=sw, sh=sh, cols=12, rows=6)
    left, top, w, h = g.cell(col=1, span=1, row=1, row_span=1)
    assert left == g.margin_x
    assert top == g.margin_top
    assert w == g.col_w
    assert h == g.row_h


def test_grid_span_absorbs_gutters():
    """A span-3 cell is wider than 3 separate col-1 cells by 2 gutters
    — the cell straddles the gutters between covered columns."""
    from pptx.util import Inches, Pt
    from co_scientist_local.tools import slide_render_helpers as srh

    g = srh.grid(sw=Inches(13.333), sh=Inches(7.5), cols=12, rows=6)
    _, _, w_single, _ = g.cell(col=1, span=1, row=1)
    _, _, w_three, _ = g.cell(col=1, span=3, row=1)
    assert w_three == 3 * w_single + 2 * g.gutter


def test_grid_row_helper_spans_all_columns():
    """`g.row(...)` is full-width across `cols` columns."""
    from pptx.util import Inches
    from co_scientist_local.tools import slide_render_helpers as srh

    g = srh.grid(sw=Inches(13.333), sh=Inches(7.5))
    full = g.row(row=2)
    same = g.cell(col=1, span=g.cols, row=2)
    assert full == same


def test_grid_rejects_out_of_range():
    """Cells outside the grid raise ValueError, not silently clip."""
    from pptx.util import Inches
    from co_scientist_local.tools import slide_render_helpers as srh

    g = srh.grid(sw=Inches(13.333), sh=Inches(7.5), cols=12, rows=6)
    with pytest.raises(ValueError, match="col must be 1..12"):
        g.cell(col=0, row=1)
    with pytest.raises(ValueError, match="exceeds 12 cols"):
        g.cell(col=10, span=5, row=1)
    with pytest.raises(ValueError, match="row must be 1..6"):
        g.cell(col=1, row=7)
    with pytest.raises(ValueError, match="exceeds 6 rows"):
        g.cell(col=1, row=5, row_span=3)


def test_export_code_slide_uses_grid(state, tmp_path, monkeypatch):
    """End-to-end: a `code` snippet builds a grid + a card grid in it,
    everything renders, the cards land at the grid coordinates."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Grid demo", notes="n", render_mode="code",
        code=(
            "g = h.grid(sw=sw, sh=sh, cols=12, rows=4)\n"
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            "left, top, w, ht = g.cell(col=1, span=6, row=1, row_span=4)\n"
            "h.bullet_list(slide, ['left half'],\n"
            "              palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "              left=left, top=top, width=w, height=ht)\n"
            "left, top, w, ht = g.cell(col=7, span=6, row=1, row_span=4)\n"
            "h.bullet_list(slide, ['right half'],\n"
            "              palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "              left=left, top=top, width=w, height=ht)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_slides"] == 1
    assert res["code_errors"] == []
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    assert "left half" in flat
    assert "right half" in flat


def test_spacing_unit_exposed_to_snippets(state, tmp_path, monkeypatch):
    """A snippet can read `h.SPACING_UNIT_PT` for 8pt-rhythm math."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Rhythm", notes="n", render_mode="code",
        code=(
            "assert h.SPACING_UNIT_PT == 8\n"
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_slides"] == 1
    assert res["code_errors"] == []


# ─── per-slide PNG export (todo 004 §A — critique loop) ─────────────────


def _fake_render_pdf_to_pngs(_pdf_path, out_dir, *, dpi=150):
    """Pretend each PDF page rendered to a PNG. Used to test the
    critique-PNG wiring without a real PDF source."""
    pages = [
        out_dir / "slide_001.png",
        out_dir / "slide_002.png",
    ]
    for p in pages:
        p.write_bytes(_PNG_1x1)
    return pages


def test_export_emits_per_slide_pngs(state, tmp_path, monkeypatch):
    """When the PDF render succeeds, the export also produces one PNG
    per slide and uploads each to Storage for the critique loop."""
    pytest.importorskip("pptx")

    slug = _setup(state)
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title="Slide 1", body="x", render_mode="text")
    decks.add_slide(state, slug, "d", slide_number=2, role="conclusion",
                    title="Slide 2", body="y", render_mode="text")
    # Fake out the soffice + PNG-rendering steps so we don't need a
    # real LibreOffice / PDF round-trip in CI.
    fake_pdf = tmp_path / "deck.pdf"
    fake_pdf.write_bytes(b"%PDF-1.4 fake")
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: fake_pdf)
    monkeypatch.setattr(deck_render, "_render_pdf_to_pngs",
                        _fake_render_pdf_to_pngs)

    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["slide_pngs_skipped"] is False
    assert len(res["slide_pngs"]) == 2
    for i, png in enumerate(res["slide_pngs"], start=1):
        assert png["slide_number"] == i
        # Local file present, blob_path includes slide_NNN.png filename.
        assert pathlib.Path(png["local_path"]).is_file()
        assert png["blob_path"].endswith(f"slide_{i:03d}.png")
        # And it was uploaded to the backend.
        assert state.backend.get_blob(png["blob_path"])


def test_render_pdf_to_pngs_returns_empty_when_pdf_missing(tmp_path):
    """Best-effort: a missing / unreadable PDF returns [] silently
    rather than raising, so the rest of the export doesn't fail."""
    out = deck_render._render_pdf_to_pngs(
        tmp_path / "does-not-exist.pdf", tmp_path,
    )
    assert out == []


def test_export_code_slide_image_from_figure(state, tmp_path, monkeypatch):
    """`h.image_figure(N, ...)` resolves a paper figure and embeds it."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    decks.add_slide(
        state, slug, "d", slide_number=1, role="result",
        title="With figure", notes="n", render_mode="code",
        code=(
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            "h.image_figure(1, left=Inches(0.7), top=Inches(2.0),\n"
            "               width=sw - Inches(1.4), height=Inches(4.0))\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_slides"] == 1
    assert res["code_errors"] == []
    pics = [
        sh for sh in Presentation(str(out)).slides[0].shapes
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pics) == 1   # one image embedded via image_figure


# ─── Pattern library (todo 004 §B) ──────────────────────────────────────


def _run_pattern_slide(state, tmp_path, monkeypatch, code: str,
                       *, title: str = "T", body: str = ""):
    """Build a single deck with one code slide whose snippet is the
    given code string. Returns the export result + the first slide of
    the loaded Presentation."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title=title, body=body, notes="n",
                    render_mode="code", code=code)
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(
        state, slug, "d", output_path=str(out),
    )
    slide = Presentation(str(out)).slides[0]
    return res, slide


def _slide_text(slide):
    return " ".join(
        sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
    )


def test_pattern_hero_with_trailing_evidence(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.hero_with_trailing_evidence(slide, headline='The thesis',\n"
        "    evidence=['point one', 'point two', 'point three'],\n"
        "    palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "    sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    assert "The thesis" in txt
    for line in ("point one", "point two", "point three"):
        assert line in txt


def test_pattern_chapter_divider(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.chapter_divider(slide, chapter_label='Era II',\n"
        "    summary='Data goes digital',\n"
        "    palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "    sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    assert "Era II" in txt
    assert "Data goes digital" in txt


def test_pattern_metric_tile_row(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.metric_tile_row(slide, tiles=[\n"
        "    ('30', 'seconds per query', 's'),\n"
        "    ('500', 'accessions', None),\n"
        "    ('4', 'modalities', None),\n"
        "], palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "   sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    for token in ("30", "500", "4", "seconds per query", "accessions",
                  "modalities"):
        assert token in txt


def test_pattern_evidence_stack(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.evidence_stack(slide, claim='AI breeding scales',\n"
        "    evidence=[\n"
        "        {'tag': 'data', 'body': '500 accessions multi-modal'},\n"
        "        {'tag': 'speed', 'body': '30s vs 2 weeks per query'},\n"
        "    ], palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "       sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    assert "AI breeding scales" in txt
    assert "DATA" in txt
    assert "SPEED" in txt


def test_pattern_flow_pipeline(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.flow_pipeline(slide, steps=[\n"
        "    {'tag': 'Collect', 'body': 'multi-modal data'},\n"
        "    {'tag': 'Model', 'body': 'train on it'},\n"
        "    {'tag': 'Deploy', 'body': 'breeder queries'},\n"
        "], palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "   sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    for token in ("Collect", "Model", "Deploy", "multi-modal data",
                  "breeder queries"):
        assert token in txt


def test_pattern_before_after_split(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.before_after_split(slide,\n"
        "    before={'title': 'Manual', 'body': 'two weeks per query'},\n"
        "    after={'title': 'MCP', 'body': '30 seconds per query'},\n"
        "    transition_label='150× faster',\n"
        "    palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "    sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    assert "Manual" in txt
    assert "MCP" in txt
    assert "150" in txt


def test_pattern_contrast_pair(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.contrast_pair(slide,\n"
        "    left_item={'title': 'In-house', 'pros': ['control'],\n"
        "               'cons': ['slow']},\n"
        "    right_item={'title': 'Cloud', 'pros': ['fast'],\n"
        "                'cons': ['lock-in']},\n"
        "    axis_label='deploy time vs control',\n"
        "    palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "    sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    assert "In-house" in txt
    assert "Cloud" in txt
    assert "control" in txt
    assert "lock-in" in txt
    assert "DEPLOY TIME VS CONTROL" in txt


def test_pattern_quadrant_map(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.quadrant_map(slide, items=[\n"
        "    {'label': 'A', 'x': 0.2, 'y': 0.8},\n"
        "    {'label': 'B', 'x': 0.8, 'y': 0.2},\n"
        "], axes={'x': 'cost', 'y': 'impact',\n"
        "         'x_low': 'cheap', 'x_high': 'expensive'},\n"
        "   palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "   sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    assert "A" in txt
    assert "B" in txt
    assert "cost" in txt
    assert "impact" in txt


def test_pattern_numbered_milestone_arc(state, tmp_path, monkeypatch):
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=(
        "p.numbered_milestone_arc(slide, milestones=[\n"
        "    {'tag': 'Era I', 'note': 'paper records'},\n"
        "    {'tag': 'Era II', 'note': 'data digital'},\n"
        "    {'tag': 'Era III', 'note': 'practice digital'},\n"
        "], palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "   sw=sw, sh=sh)\n"
    ))
    assert res["code_errors"] == []
    txt = _slide_text(slide)
    for token in ("Era I", "Era II", "Era III", "paper records",
                  "data digital", "practice digital", "1", "2", "3"):
        assert token in txt


# ─── Pattern QA: contracts + bounds (todo 005 §D) ───────────────────────
#
# Every under-title pattern must coexist with the standard agent
# preamble (h.accent_stripe + h.title_block). The preamble adds exactly
# three shapes in the top "title zone" (y < ~Inches(1.65)):
#     1. accent_stripe rectangle (top=0)
#     2. title textbox (top=Inches(0.45))
#     3. accent rule rectangle (top=Inches(1.52))
# Any other shape with top < title_zone_bottom is a title-collision bug.
#
# chapter_divider owns the whole slide; we test it without the preamble.


_PREAMBLE = (
    "h.accent_stripe(slide, palette=palette, sw=sw)\n"
    "h.title_block(slide, 'My Title', palette=palette, fonts=fonts,\n"
    "              type_scale=type_scale, sw=sw, sh=sh)\n"
)


def _shape_in_bounds(sh, sw, sh_total, tol_emu: int = 200):
    """All shapes must sit within slide bounds (small EMU tolerance)."""
    return (
        sh.left is not None and sh.top is not None
        and sh.left >= -tol_emu
        and sh.top >= -tol_emu
        and sh.left + sh.width <= sw + tol_emu
        and sh.top + sh.height <= sh_total + tol_emu
    )


def _count_shapes_in_title_zone(slide):
    """Pixels-equivalent: count shapes whose top is above ~Inches(1.65)."""
    from pptx.util import Inches as _I
    return sum(1 for s in slide.shapes
               if s.top is not None and s.top < _I(1.65))


def _build_pattern_slide(state, tmp_path, monkeypatch, *,
                        snippet: str, with_preamble: bool = True):
    """One-slide deck running snippet (optionally with preamble) and
    return the loaded Presentation slide + raw bounds for assertion."""
    from pptx import Presentation
    pytest.importorskip("pptx")

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    code = (_PREAMBLE + snippet) if with_preamble else snippet
    decks.add_slide(state, slug, "d", slide_number=1, role="background",
                    title="My Title", notes="n", render_mode="code",
                    code=code)
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(
        state, slug, "d", output_path=str(out),
    )
    assert res["code_errors"] == [], res["code_errors"]
    prs = Presentation(str(out))
    return prs.slides[0], prs.slide_width, prs.slide_height


# Each (name, snippet, owns_whole_slide) — pattern + Korean-ish content
_UNDER_TITLE_PATTERNS = [
    ("hero_with_trailing_evidence",
     "p.hero_with_trailing_evidence(slide,\n"
     "  headline='앞으로 10년의 농업 AI는 데이터에 말을 거는 일이다.',\n"
     "  evidence=['500종 통합', '4 모달리티 결합', '쿼리당 30초'],\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("metric_tile_row",
     "p.metric_tile_row(slide, tiles=[('30','초당 쿼리','s'),\n"
     "  ('500','accession',None),('4','modality',None)],\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("evidence_stack",
     "p.evidence_stack(slide, claim='MCP가 새 파이프라인이다.',\n"
     "  evidence=[{'tag': '속도', 'body': '2주가 30초로'},\n"
     "            {'tag': '범위', 'body': '500 accession × 4 modality'},\n"
     "            {'tag': '재현', 'body': 'provenance trail이 모든 단계'}],\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("flow_pipeline",
     "p.flow_pipeline(slide, steps=[\n"
     "  {'tag': '수집', 'body': '다중 모달 accession 데이터'},\n"
     "  {'tag': '정제', 'body': 'DOI 도장 + 결측 처리'},\n"
     "  {'tag': '학습', 'body': '단일 임베딩 학습'},\n"
     "  {'tag': '배포', 'body': '자연어 쿼리 인터페이스'}],\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("before_after_split",
     "p.before_after_split(slide,\n"
     "  before={'title': '맞춤 파이프라인', 'body': '쿼리당 2주: 컬럼 손으로 조인, 플롯 손으로 코딩.'},\n"
     "  after={'title': 'MCP 쿼리', 'body': '쿼리당 30초: 자연어 입력 → provenance trail.'},\n"
     "  transition_label='150× 더 빠름',\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("contrast_pair",
     "p.contrast_pair(slide,\n"
     "  left_item={'title': '사내 서버',\n"
     "             'pros': ['완전한 통제','쿼리당 비용 없음'],\n"
     "             'cons': ['운영 부담','스케일 느림']},\n"
     "  right_item={'title': '클라우드 MCP',\n"
     "              'pros': ['분 단위 스케일','운영 부담 없음'],\n"
     "              'cons': ['쿼리당 비용','벤더 lock-in']},\n"
     "  axis_label='배포 시간 vs 운영 부담',\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("quadrant_map",
     "p.quadrant_map(slide,\n"
     "  items=[{'label':'A','x':0.2,'y':0.8},{'label':'B','x':0.8,'y':0.3},\n"
     "         {'label':'C','x':0.1,'y':0.2},{'label':'D','x':0.7,'y':0.6}],\n"
     "  axes={'x':'배포 시간','y':'재현성','x_low':'일 단위','x_high':'분 단위'},\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("numbered_milestone_arc",
     "p.numbered_milestone_arc(slide, milestones=[\n"
     "  {'tag':'Era I (2010–)','note':'시퀀스 예측. 유전체 안에서.'},\n"
     "  {'tag':'2018 · DL-on-SNP','note':'딥러닝이 SNP를 받아들이기 시작'},\n"
     "  {'tag':'2024 · Genomic LM','note':'언어모델이 유전체 안으로 진입'},\n"
     "  {'tag':'Now','note':'육종가가 데이터에 말을 거는 시대'}],\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    # Structural patterns (todo 006). title_slide is owns-slide → tested
    # separately. The 3 under-title structural patterns join the suite:
    ("title_and_body",
     "p.title_and_body(slide, title='제목과 본문',\n"
     "  body=['첫 번째 핵심 포인트', '두 번째 보완 포인트',\n"
     "        '세 번째 결론적인 한 줄'],\n"
     "  lead='AI 육종은 데이터에 말을 거는 일.',\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
    ("title_two_content",
     "p.title_two_content(slide, title='2-column generic',\n"
     "  left={'heading': '맞춤 파이프라인',\n"
     "        'bullets': ['컬럼 손으로 조인','플롯 손으로 코딩','쿼리당 2주']},\n"
     "  right={'heading': 'MCP 쿼리',\n"
     "         'bullets': ['자연어 입력','provenance trail','쿼리당 30초']},\n"
     "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"),
]


@pytest.mark.parametrize("name,snippet", _UNDER_TITLE_PATTERNS,
                         ids=[n for n, _ in _UNDER_TITLE_PATTERNS])
def test_pattern_no_title_collision(name, snippet, state, tmp_path,
                                    monkeypatch):
    """Each under-title pattern, run with the standard preamble:
    exactly 3 shapes sit in the title zone (the preamble chrome — top
    stripe, title textbox, title accent rule). Pattern content stays
    below."""
    slide, sw, sh_total = _build_pattern_slide(
        state, tmp_path, monkeypatch, snippet=snippet,
    )
    in_title = _count_shapes_in_title_zone(slide)
    assert in_title == 3, (
        f"{name}: expected 3 title-zone shapes (preamble chrome) but "
        f"got {in_title} — pattern leaks shapes into the title area"
    )


@pytest.mark.parametrize("name,snippet", _UNDER_TITLE_PATTERNS,
                         ids=[n for n, _ in _UNDER_TITLE_PATTERNS])
def test_pattern_stays_within_slide_bounds(name, snippet, state, tmp_path,
                                            monkeypatch):
    """Each pattern + preamble: every shape sits fully within the slide.
    Catches the Bug-A class of "first/last marker extends past edges"."""
    slide, sw, sh_total = _build_pattern_slide(
        state, tmp_path, monkeypatch, snippet=snippet,
    )
    out_of_bounds = [
        sh for sh in slide.shapes
        if not _shape_in_bounds(sh, sw, sh_total)
    ]
    assert not out_of_bounds, (
        f"{name}: {len(out_of_bounds)} shape(s) extend past slide edges. "
        f"First: left={out_of_bounds[0].left}, top={out_of_bounds[0].top}, "
        f"w={out_of_bounds[0].width}, h={out_of_bounds[0].height}, "
        f"slide={sw}×{sh_total}"
    )


def test_chapter_divider_owns_whole_slide(state, tmp_path, monkeypatch):
    """chapter_divider is the 'owns whole slide' pattern — no preamble.
    Bug C check: it must NOT leave an orphan accent stripe in the title
    zone without an anchoring label above. Label should sit near vertical
    center."""
    slide, sw, sh_total = _build_pattern_slide(
        state, tmp_path, monkeypatch, with_preamble=False, snippet=(
            "p.chapter_divider(slide, chapter_label='Era II',\n"
            "  summary='데이터가 디지털로 — 육종가는 종이 위에 있다.',\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "  sw=sw, sh=sh)\n"
        ),
    )
    # All shapes within slide bounds
    for sh in slide.shapes:
        assert _shape_in_bounds(sh, sw, sh_total), \
            f"chapter_divider shape out of bounds: top={sh.top}, " \
            f"left={sh.left}, w={sh.width}, h={sh.height}"
    # No shape in the title zone (pattern owns the slide; no orphan stripe).
    from pptx.util import Inches as _I
    in_title = _count_shapes_in_title_zone(slide)
    assert in_title == 0, (
        f"chapter_divider leaked {in_title} shape(s) into the title "
        "zone — it should own the whole canvas, not place orphan "
        "stripes above the label"
    )
    # Chapter label should sit near vertical center, not at ¾ down.
    text_frames = [sh for sh in slide.shapes if sh.has_text_frame
                   and "Era II" in sh.text_frame.text]
    assert text_frames
    label = text_frames[0]
    center_y = sh_total // 2
    # Label's vertical mid-point should be within ~Inches(1.0) of center
    label_mid = label.top + label.height // 2
    assert abs(label_mid - center_y) < _I(1.0), (
        f"chapter_divider label centered poorly: label_mid={label_mid}, "
        f"slide_center={center_y}"
    )


def test_pattern_title_slide_owns_whole_slide(state, tmp_path, monkeypatch):
    """title_slide is the 'owns whole slide' opener — no preamble.
    Eyebrow + big title + accent rule + subtitle, all centered."""
    from pptx.util import Inches as _I
    slide, sw, sh_total = _build_pattern_slide(
        state, tmp_path, monkeypatch, with_preamble=False, snippet=(
            "p.title_slide(slide,\n"
            "  title='다중 모달 농업 AI',\n"
            "  subtitle='강양재 · 한국육종학회 2026.06.05',\n"
            "  eyebrow='Lab seminar',\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "  sw=sw, sh=sh)\n"
        ),
    )
    # No shape leaks into the title-bar zone, since the pattern owns
    # the canvas — all content lives in a centered block.
    assert _count_shapes_in_title_zone(slide) == 0
    # All shapes within bounds.
    for sh in slide.shapes:
        assert _shape_in_bounds(sh, sw, sh_total)
    # The pattern wrote all three text bits.
    flat = " ".join(s.text_frame.text for s in slide.shapes
                    if s.has_text_frame)
    assert "다중 모달 농업 AI" in flat
    assert "강양재" in flat
    assert "LAB SEMINAR" in flat   # eyebrow renders upper-case


def test_pattern_title_and_image_grid(state, tmp_path, monkeypatch):
    """4-image grid with captions. Pattern is under-title — runs with
    the standard preamble; all 4 images embed and captions survive."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Stage 4 PNGs the pattern will embed
    paths = []
    for i in range(4):
        p = tmp_path / f"tile_{i}.png"
        p.write_bytes(_PNG_1x1)
        paths.append(str(p))

    images_repr = ",\n  ".join(
        f"{{'path': {p!r}, 'caption': '캡션 {i + 1}'}}"
        for i, p in enumerate(paths)
    )
    snippet = (
        "p.title_and_image_grid(slide, title='이미지 그리드',\n"
        f"  images=[\n  {images_repr}],\n"
        "  cols=2,\n"
        "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "  sw=sw, sh=sh)\n"
    )
    slide, sw, sh_total = _build_pattern_slide(
        state, tmp_path, monkeypatch, snippet=snippet,
    )
    # 3 preamble chrome shapes in title zone, none from the grid pattern
    assert _count_shapes_in_title_zone(slide) == 3
    # 4 PICTURE shapes embedded
    pics = [sh for sh in slide.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 4, f"expected 4 image tiles, got {len(pics)}"
    # Every shape within slide bounds
    for sh in slide.shapes:
        assert _shape_in_bounds(sh, sw, sh_total), \
            f"out-of-bounds shape at top={sh.top}, left={sh.left}"
    # All 4 captions survive
    flat = " ".join(s.text_frame.text for s in slide.shapes
                    if s.has_text_frame)
    for i in range(4):
        assert f"캡션 {i + 1}" in flat


# ─── API consistency follow-ups (todo 007) ──────────────────────────────


def test_extended_roles_accept_common_deck_vocabulary():
    """`hook`, `thesis`, `section`, `figure`, `image`, `content` are
    now accepted alongside the canonical 9 roles (todo 007 #1)."""
    from co_scientist_local.tools import papers
    from co_scientist_local.backends.memory import InMemoryBackend
    from co_scientist_local.state import State
    state = State(project_id="p", owner_uid="u", backend=InMemoryBackend())
    papers.create_paper(state, title="x")
    decks.create_deck(state, "x", title="d", deck_id="d")
    for r in ("hook", "thesis", "section", "figure", "image", "content"):
        decks.add_slide(state, "x", "d",
                        slide_number=hash(r) % 1000 + 1,
                        role=r, title=r, render_mode="text")


def test_grid_cell_returns_cell_namedtuple():
    """`g.cell(...)` returns a Cell namedtuple — supports tuple-
    unpacking AND attribute access (todo 007 axis 3)."""
    from co_scientist_local.tools import slide_render_helpers as srh
    from pptx.util import Inches
    g = srh.grid(sw=Inches(13.333), sh=Inches(7.5))
    cell = g.cell(col=1, span=2, row=1, row_span=1)
    # Attribute access
    assert cell.left == g.margin_x
    assert cell.top == g.margin_top
    assert cell.width == g.col_w * 2 + g.gutter
    assert cell.height == g.row_h
    # Tuple-unpacking still works (back-compat with old tuple form)
    left, top, w, h = cell
    assert (left, top, w, h) == (cell.left, cell.top, cell.width, cell.height)


def test_h_text_helper_one_call_textbox(state, tmp_path, monkeypatch):
    """`h.text(slide, text, ...)` is the one-call textbox convenience
    (todo 007 §D) — drops the add_textbox + text_frame + paragraph +
    run + font boilerplate to a single call."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="content",
        title="text helper", notes="n", render_mode="code",
        code=(
            "h.text(slide, '소제목',\n"
            "       left=Inches(0.7), top=Inches(2.0),\n"
            "       width=Inches(8), height=Inches(0.6),\n"
            "       palette=palette, size_pt=24, bold=True)\n"
            "h.text(slide, 'caption — muted',\n"
            "       left=Inches(0.7), top=Inches(2.7),\n"
            "       width=Inches(8), height=Inches(0.4),\n"
            "       palette=palette, size_pt=12,\n"
            "       color=palette['muted'], italic=True)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == []
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    assert "소제목" in flat
    assert "caption — muted" in flat


def test_palette_full_seven_keys_in_code_namespace(state, tmp_path, monkeypatch):
    """The exec namespace's `palette` exposes all 7 keys (todo 007
    axis 4) — accent / background / foreground / surface / muted /
    secondary / highlight. Each is an RGBColor."""
    pytest.importorskip("pptx")
    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept=(
        "Palette:\n"
        "  bg: #fff  text: #111  accent: #b58900\n"
        "  muted: #6c757d  secondary: #226699  highlight: #ee5500\n"
    ))
    decks.add_slide(
        state, slug, "d", slide_number=1, role="content",
        title="palette test", notes="n", render_mode="code",
        code=(
            "from pptx.dml.color import RGBColor\n"
            "for key in ('accent', 'background', 'foreground', 'surface',\n"
            "            'muted', 'secondary', 'highlight'):\n"
            "    assert key in palette, f'missing {key}'\n"
            "    assert isinstance(palette[key], RGBColor), key\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    assert res["code_slides"] == 1


def test_image_figure_accepts_slide_first_positional(state, tmp_path, monkeypatch):
    """Image helpers accept `slide` as the first positional argument
    for consistency with every other helper / pattern (todo 007 axis
    1). The shorter closure form still works for back-compat."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    p = tmp_path / "fig.png"
    p.write_bytes(_PNG_1x1)
    figures.add_figure(state, slug, figure_number=1, title="F", local_path=str(p))
    decks.add_slide(
        state, slug, "d", slide_number=1, role="figure",
        title="image", notes="n", render_mode="code",
        code=(
            # Slide-first form (the consistent one)
            "h.image_figure(slide, 1, left=Inches(0.5), top=Inches(2),\n"
            "               width=Inches(5), height=Inches(4))\n"
            # Closure / back-compat form (no slide)
            "h.image_figure(1, left=Inches(6), top=Inches(2),\n"
            "               width=Inches(5), height=Inches(4))\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    pics = [
        sh for sh in Presentation(str(out)).slides[0].shapes
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pics) == 2  # both forms succeeded


# ─── Tier 2: items canonical + dict shape detection (todo 007 axis 2) ──


def test_patterns_accept_items_canonical_name(state, tmp_path, monkeypatch):
    """Every list-of-items pattern accepts `items=` as the canonical
    kwarg (todo 007 axis 2). Legacy names (evidence / tiles / steps /
    milestones) still work as aliases — covered elsewhere by the
    `_UNDER_TITLE_PATTERNS` suite."""
    pytest.importorskip("pptx")
    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="content",
        title="items canonical", notes="n", render_mode="code",
        code=(
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            "# hero — items as list[str]\n"
            "p.hero_with_trailing_evidence(slide,\n"
            "  headline='canonical hero', items=['e1', 'e2', 'e3'],\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"
        ),
    )
    decks.add_slide(
        state, slug, "d", slide_number=2, role="content",
        title="items via flow", notes="n", render_mode="code",
        code=(
            "p.flow_pipeline(slide, items=[\n"
            "  {'tag': 'A', 'body': 'first'},\n"
            "  {'tag': 'B', 'body': 'second'},\n"
            "  {'tag': 'C', 'body': 'third'}],\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    assert res["code_slides"] == 2


def test_metric_tile_row_accepts_dict_items(state, tmp_path, monkeypatch):
    """`metric_tile_row` accepts dict-shaped items alongside tuples —
    {value, label, unit} OR canonical {tag, body, unit} (todo 007 axis 2)."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    slug = _setup(state)
    decks.add_slide(
        state, slug, "d", slide_number=1, role="result",
        title="metric dicts", notes="n", render_mode="code",
        code=(
            "p.metric_tile_row(slide, items=[\n"
            "  {'value': '30', 'label': 'seconds', 'unit': 's'},\n"
            "  {'tag': '500', 'body': 'accessions'},\n"          # canonical
            "  ('150', 'fold speedup', '×'),\n"                    # tuple ok
            "], palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "   sw=sw, sh=sh)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    for token in ("30", "seconds", "500", "accessions", "150", "fold speedup"):
        assert token in flat


def test_passing_both_items_and_legacy_alias_raises(state):
    """Loud error when caller passes both `items` and the legacy alias
    (catches typo / partial-rename mistakes)."""
    from co_scientist_local.tools import slide_patterns
    from pptx import Presentation
    pytest.importorskip("pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(TypeError, match="pass `items` OR"):
        slide_patterns.flow_pipeline(
            slide, items=[{"tag": "A", "body": "x"}],
            steps=[{"tag": "B", "body": "y"}],
            palette={}, fonts={}, type_scale={}, sw=0, sh=0,
        )


def test_numbered_milestone_arc_canonical_body_field(state, tmp_path, monkeypatch):
    """numbered_milestone_arc accepts both `note` (legacy) and `body`
    (canonical) inside each item dict (todo 007 axis 2)."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    slug = _setup(state)
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="milestone canonical body", notes="n", render_mode="code",
        code=(
            "p.numbered_milestone_arc(slide, items=[\n"
            "  {'tag': 'A', 'body': 'canonical body field'},\n"   # canonical
            "  {'tag': 'B', 'note': 'legacy note field'},\n"      # legacy
            "  {'tag': 'C', 'body': 'mixed works'}],\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "  sw=sw, sh=sh)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    assert "canonical body field" in flat
    assert "legacy note field" in flat


# ─── Text autofit (Korean-aware) ─────────────────────────────────────────


def test_estimate_text_width_pt_korean_wider_than_ascii():
    """Korean chars take ~1.0 × font_pt; ASCII letters ~0.55 × font_pt
    — so the same character count of Korean is much wider than ASCII."""
    from co_scientist_local.tools import slide_render_helpers as srh
    en = srh.estimate_text_width_pt("hello world", font_pt=20)
    ko = srh.estimate_text_width_pt("안녕하세요세상", font_pt=20)
    # Korean (7 chars wide) should be visibly wider than English
    # (11 chars but mostly narrow letters)
    assert ko > en
    # And ~20pt × 7 chars ≈ 140pt
    assert 100 < ko < 200


def test_autofit_pt_shrinks_when_text_overflows():
    """A long Korean paragraph in a narrow short box → autofit returns
    a font size smaller than start_pt."""
    from co_scientist_local.tools import slide_render_helpers as srh
    from pptx.util import Inches as _I
    text = "한국어 텍스트가 길어서 박스 안에 한 줄로 들어갈 수 없는 경우 " * 3
    fit = srh.autofit_pt(
        text,
        max_width_emu=_I(3), max_height_emu=_I(1),
        start_pt=24, line_spacing=1.22, min_pt=10,
    )
    assert fit < 24, "autofit should have shrunk; got the start size back"
    assert fit >= 10


def test_autofit_pt_keeps_start_size_when_text_fits():
    """Short text in a big box → autofit returns start_pt unchanged."""
    from co_scientist_local.tools import slide_render_helpers as srh
    from pptx.util import Inches as _I
    fit = srh.autofit_pt(
        "Hi", max_width_emu=_I(10), max_height_emu=_I(4),
        start_pt=20, line_spacing=1.22, min_pt=10,
    )
    assert fit == 20


def test_autofit_exposed_on_helpers_namespace(state, tmp_path, monkeypatch):
    """The exec namespace exposes h.autofit_pt so slide code can autofit
    its own custom text emits."""
    pytest.importorskip("pptx")
    slug = _setup(state)
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="autofit hook", notes="n", render_mode="code",
        code=(
            "assert h.autofit_pt('abc', max_width_emu=Inches(5),\n"
            "                    max_height_emu=Inches(2), start_pt=24) == 24\n"
            "# Long Korean text in a narrow box must shrink\n"
            "fit = h.autofit_pt('한국어 텍스트 ' * 30,\n"
            "                   max_width_emu=Inches(3), max_height_emu=Inches(1),\n"
            "                   start_pt=24, min_pt=10)\n"
            "assert 10 <= fit < 24, f'expected shrunk size, got {fit}'\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    assert res["code_slides"] == 1


# ─── Iconography (todo 004 §C) ──────────────────────────────────────────


def test_icon_names_lists_known_vocabulary():
    """h.icon_names() returns the full sorted vocabulary the agent can
    pick from. Must include at least one shape-based + one glyph-based
    entry so the agent knows both backends exist."""
    from co_scientist_local.tools import slide_render_helpers as srh
    names = srh.icon_names()
    assert isinstance(names, list)
    assert names == sorted(names)
    # Spot-checks: shape + glyph fallbacks both present
    assert "arrow-right" in names    # MSO_SHAPE built-in
    assert "lightning" in names      # MSO_SHAPE built-in
    assert "check" in names          # Unicode glyph fallback
    assert "info" in names           # Unicode glyph fallback


def test_icon_unknown_name_raises():
    """h.icon('nonexistent', ...) raises ValueError, not silently
    nothing — so typos are loud."""
    from co_scientist_local.tools import slide_render_helpers as srh
    from pptx import Presentation
    from pptx.util import Inches
    pytest.importorskip("pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(ValueError, match="unknown icon name"):
        srh.icon(slide, "totally-not-an-icon",
                 left=Inches(1), top=Inches(1), size=Inches(0.5),
                 palette={"accent": _h_accent()})


def _h_accent():
    """1-line helper — RGBColor for tests that don't need a full theme."""
    from pptx.dml.color import RGBColor
    return RGBColor(0xB5, 0x89, 0x00)


def test_icon_shape_renders_as_native_autoshape(state, tmp_path, monkeypatch):
    """h.icon('lightning', ...) produces an MSO_SHAPE auto-shape (not
    a picture, not a textbox) — the icon stays native + recolorable in
    PowerPoint."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Icons", notes="n", render_mode="code",
        code=(
            "h.icon(slide, 'lightning', left=Inches(1), top=Inches(2),\n"
            "       size=Inches(1), palette=palette)\n"
            "h.icon(slide, 'database', left=Inches(3), top=Inches(2),\n"
            "       size=Inches(1), palette=palette)\n"
            "h.icon(slide, 'warning', left=Inches(5), top=Inches(2),\n"
            "       size=Inches(1), palette=palette)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == []
    slide = Presentation(str(out)).slides[0]
    auto_shapes = [
        sh for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]
    # 3 icon shapes (and no chrome was added by this snippet)
    assert len(auto_shapes) == 3


def test_icon_glyph_renders_as_textbox(state, tmp_path, monkeypatch):
    """h.icon('check', ...) — for icons python-pptx doesn't ship as a
    shape — renders as a textbox with the Unicode glyph centered."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    slug = _setup(state)
    decks.update_deck(state, slug, "d", concept="accent: #b58900")
    decks.add_slide(
        state, slug, "d", slide_number=1, role="background",
        title="Check + X", notes="n", render_mode="code",
        code=(
            "h.icon(slide, 'check', left=Inches(1), top=Inches(2),\n"
            "       size=Inches(0.6), palette=palette)\n"
            "h.icon(slide, 'x', left=Inches(3), top=Inches(2),\n"
            "       size=Inches(0.6), palette=palette)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == []
    flat = " ".join(
        sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes
        if sh.has_text_frame
    )
    assert "✓" in flat
    assert "✗" in flat


# ─── Semantic type roles (todo 004 §E) ──────────────────────────────────


def test_type_scale_includes_semantic_role_keys():
    """The default type_scale exposes semantic role names alongside the
    legacy keys, so new code can pick by role (display_hero, body_small,
    label_tag) instead of arbitrary names."""
    ts = deck_render._theme_type_scale(None)
    for role in (
        "display_cover", "display_hero", "display_chapter",
        "headline_section", "title_slide",
        "body_large", "body_standard", "body_small",
        "label_tag", "label_caption", "scale_ratio",
    ):
        assert role in ts, f"missing semantic role: {role}"
    # scale_ratio is a float; others are ints.
    assert isinstance(ts["scale_ratio"], float)
    assert isinstance(ts["display_hero"], int)
    # Sanity ordering:
    # chapter > cover > hero > slide-title > headline-section >
    # body_large > body_standard > body_small > label_*
    assert ts["display_chapter"] > ts["display_cover"]
    assert ts["display_cover"] > ts["display_hero"]
    assert ts["display_hero"] > ts["title_slide"]
    assert ts["title_slide"] > ts["headline_section"]
    assert ts["body_large"] >= ts["body_standard"] >= ts["body_small"]
    assert ts["body_small"] > ts["label_tag"]


def test_type_scale_semantic_role_override(state, tmp_path, monkeypatch):
    """A concept's Type-scale block can override the new semantic keys
    too (the parser auto-handles all type_scale keys)."""
    concept = (
        "Type scale:\n"
        "  display_hero: 36  body_small: 14  scale_ratio: 1.2\n"
    )
    ts = deck_render._theme_type_scale(concept)
    assert ts["display_hero"] == 36
    assert ts["body_small"] == 14
    assert ts["scale_ratio"] == pytest.approx(1.2)
    # unspecified keys keep defaults
    assert ts["display_cover"] == 48


def test_pattern_zoom_in_callout(state, tmp_path, monkeypatch):
    """zoom_in_callout needs an image source; smoke-test with a 1×1 PNG."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    src = tmp_path / "ctx.png"
    src.write_bytes(_PNG_1x1)
    code = (
        f"p.zoom_in_callout(slide, context_image_path={str(src)!r},\n"
        "    callout={'x': 0.3, 'y': 0.3, 'w': 0.3, 'h': 0.3},\n"
        "    note='zoom note',\n"
        "    palette=palette, fonts=fonts, type_scale=type_scale,\n"
        "    sw=sw, sh=sh)\n"
    )
    res, slide = _run_pattern_slide(state, tmp_path, monkeypatch, code=code)
    assert res["code_errors"] == []
    pics = [sh for sh in slide.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    # Two pictures: full context + zoomed inset
    assert len(pics) == 2
    txt = _slide_text(slide)
    assert "zoom note" in txt


# ─── figure_full + hero adaptive (todo 008) ─────────────────────────────


def test_figure_full_renders_image_and_caption(state, tmp_path, monkeypatch):
    """`p.figure_full` embeds one PICTURE that fills the body area + a
    caption textbox in the bottom margin (todo 008 §A)."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches as _I

    slug = _setup(state)
    img = tmp_path / "fig.png"
    img.write_bytes(_PNG_1x1)
    decks.add_slide(
        state, slug, "d", slide_number=1, role="figure",
        title="figure-full demo", notes="n", render_mode="code",
        code=(
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            f"p.figure_full(slide, image_path={str(img)!r},\n"
            "  caption='Fig 1 · the only figure on this slide',\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "  sw=sw, sh=sh)\n"
        ),
    )
    monkeypatch.setattr(deck_render, "_pdf_via_soffice", lambda _p: None)
    out = tmp_path / "deck.pptx"
    res = deck_render.export_deck_to_pptx(state, slug, "d", output_path=str(out))
    assert res["code_errors"] == [], res["code_errors"]
    slide = Presentation(str(out)).slides[0]
    # Exactly one picture (the figure), one caption textbox containing
    # "Fig 1 · …"
    pics = [sh for sh in slide.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    # The picture occupies more than 60% of the slide height (full-grid).
    sh_total = Presentation(str(out)).slide_height
    assert pics[0].height > sh_total * 0.55, \
        f"figure should fill the body area; got {pics[0].height}/{sh_total}"
    flat = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Fig 1 · the only figure on this slide" in flat


def test_figure_full_requires_exactly_one_image_source(state, tmp_path):
    """Loud error when caller passes neither or both image-source kwargs."""
    from co_scientist_local.tools import slide_patterns
    from pptx import Presentation
    pytest.importorskip("pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(TypeError, match="EXACTLY ONE"):
        slide_patterns.figure_full(
            slide, palette={"foreground": None, "accent": None},
            fonts={}, type_scale={}, sw=0, sh=0,
        )


def test_hero_with_trailing_evidence_tight_on_short_content():
    """Per todo 008 §C3 — `hero_with_trailing_evidence` should NOT
    spread short items across the full slide height. The pattern
    auto-switches to a tight per-item row height when avg item
    length < 50 chars."""
    # We can't easily measure shape coordinates without rendering, so
    # use the export path and check that the right-column textboxes are
    # closer together than the slide-height-divided naïve placement.
    pytest.importorskip("pptx")
    from pptx import Presentation
    from co_scientist_local.backends.memory import InMemoryBackend
    from co_scientist_local.state import State
    s = State(project_id="x", owner_uid="u", backend=InMemoryBackend())
    papers.create_paper(s, title="t")
    decks.create_deck(s, "t", title="d", deck_id="d")
    decks.add_slide(
        s, "t", "d", slide_number=1, role="thesis",
        title="hero short", notes="n", render_mode="code",
        code=(
            "h.accent_stripe(slide, palette=palette, sw=sw)\n"
            "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
            "              type_scale=type_scale, sw=sw, sh=sh)\n"
            "p.hero_with_trailing_evidence(slide,\n"
            "  headline='Short thesis line.',\n"
            "  items=['Speed', 'Breadth', 'Provenance'],\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "  sw=sw, sh=sh)\n"
        ),
    )
    import tempfile, pathlib as _p
    with tempfile.TemporaryDirectory() as tmpd:
        out = _p.Path(tmpd) / "deck.pptx"
        deck_render._pdf_via_soffice = lambda _x: None   # monkey-skip
        res = deck_render.export_deck_to_pptx(
            s, "t", "d", output_path=str(out),
        )
        assert res["code_errors"] == [], res["code_errors"]
        slide = Presentation(str(out)).slides[0]
        # Collect right-column evidence number boxes ("01" / "02" / "03")
        ev_runs = [
            sh for sh in slide.shapes
            if sh.has_text_frame
            and sh.text_frame.text.strip() in {"01", "02", "03"}
        ]
        assert len(ev_runs) == 3
        ys = sorted(sh.top for sh in ev_runs)
        gap_01_02 = ys[1] - ys[0]
        # Without the adaptive fix, gap ≈ body_h / 3 ≈ Inches(1.7) = ~1.5M EMU.
        # With the adaptive fix (short content → row_h ≈ body_pt * 4.2pt =
        # ~84pt ≈ 1.07M EMU), gap should be substantially smaller.
        from pptx.util import Inches as _I
        assert gap_01_02 < _I(1.4), (
            f"hero should tighten row gap for short content; got "
            f"gap={gap_01_02} EMU between '01' and '02'"
        )


# ─── Reference corpus (todo 004 §F) ─────────────────────────────────────


def test_reference_corpus_manifest_well_formed():
    """The shipped corpus manifest lists every pattern with `file`, `do`,
    and `dont` keys; every referenced file exists on disk."""
    import json
    corpus_dir = pathlib.Path(__file__).resolve().parents[1] \
        / "packages" / "skills" / "paper-deck" / "reference_corpus"
    manifest_path = corpus_dir / "manifest.json"
    assert manifest_path.is_file(), "reference corpus manifest missing"
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert "patterns" in data and isinstance(data["patterns"], dict)
    assert len(data["patterns"]) >= 8, \
        f"expected 8+ patterns in corpus, got {len(data['patterns'])}"
    for pattern, entry in data["patterns"].items():
        assert "file" in entry, f"{pattern}: missing file"
        assert "do" in entry and entry["do"], f"{pattern}: missing do"
        assert "dont" in entry and entry["dont"], f"{pattern}: missing dont"
        png_path = corpus_dir / entry["file"]
        assert png_path.is_file(), \
            f"{pattern}: PNG {entry['file']} not on disk"
        # PNG header check (8 bytes)
        head = png_path.read_bytes()[:8]
        assert head == b"\x89PNG\r\n\x1a\n", \
            f"{pattern}: {entry['file']} isn't a PNG"

