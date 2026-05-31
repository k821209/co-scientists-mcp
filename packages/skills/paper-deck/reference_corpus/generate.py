"""Regenerate the reference corpus PNGs (todo 004 §F).

Each entry is one canonical pattern + a curated content example, exported
through the same `export_deck_to_pptx` → soffice → PyMuPDF pipeline an
agent uses, so the references are pixel-identical to what the agent will
produce. The manifest.json records the pattern, content, do, and dont
for each example so the agent can grep for "show me a good flow_pipeline"
and Read the matching PNG.

Run:

    PYTHONPATH=apps/local-mcp python \\
        packages/skills/paper-deck/reference_corpus/generate.py

Requires soffice + pymupdf in scope (same as the export pipeline).
"""
from __future__ import annotations

import json
import pathlib
import shutil
import sys

# Allow running from the repo root: PYTHONPATH=apps/local-mcp
_THIS = pathlib.Path(__file__).resolve()
_REPO = _THIS.parents[4]
sys.path.insert(0, str(_REPO / "apps" / "local-mcp"))

from co_scientist_local.backends.memory import InMemoryBackend
from co_scientist_local.state import State
from co_scientist_local.tools import deck_render, decks, papers
from PIL import Image


PREAMBLE = (
    # Deck chrome defaults — each entry can override before PREAMBLE
    # by assigning EYE / FOOTER / PAGE / TOTAL. (todo 011 — every
    # exemplar carries chrome so the agent learns chrome is part of
    # every slide, not an optional add-on.)
    "try:\n  EYE\nexcept NameError:\n  EYE = ''\n"
    "try:\n  FOOTER\nexcept NameError:\n  FOOTER = 'Reference corpus · paper-deck'\n"
    "try:\n  PAGE\nexcept NameError:\n  PAGE = 1\n"
    "try:\n  TOTAL\nexcept NameError:\n  TOTAL = 12\n"
    "h.accent_stripe(slide, palette=palette, sw=sw)\n"
    "h.title_block(slide, title, palette=palette, fonts=fonts,\n"
    "              type_scale=type_scale, sw=sw, sh=sh)\n"
    "h.deck_chrome(slide, palette=palette, fonts=fonts,\n"
    "              type_scale=type_scale, sw=sw, sh=sh,\n"
    "              eyebrow=EYE, footer=FOOTER,\n"
    "              page_number=PAGE, total=TOTAL)\n"
)


# Each entry: (filename, pattern, title, code, owns_slide, do, dont)
REFERENCES = [
    (
        "title_slide.png", "title_slide",
        "Multi-modal AI for plant breeding",
        "p.title_slide(slide,\n"
        "  title='Multi-modal AI for plant breeding',\n"
        "  subtitle='Yang Jae Kang · Korean Breeding Society 2026',\n"
        "  eyebrow='Lab seminar · 30 min',\n"
        "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n",
        True,
        "Eyebrow (kicker label) above the big title, accent rule under it, "
        "subtitle (author + venue + date) on a third line. All centered, "
        "all vertically balanced. Use for the *opening* slide of a deck.",
        "Don't add body bullets or a closing tagline here — keep the slide "
        "to four short text lines max. The title IS the slide; whitespace "
        "carries the rest. For a centered emphasis MID-deck, use "
        "`chapter_divider` instead.",
    ),
    (
        "chapter_divider.png", "chapter_divider",
        "Section opener",
        "p.chapter_divider(slide, chapter_label='Era II',\n"
        "  summary='Data goes digital — but the breeder still works on paper.',\n"
        "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n",
        True,
        "Massive 56pt+ chapter label vertically-centered, short accent rule "
        "*under* the label (visual anchor), italic summary as a tagline. "
        "Use between sections of a long talk so the audience feels the break.",
        "Don't repeat the deck-level accent stripe at the top — this pattern "
        "owns the canvas. Don't make the summary a paragraph; one tight "
        "sentence (≤ 50 chars) carries the rhythm.",
    ),
    (
        "metric_tile_row.png", "metric_tile_row",
        "Quantitative summary",
        PREAMBLE + (
            "p.metric_tile_row(slide, items=[\n"
            "  ('30', 'seconds per query', 's'),\n"
            "  ('500', 'accessions', None),\n"
            "  ('4', 'modalities', None),\n"
            "  ('150', 'fold speedup', '×'),\n"
            "], palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "   sw=sw, sh=sh)\n"
        ),
        False,
        "3–5 big numbers (display type, accent color) with thin labels "
        "and optional units. Numbers carry the headline emotion; labels "
        "carry the meaning. Use for KPI / result slides where the "
        "magnitudes themselves are the message.",
        "Don't put 6+ tiles — the row gets cramped and each number "
        "loses weight. Don't write a long label; ≤ ~24 chars per tile, "
        "one line only.",
    ),
    (
        "flow_pipeline.png", "flow_pipeline",
        "Process pipeline",
        PREAMBLE + (
            "p.flow_pipeline(slide, items=[\n"
            "  {'tag': 'Collect',  'body': 'multi-modal accession data'},\n"
            "  {'tag': 'Curate',   'body': 'tidy + DOI-stamp every field'},\n"
            "  {'tag': 'Embed',    'body': 'train a single embedding'},\n"
            "  {'tag': 'Deploy',   'body': 'natural-language query interface'},\n"
            "], palette=palette, fonts=fonts, type_scale=type_scale,\n"
            "   sw=sw, sh=sh)\n"
        ),
        False,
        "Numbered cards (01 / 02 / 03 / 04) connected by right-arrows. "
        "Use for *sequential* processes — method, workflow, lifecycle. "
        "3–5 steps is the readable range.",
        "Don't use this for non-sequential lists (peers without an order) "
        "— compose a bespoke grid from `h.*` instead. Don't write a "
        "long body per step; the readability comes from short headings "
        "and quick descriptions.",
    ),
    (
        "figure_full.png", "figure_full",
        "Reference architecture (figure-only slide)",
        PREAMBLE + (
            "p.figure_full(slide, image_path=REF_TILE_BLUE,\n"
            "  caption='Fig 3 · LLM (사고) ↔ MCP (도구) ↔ Harness (운영). 자연어가 결정, MCP가 다리, harness가 운영.',\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"
        ),
        False,
        "Image owns the **full grid** (rows 1–6 ≈ 85% of slide height). "
        "Caption tucks into the 0.6\" bottom-margin strip OUTSIDE the "
        "grid, in muted body italics. Wins back ~17% of figure area "
        "vs the old row_span=4 + caption-in-row-6 layout (todo 008 §A).",
        "Don't waste a grid row as 'breathing space' between figure "
        "and caption — the bottom margin already does that job. Don't "
        "let the caption run > 120 chars; 1 line preferred, 2 max.",
    ),
    (
        "gantt_chart.png", "gantt_chart",
        "추진 일정 (M1–M8) 및 업무 분장",
        ("EYE = 'WHEN · 7개월 병렬 운영'\n"
         "FOOTER = '참조 사업 · ㈜디보'\n"
         "PAGE = 9; TOTAL = 13\n"
         + PREAMBLE
         + "p.gantt_chart(slide, items=[\n"
           "  {'label': '착수보고 및 계획 확정',    'start': 1, 'span': 1},\n"
           "  {'label': '시료 QC, HMW DNA 추출', 'start': 1, 'span': 2},\n"
           "  {'label': 'PacBio HiFi 시퀀싱',     'start': 2, 'span': 2},\n"
           "  {'label': 'ONT PromethION 시퀀싱',  'start': 2, 'span': 2},\n"
           "  {'label': 'DNBSEQ-G99 시퀀싱',      'start': 2, 'span': 3},\n"
           "  {'label': 'K-mer / 어셈블리',        'start': 3, 'span': 2},\n"
           "  {'label': 'Assembly QC',            'start': 4, 'span': 2},\n"
           "  {'label': '매핑 / 변이 호출',        'start': 4, 'span': 3},\n"
           "  {'label': '중간보고',                'start': 5, 'span': 1},\n"
           "  {'label': 'Introgression / Fst',    'start': 6, 'span': 2},\n"
           "  {'label': '종 특이적 마커 도출',     'start': 7, 'span': 2},\n"
           "  {'label': '최종 보고서',             'start': 7, 'span': 2}],\n"
           "  period_count=8,\n"
           "  palette=palette, fonts=fonts, type_scale=type_scale,\n"
           "  sw=sw, sh=sh)\n"),
        False,
        "Activity rows × period columns with accent-colored bars at "
        "each row's `{start, span}`. Zebra row backgrounds + period "
        "labels across the top + label column on the left. Pair with "
        "`h.deck_chrome` (eyebrow / footer / page number) for "
        "proposal-grade rhythm. Use for project timelines, parallel "
        "workstreams, multi-month delivery plans.",
        "Don't squeeze long Korean activity labels into a narrow label "
        "column — keep labels ≤ ~20 chars or let them wrap. Don't "
        "stack > ~12 rows; readability falls off.",
    ),
    (
        "title_and_image_grid.png", "title_and_image_grid",
        "Multi-image overview",
        PREAMBLE + (
            "p.title_and_image_grid(slide, title=title,\n"
            "  images=[{'path': REF_TILE_RED, 'caption': 'Era I — paper records'},\n"
            "          {'path': REF_TILE_BLUE, 'caption': 'Era II — data digital'},\n"
            "          {'path': REF_TILE_GREEN, 'caption': 'Era III — practice digital'},\n"
            "          {'path': REF_TILE_PURPLE, 'caption': 'Now — natural-language interface'}],\n"
            "  cols=2,\n"
            "  palette=palette, fonts=fonts, type_scale=type_scale, sw=sw, sh=sh)\n"
        ),
        False,
        "N images in a `cols`-column grid with italic captions below each. "
        "2×2 = comparison (peers, no order). N×1 row = progression "
        "(left → right reads as sequence). Captions stay in body type, "
        "never display.",
        "Don't use a grid for one hero image (`h.image_*` full-bleed). "
        "Don't pile in 6+ tiles — they get small and lose weight. Don't "
        "skip captions; the grid needs labeling to land.",
    ),
    # ─── Dense bespoke exemplars (todo 011) ───────────────────────────
    # These two slides show what a proposal-grade content slide actually
    # looks like — 40–60 shapes, hand-composed from h.* primitives.
    # NO single pattern call. The agent reads these and learns: when the
    # content asks for density, GO BESPOKE.
    (
        "proposal_dense.png", "proposal_dense",
        "Part 1 — 참조 유전체 구축 (3중 플랫폼 전략)",
        ("EYE = 'HOW · 추진 방법'\n"
         "FOOTER = '기러기류 종 특이적 마커 발굴 · ㈜디보'\n"
         "PAGE = 5; TOTAL = 13\n"
         + PREAMBLE
         + """
# ── 3 platform comparison cards (top half) ────────────────────────────
g = h.grid(sw=sw, sh=sh, cols=12, rows=6,
           margin_top=Inches(2.0), margin_bot=Inches(0.6))
plats = [
    {"name": "PacBio HiFi",     "mode": "외주 시퀀싱", "spec": "≥30×, Q30+",   "tag": "단일 리드 정확도"},
    {"name": "ONT PromethION",  "mode": "자체 운영",   "spec": "≥20×, R10.4.1, Q20+", "tag": "Ultra-long (>100 kb), 복원서열 통과"},
    {"name": "DNBSEQ-G99",      "mode": "자체 운영",   "spec": "Short-read",    "tag": "Hybrid polishing + k-mer QC"},
]
for i, plat in enumerate(plats):
    cell = g.cell(col=1 + i*4, span=4, row=1, row_span=2)
    # Card body + top accent stripe (hand-composed; NOT a pattern call)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        cell.left, cell.top, cell.width, cell.height)
    box.fill.solid(); box.fill.fore_color.rgb = palette["surface"]
    box.line.color.rgb = palette["accent"]; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        cell.left, cell.top, cell.width, Pt(4))
    stripe.line.fill.background()
    stripe.fill.solid(); stripe.fill.fore_color.rgb = palette["accent"]
    stripe.shadow.inherit = False
    h.text(slide, plat["name"],
           left=cell.left + Pt(10), top=cell.top + Pt(10),
           width=cell.width - Pt(20), height=Pt(28),
           palette=palette, size_pt=18, bold=True,
           font_name=fonts.get("display"))
    h.text(slide, plat["mode"],
           left=cell.left + Pt(10), top=cell.top + Pt(38),
           width=cell.width - Pt(20), height=Pt(18),
           palette=palette, size_pt=11, color=palette["muted"],
           font_name=fonts.get("body"))
    h.text(slide, plat["spec"],
           left=cell.left + Pt(10), top=cell.top + Pt(62),
           width=cell.width - Pt(20), height=Pt(20),
           palette=palette, size_pt=14, bold=True,
           font_name=fonts.get("body"))
    h.text(slide, plat["tag"],
           left=cell.left + Pt(10), top=cell.top + Pt(88),
           width=cell.width - Pt(20), height=Pt(40),
           palette=palette, size_pt=11, color=palette["muted"],
           font_name=fonts.get("body"))

# ── 5-stage pipeline row (hand-composed) ──────────────────────────────
stages = ["시료 QC HMW DNA", "라이브러리 시퀀싱", "K-mer (Meryl/GS2)",
          "de novo (Verkko)", "Hybrid Polishing"]
h.text(slide, "어셈블리 파이프라인 (척추동물 reference-grade 표준 접근)",
       left=Inches(0.6), top=Inches(4.5),
       width=Inches(12), height=Pt(20),
       palette=palette, size_pt=13, bold=True,
       font_name=fonts.get("display"))
stage_top = Inches(4.85)
stage_h = Inches(0.7)
stage_total_w = sw - Inches(1.2)
stage_w = (stage_total_w - Pt(8) * (len(stages) - 1)) // len(stages)
for i, s in enumerate(stages):
    sx = Inches(0.6) + i * (stage_w + Pt(8))
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        sx, stage_top, stage_w, stage_h)
    box.fill.solid(); box.fill.fore_color.rgb = palette["background"]
    box.line.color.rgb = palette["secondary"]; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    h.text(slide, f"{i+1}", left=sx + Pt(6), top=stage_top + Pt(4),
           width=Pt(20), height=Pt(20),
           palette=palette, size_pt=10, bold=True,
           color=palette["secondary"], font_name=fonts.get("display"))
    h.text(slide, s, left=sx + Pt(6), top=stage_top + Pt(22),
           width=stage_w - Pt(12), height=stage_h - Pt(24),
           palette=palette, size_pt=11, font_name=fonts.get("body"))

# ── Target metrics line at the bottom ─────────────────────────────────
h.text(slide, "▶ 목표 산출물",
       left=Inches(0.6), top=Inches(5.8),
       width=Inches(2.0), height=Pt(20),
       palette=palette, size_pt=12, bold=True,
       color=palette["accent"], font_name=fonts.get("display"))
h.text(slide, "Contig N50 ≥ 5 Mb · BUSCO(aves_odb10) C ≥ 95% · Mercury QV ≥ 40 · scaffold/T2T-급 chromosome-level reference + Mitogenome",
       left=Inches(2.6), top=Inches(5.8),
       width=sw - Inches(3.2), height=Pt(34),
       palette=palette, size_pt=11, color=palette["foreground"],
       font_name=fonts.get("body"))
"""),
        False,
        "**Bespoke** — NOT a pattern call. ~50 shapes hand-composed "
        "from h.* primitives: 3-platform comparison cards (top half) + "
        "5-stage pipeline row + target-metrics line + deck chrome. "
        "Use as the reference rendering whenever the slide's content "
        "asks for *proposal-grade density* (3+ structured comparison "
        "sections / multiple coordinated layouts / equipment lists / "
        "personnel pages). Read this exemplar's source in "
        "`reference_corpus/generate.py` for the canonical density "
        "vocabulary.",
        "Don't shoehorn this kind of content into a single pattern "
        "(`card_grid` / `metric_tile_row` / etc.). Pattern calls top "
        "out at ~15 shapes; the content here needs ~50. The "
        "right answer is to compose directly — that's the whole point "
        "of `code` mode.",
    ),
    (
        "personnel_equipment.png", "personnel_equipment",
        "수행 인력 및 보유 장비",
        ("EYE = 'WHO · 사업 수행 능력'\n"
         "FOOTER = '기러기류 종 특이적 마커 발굴 · ㈜디보'\n"
         "PAGE = 10; TOTAL = 13\n"
         + PREAMBLE
         + """
# ── LEFT half: personnel table (hand-composed) ────────────────────────
h.text(slide, "수행 조직 (4명 : 박사 2 + 석사 2)",
       left=Inches(0.6), top=Inches(2.0),
       width=Inches(5.8), height=Pt(20),
       palette=palette, size_pt=13, bold=True,
       font_name=fonts.get("display"))
people = [
    {"name": "강양재",   "role": "연구책임자 (PI)", "pct": "20%",
     "expertise": "박사 · 농학 · 분석 총괄 → 어셈블리 · 변이 · 마커 발굴 검토"},
    {"name": "박진하",   "role": "시퀀싱 · 실험 책임자", "pct": "50%",
     "expertise": "석사 · DNA QC · 라이브러리 · ONT/DNBSEQ 직접 운영"},
    {"name": "유준일",   "role": "자문 · 분석 백업", "pct": "20%",
     "expertise": "박사 · 정형성 데이터 처리 · 변이 호출 보조"},
    {"name": "이유희",   "role": "행정 지원 · 실험 보조", "pct": "10%",
     "expertise": "석사 · 보고서 작성 · 시료 운영 · 보관"},
]
y = Inches(2.4)
for p in people:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.6), y, Inches(5.8), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = palette["surface"]
    box.line.color.rgb = palette["muted"]; box.line.width = Pt(0.5)
    box.shadow.inherit = False
    # Top-row: name + percentage; bottom-row: role + expertise
    h.text(slide, p["name"],
           left=Inches(0.7), top=y + Pt(4),
           width=Inches(1.4), height=Pt(20),
           palette=palette, size_pt=13, bold=True,
           font_name=fonts.get("display"))
    h.text(slide, p["pct"],
           left=Inches(5.5), top=y + Pt(4),
           width=Inches(0.7), height=Pt(20),
           palette=palette, size_pt=14, bold=True,
           color=palette["accent"], font_name=fonts.get("display"))
    h.text(slide, p["role"],
           left=Inches(2.1), top=y + Pt(4),
           width=Inches(3.3), height=Pt(20),
           palette=palette, size_pt=11, color=palette["muted"],
           font_name=fonts.get("body"))
    h.text(slide, p["expertise"],
           left=Inches(0.7), top=y + Pt(26),
           width=Inches(5.5), height=Inches(0.45),
           palette=palette, size_pt=10, color=palette["foreground"],
           font_name=fonts.get("body"))
    y = y + Inches(0.95)

# ── RIGHT half: equipment list with sections (hand-composed) ──────────
h.text(slide, "보유 장비 — 시퀀싱 · HPC · 스토리지",
       left=Inches(6.9), top=Inches(2.0),
       width=Inches(6.0), height=Pt(20),
       palette=palette, size_pt=13, bold=True,
       font_name=fonts.get("display"))
sections = [
    ("A. 시퀀싱", [
        ("ONT PromethION P24", "Long-read · 참조유전체"),
        ("MGI DNBSEQ-G99",     "Short-read · 6종 재서열"),
        ("DELL Precision 3680 + RTX", "Nanopore 베이스콜링"),
    ]),
    ("B. HPC (Flagship)", [
        ("ASUS ESC8000A-E12",     "512 cores · 1024 GiB RAM"),
        ("RTX 6000 · 6 노트",      "GATK GenomicsDB joint"),
    ]),
    ("C. 스토리지", [
        ("Synology RS3621",     "146 TiB · 6 노드"),
        ("DELL MD1200",          "90 TiB"),
        ("DELL MD1400",          "164 TiB"),
    ]),
]
y = Inches(2.4)
for title_text, items in sections:
    h.text(slide, title_text,
           left=Inches(6.9), top=y,
           width=Inches(5.8), height=Pt(16),
           palette=palette, size_pt=11, bold=True,
           color=palette["secondary"], font_name=fonts.get("display"))
    y = y + Pt(18)
    for name, spec in items:
        h.text(slide, name,
               left=Inches(7.05), top=y,
               width=Inches(2.6), height=Pt(16),
               palette=palette, size_pt=10, bold=True,
               font_name=fonts.get("body"))
        h.text(slide, spec,
               left=Inches(9.7), top=y,
               width=Inches(3.0), height=Pt(16),
               palette=palette, size_pt=10, color=palette["muted"],
               font_name=fonts.get("body"))
        y = y + Pt(15)
    y = y + Pt(6)
"""),
        False,
        "**Bespoke** — NOT a pattern call. ~60 shapes hand-composed: "
        "LEFT half is a 4-row personnel table (name + role + percentage "
        "+ expertise), RIGHT half is a multi-section equipment list "
        "(sequencing / HPC / storage). Use as the reference whenever "
        "the slide carries TWO independent dense compositions side-by-"
        "side (personnel + infrastructure, methodology + outputs, "
        "criteria + verification).",
        "Don't try to fit two independent tabular blocks into a single "
        "canned two-panel shape — when each side has its own internal "
        "multi-row layout, compose them bespoke side-by-side from h.* "
        "primitives.",
    ),
]


def _stage_tiles(tmp_dir: pathlib.Path) -> dict:
    """Solid-color PNG tiles standing in for real figures in the image-
    grid example (so the example renders without external assets)."""
    tiles = {
        "REF_TILE_RED":    (200, 90, 90),
        "REF_TILE_BLUE":   (90, 130, 180),
        "REF_TILE_GREEN":  (130, 170, 90),
        "REF_TILE_PURPLE": (160, 100, 170),
    }
    out = {}
    for name, color in tiles.items():
        p = tmp_dir / f"{name.lower()}.png"
        Image.new("RGB", (640, 480), color).save(p)
        out[name] = p
    return out


def main() -> int:
    corpus_dir = _THIS.parent
    tmp_dir = corpus_dir / "_tmp"
    tmp_dir.mkdir(exist_ok=True)
    tiles = _stage_tiles(tmp_dir)

    state = State(project_id="reference-corpus", owner_uid="me",
                  backend=InMemoryBackend())
    paper = papers.create_paper(state, title="Reference corpus")
    slug = paper["slug"]
    decks.create_deck(state, slug, title="corpus", deck_id="d")
    decks.update_deck(state, slug, "d", concept=(
        "Palette:\n"
        "  bg: #fafaf7  surface: #ffffff  text: #1a1a1a  accent: #b58900\n"
        "  muted: #6c757d  secondary: #2e7d32  highlight: #c0392b\n"
        "Typography:\n"
        "  display: Inter Bold  body: Inter Regular  mono: JetBrains Mono\n"
    ))

    # Substitute tile paths into the image-grid code
    manifest = {
        "generated_at": "2026-05-27",
        "generator": "packages/skills/paper-deck/reference_corpus/generate.py",
        "note": (
            "Body slides are bespoke by default — compose them directly "
            "from h.* primitives (see proposal_dense / personnel_equipment "
            "for the density vocabulary). The named patterns here are "
            "MECHANICAL scaffolds only: structural layouts (gantt, image "
            "grids, metric rows, pipelines, dividers, figure-only, title) "
            "that are tedious to hand-roll and content-neutral. Do NOT "
            "reach for a pattern to frame argument/evidence/comparison "
            "content — design that layout yourself. Agent flow: read "
            "manifest.json; if the slide is one of the mechanical "
            "structures below, Read its PNG and call the pattern; "
            "otherwise compose bespoke from h.*."
        ),
        "patterns": {},
    }
    for i, (fname, pattern, title, code, owns_slide, do, dont) in enumerate(REFERENCES, start=1):
        # Inline tile paths if needed
        for name, p in tiles.items():
            code = code.replace(name, repr(str(p)))
        decks.add_slide(
            state, slug, "d", slide_number=i, role="background",
            title=title, body="", notes="n", render_mode="code",
            code=code,
        )
        manifest["patterns"][pattern] = {
            "file": fname,
            "title": title,
            "owns_slide": owns_slide,
            "do": do,
            "dont": dont,
        }

    pptx_out = tmp_dir / "corpus.pptx"
    res = deck_render.export_deck_to_pptx(
        state, slug, "d", output_path=str(pptx_out),
    )
    if res["code_errors"]:
        print("code_errors:", res["code_errors"])
        return 1
    if res["pdf_skipped"]:
        print("PDF was skipped — soffice missing. Cannot render PNGs.")
        return 1

    # Move + rename PNGs from slide_001.png etc. → pattern-name.png
    # Also stamp shape_count on each manifest entry so the agent sees
    # density as a measurable property (todo 011 — agent reads the
    # manifest and learns "proposal_dense=~50 shapes, evidence_stack=
    # ~12 shapes" before authoring its own slide).
    from pptx import Presentation as _P
    pptx_doc = _P(str(pptx_out))
    for i, (fname, pattern, *_rest) in enumerate(REFERENCES, start=1):
        src = pathlib.Path(res["slide_pngs"][i - 1]["local_path"])
        dst = corpus_dir / fname
        shutil.copy(src, dst)
        # Total shape count on the corresponding slide (1-indexed →
        # zero-indexed in python-pptx)
        manifest["patterns"][pattern]["shape_count"] = len(
            pptx_doc.slides[i - 1].shapes
        )
    (corpus_dir / "manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    # Tidy temp
    shutil.rmtree(tmp_dir, ignore_errors=True)
    print(f"wrote {len(REFERENCES)} reference slides + manifest.json to "
          f"{corpus_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
