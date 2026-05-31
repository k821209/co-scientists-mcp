"""Structural slide layouts for `code` slides.

These are *mechanical* layouts only — content-neutral scaffolds whose
value is precise, tedious-to-hand-roll placement (timeline bar math,
aspect-fit image grids, evenly-spaced tiles). They impose NO narrative
or argument shape on the content.

Content slides — theses, evidence, comparisons, takeaways — are
authored BESPOKE: the agent designs the layout for THAT slide from
`h.*` primitives. There are deliberately no "thesis + evidence",
"before/after", "two competing options", etc. patterns: slotting
content into a canned frame is what made every deck look templated.
If you reach for a pattern to shape an *argument*, compose it from
primitives instead.

Patterns are bound as `p.*` in the `code` slide exec namespace. Each
function paints directly onto the slide and returns None (or a small
summary dict). All theme inputs (`palette`, `fonts`, `type_scale`,
`sw`, `sh`) follow the same calling convention as `h.*` helpers.

Layout catalog (mechanical scaffolds only):

  title_slide          — deck cover
  chapter_divider      — section opener
  metric_tile_row      — KPI / quantitative summary (evenly-spaced tiles)
  flow_pipeline        — workflow / process steps (arrow chain)
  gantt_chart          — timeline with period bars
  figure_full          — full-bleed aspect-fit figure
  title_and_image_grid — N images in a column grid
"""
from __future__ import annotations

from pptx.util import Inches, Pt          # type: ignore
from pptx.enum.shapes import MSO_SHAPE    # type: ignore
from pptx.enum.text import (              # type: ignore
    MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN,
)
from pptx.dml.color import RGBColor       # type: ignore

from . import slide_render_helpers as _h


# Where the body content starts for patterns that sit under an agent-
# drawn h.title_block (which spans ~Inches(0.45) .. Inches(1.5) plus a
# short accent rule below it). Inches(1.85) leaves a small clearance.
# Patterns that "own the whole slide" (e.g. chapter_divider) ignore this.
# (todo 005 Bug B / meta-issue: patterns need a contract about caller state.)
_BODY_TOP = Inches(1.85)
_SIDE_MARGIN = Inches(0.7)
_BOTTOM_MARGIN = Inches(0.6)


# ─── shared internals ────────────────────────────────────────────────────


def _muted(fg) -> RGBColor:
    """A muted version of the foreground color for secondary labels.
    We can't mix RGB without knowing intent, so we approximate: pull
    luminance halfway toward gray. Falls back to #6c757d."""
    try:
        r, g, b = fg
        mid = (r + g + b) // 3
        return RGBColor(
            int((r + mid * 2) // 3),
            int((g + mid * 2) // 3),
            int((b + mid * 2) // 3),
        )
    except Exception:
        return RGBColor(0x6C, 0x75, 0x7D)


def _emit_text(slide, text, *, left, top, width, height,
               size_pt: int, color, font_name=None, bold: bool = False,
               italic: bool = False, align=None, anchor=None,
               line_spacing: float = 1.2, autofit: bool = True,
               min_pt: int = 10):
    """Add a textbox with a single styled run.

    Autofits the font size to the box BEFORE rendering (Korean-aware
    width estimate) so the PNG preview matches the slide — soffice
    doesn't fully honor TEXT_TO_SHAPE auto-shrink, so we shrink up-
    front. Pass `autofit=False` to opt out (e.g. when the caller wants
    a guaranteed font size for visual hierarchy).
    """
    actual_pt = (
        _h.autofit_pt(text or "", max_width_emu=width,
                      max_height_emu=height,
                      start_pt=size_pt, line_spacing=line_spacing,
                      min_pt=min_pt)
        if autofit else size_pt
    )
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    _h._autoshrink(tf)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(actual_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return tb


def _resolve_items(name: str, items, **legacy_aliases):
    """Per todo 007 Tier 2 — every list-of-items parameter is canonically
    called `items`. Patterns that historically named the slot something
    else (evidence / tiles / steps / milestones) accept the old name as
    an alias. Raise TypeError loudly when both are passed (caller
    confusion) or neither (missing required input).

    Returns the resolved list.
    """
    explicit_legacy = {k: v for k, v in legacy_aliases.items()
                       if v is not None}
    if items is not None and explicit_legacy:
        raise TypeError(
            f"{name}: pass `items` OR one of "
            f"{sorted(legacy_aliases)} — got both"
        )
    if items is not None:
        return list(items)
    if len(explicit_legacy) == 1:
        return list(next(iter(explicit_legacy.values())))
    if len(explicit_legacy) > 1:
        raise TypeError(
            f"{name}: multiple legacy aliases passed — pick one of "
            f"{sorted(explicit_legacy)} or `items`"
        )
    raise TypeError(
        f"{name}() requires `items` (or legacy alias: "
        f"{', '.join(sorted(legacy_aliases))})"
    )


def _item_get(item, *keys, default=""):
    """For canonical-shape items: walk a tuple of key aliases and return
    the first non-None value. e.g. `_item_get(it, 'tag', 'title', 'value')`
    accepts `{tag: …}`, `{title: …}`, or `{value: …}` interchangeably so
    the same dict shape works across patterns with different historical
    naming."""
    if not isinstance(item, dict):
        return default
    for k in keys:
        v = item.get(k)
        if v is not None and v != "":
            return v
    return default


def _accent_rule(slide, *, left, top, width, height, color):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect


# ─── 1. chapter_divider ───────────────────────────────────────────────────


def chapter_divider(slide, *, chapter_label: str, summary: str = "",
                    palette, fonts, type_scale, sw, sh):
    """Section opener — Era I/II/III. Massive chapter label vertically
    centered + a short accent rule UNDER the label + optional summary
    tagline below. Distinct rhythm from interior slides — the audience
    feels the section break.

    Contract: **Owns the whole slide.** Do NOT call h.accent_stripe()
    or h.title_block() before this pattern. This pattern is the slide;
    the centred label is the only top-level element.

    Content limits:
        chapter_label: ≤ ~12 chars at chapter size (>= 56pt). Korean
            counts as 2× width per char.
        summary: ≤ ~50 chars; wraps to 2 lines max.
    """
    # Use a clearly chapter-sized type (>= 56pt) so the label reads as a
    # section break, not a slide title (todo 005 Bug C).
    label_pt = max(56, type_scale.get("cover_title", 40) + 16)
    summary_pt = type_scale.get("head", 26)

    fg = palette["foreground"]
    fg_muted = _muted(fg)
    accent = palette["accent"]
    margin_x = sw // 14
    box_w = sw - 2 * margin_x

    # Build a centred vertical block: label + rule + summary.
    label_h = Pt(label_pt * 1.25)
    rule_gap = Pt(_h.SPACING_UNIT_PT * 2)   # 16pt above + below the rule
    rule_h = Pt(8)
    summary_h = Pt(summary_pt * 2.4) if summary else 0
    block_h = label_h + rule_gap + rule_h + (rule_gap + summary_h if summary else 0)
    block_top = (sh - block_h) // 2

    _emit_text(slide, chapter_label,
               left=margin_x, top=block_top,
               width=box_w, height=label_h,
               size_pt=label_pt, color=fg,
               font_name=fonts.get("display"), bold=True,
               align=PP_ALIGN.CENTER, line_spacing=1.05)
    # Accent rule UNDER the chapter label — anchors the title visually.
    rule_w = max(Inches(1.2), sw // 8)
    _accent_rule(slide,
                 left=(sw - rule_w) // 2,
                 top=block_top + label_h + rule_gap,
                 width=rule_w, height=rule_h, color=accent)
    if summary:
        _emit_text(slide, summary,
                   left=margin_x,
                   top=block_top + label_h + rule_gap + rule_h + rule_gap,
                   width=box_w, height=summary_h,
                   size_pt=summary_pt, color=fg_muted,
                   font_name=fonts.get("body"),
                   italic=True, align=PP_ALIGN.CENTER,
                   line_spacing=1.25)


# ─── 2. metric_tile_row ───────────────────────────────────────────────────


def metric_tile_row(slide, *, items=None, tiles=None,
                    palette, fonts, type_scale, sw, sh,
                    top=None, height=None):
    """KPI / quantitative summary — a row of large numbers with thin
    labels (and optional units) under each.

    Contract: **Goes under a title_block.** Call h.accent_stripe + h.
    title_block BEFORE. Pattern starts at y ≈ `_BODY_TOP + Inches(0.5)`
    by default (use `top=` to override).

    Content limits:
        items (canonical) or tiles (legacy alias): 3–5 items, each in
            ONE of these shapes (auto-detected):
                - tuple `(value, label)` or `(value, label, unit)`
                - dict `{value, label, unit?}`
                - dict `{tag, body, unit?}` (canonical aligned shape —
                  tag = value, body = label)
            Label ≤ ~24 chars (1 line).
    """
    items = _resolve_items("metric_tile_row", items, tiles=tiles)
    if not items:
        return
    top = _BODY_TOP + Inches(0.5) if top is None else top
    height = Inches(3.0) if height is None else height
    side_margin = _SIDE_MARGIN
    usable = sw - 2 * side_margin
    gap = Pt(_h.SPACING_UNIT_PT * 3)
    tile_w = (usable - gap * (len(items) - 1)) // len(items)

    value_pt = max(48, type_scale.get("cover_title", 40) + 8)
    unit_pt = type_scale.get("head", 26)
    label_pt = type_scale.get("caption", 12)

    fg_muted = _muted(palette["foreground"])
    for i, tile in enumerate(items):
        # Shape detection: tuple-as-tile OR dict-as-tile (canonical)
        if isinstance(tile, (tuple, list)):
            if len(tile) == 3:
                value, label, unit = tile
            else:
                value, label = tile[0], tile[1]
                unit = ""
        else:
            value = _item_get(tile, "value", "tag")
            label = _item_get(tile, "label", "body")
            unit = _item_get(tile, "unit", default="")
        unit = unit or ""
        left = side_margin + i * (tile_w + gap)
        # Value (huge number, accent color)
        _emit_text(slide, str(value),
                   left=left, top=top,
                   width=tile_w, height=Pt(value_pt * 1.2),
                   size_pt=value_pt, color=palette["accent"],
                   font_name=fonts.get("display"), bold=True,
                   align=PP_ALIGN.CENTER, line_spacing=1.0)
        if unit:
            _emit_text(slide, unit,
                       left=left, top=top + Pt(value_pt * 1.3),
                       width=tile_w, height=Pt(unit_pt * 1.5),
                       size_pt=unit_pt, color=palette["foreground"],
                       font_name=fonts.get("body"),
                       align=PP_ALIGN.CENTER)
        # Label (small, muted)
        label_top = top + height - Pt(label_pt * 2.5)
        _emit_text(slide, label,
                   left=left, top=label_top,
                   width=tile_w, height=Pt(label_pt * 2.5),
                   size_pt=label_pt, color=fg_muted,
                   font_name=fonts.get("body"),
                   align=PP_ALIGN.CENTER)


# ─── 3. flow_pipeline ─────────────────────────────────────────────────────


def flow_pipeline(slide, *, items=None, steps=None,
                  palette, fonts, type_scale, sw, sh):
    """Horizontal process flow. Each step is a small card with a
    numbered tag + body; right-pointing arrows connect consecutive
    steps.

    Contract: **Goes under a title_block.** Call h.accent_stripe + h.
    title_block BEFORE. Pattern starts at `_BODY_TOP + 0.4"`.

    Content limits:
        items (canonical) or steps (legacy alias): 3–5 items, each
            `{tag (≤ ~12 chars), body (≤ ~80 chars, 3 lines)}`.
    """
    items = _resolve_items("flow_pipeline", items, steps=steps)
    if not items:
        return
    side = _SIDE_MARGIN
    top = _BODY_TOP + Inches(0.4)
    height = sh - top - _BOTTOM_MARGIN
    n = len(items)
    gap = Inches(0.3)
    arrow_w = Pt(_h.SPACING_UNIT_PT * 5)
    usable = sw - 2 * side - arrow_w * (n - 1) - gap * (n - 1)
    step_w = usable // n

    tag_pt = type_scale.get("head", 26) - 4
    body_pt = type_scale.get("body", 20) - 2
    fg_muted = _muted(palette["foreground"])

    for i, step in enumerate(items):
        left = side + i * (step_w + arrow_w + gap)
        # Step card
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, step_w, height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = palette["background"]
        card.line.color.rgb = palette["accent"]
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # Top accent stripe
        _accent_rule(slide, left=left, top=top,
                     width=step_w, height=Pt(4),
                     color=palette["accent"])
        # Step number
        _emit_text(slide, f"{i + 1:02d}",
                   left=left + Pt(12), top=top + Pt(12),
                   width=step_w - Pt(24), height=Pt(tag_pt * 1.2),
                   size_pt=int(tag_pt * 0.7),
                   color=palette["accent"],
                   font_name=fonts.get("display"), bold=True)
        # Tag
        _emit_text(slide, _item_get(step, "tag", "title"),
                   left=left + Pt(12), top=top + Pt(tag_pt * 1.6),
                   width=step_w - Pt(24), height=Pt(tag_pt * 1.5),
                   size_pt=tag_pt, color=palette["foreground"],
                   font_name=fonts.get("display"), bold=True)
        # Body
        _emit_text(slide, _item_get(step, "body", "note", "text"),
                   left=left + Pt(12), top=top + Pt(tag_pt * 3.4),
                   width=step_w - Pt(24),
                   height=height - Pt(tag_pt * 3.6),
                   size_pt=body_pt, color=fg_muted,
                   font_name=fonts.get("body"))
        # Arrow to the next step
        if i < n - 1:
            ax = left + step_w + gap // 2
            ay = top + height // 2 - Pt(8)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                ax, ay, arrow_w, Pt(16),
            )
            arrow.line.fill.background()
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = palette["accent"]
            arrow.shadow.inherit = False


# ─── gantt_chart (todo 009 D) ─────────────────────────────────────────────


def gantt_chart(slide, *, items=None, activities=None,
                periods: list[str] = None, period_count: int = None,
                palette, fonts, type_scale, sw, sh):
    """Horizontal Gantt chart — activity rows × period columns, with
    accent-colored bars positioned at each activity's start + span
    (todo 009 D). Use for project timelines, multi-month work plans,
    parallel-stage rollouts.

    Contract: **Goes under a title_block.** Pattern starts at
    `_BODY_TOP`. Pair with `h.deck_chrome` (eyebrow + page number)
    above for proposal-grade rhythm.

    Items (canonical, or legacy `activities=`):
        list[{label, start, span}] where `start` is 1-indexed period
        and `span` is the number of periods the bar covers.
        Example: {"label": "착수보고 및 계획 확정", "start": 1, "span": 1}

    Period labels: pass either `periods=["M1", "M2", …]` directly OR
    `period_count=8` (auto-labels M1..MN). Defaults to M1..M8.
    """
    items = _resolve_items("gantt_chart", items, activities=activities)
    if not items:
        return
    if periods is None:
        n = int(period_count) if period_count else 8
        periods = [f"M{i + 1}" for i in range(n)]
    n_periods = len(periods)

    body_top = _BODY_TOP
    body_left = _SIDE_MARGIN
    body_w = sw - 2 * _SIDE_MARGIN
    body_bottom = sh - _BOTTOM_MARGIN - Inches(0.3)
    body_h = body_bottom - body_top

    # Label column on the left = 32% of body width
    label_w = int(body_w * 0.32)
    timeline_left = body_left + label_w + Pt(8)
    timeline_w = body_w - label_w - Pt(8)
    col_w = timeline_w / n_periods

    # Header row height = ~1.6× label_pt
    label_pt = type_scale.get("label_tag", 12)
    body_pt = max(11, type_scale.get("body_small",
                                       type_scale.get("body", 20) - 4))
    header_h = Pt(label_pt * 2)
    rows_top = body_top + header_h + Pt(6)
    row_h = (body_bottom - rows_top) // max(1, len(items))

    fg = palette["foreground"]
    fg_muted = _muted(fg)
    accent = palette["accent"]

    # Period labels across the top
    for c, label in enumerate(periods):
        cx = timeline_left + int(col_w * c)
        _emit_text(slide, label,
                   left=cx, top=body_top,
                   width=int(col_w), height=header_h,
                   size_pt=label_pt, color=fg,
                   font_name=fonts.get("display"), bold=True,
                   align=PP_ALIGN.CENTER)

    # Activity rows
    for r, item in enumerate(items):
        label_text = _item_get(item, "label", "tag", "title", "text",
                                default="")
        start = int(item.get("start", 1)) if isinstance(item, dict) else 1
        span = int(item.get("span", 1)) if isinstance(item, dict) else 1
        start = max(1, min(n_periods, start))
        span = max(1, min(n_periods - start + 1, span))

        ry = rows_top + r * row_h
        # Zebra: alternate rows get a subtle muted-tint background
        if r % 2 == 1:
            zebra = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, body_left, ry, body_w, row_h,
            )
            zebra.line.fill.background()
            zebra.fill.solid()
            zebra.fill.fore_color.rgb = palette.get("surface", palette["background"])
            zebra.shadow.inherit = False
        # Label
        _emit_text(slide, label_text,
                   left=body_left + Pt(4), top=ry,
                   width=label_w - Pt(8), height=row_h,
                   size_pt=body_pt, color=fg,
                   font_name=fonts.get("body"),
                   anchor=MSO_ANCHOR.MIDDLE)
        # Bar
        bx = timeline_left + int(col_w * (start - 1)) + Pt(2)
        bw = int(col_w * span) - Pt(4)
        bar_h = max(Pt(10), int(row_h * 0.45))
        by = ry + (row_h - bar_h) // 2
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bar_h,
        )
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.shadow.inherit = False


# ─── figure_full (todo 008 §A) ────────────────────────────────────────────


def figure_full(slide, *, image_path: str = None, image_callable=None,
                caption: str = "",
                palette, fonts, type_scale, sw, sh):
    """Figure-only slide. The image fills the entire body area (rows
    1–6 of the 12-col grid — ~85% of the slide height); the caption
    rides in the bottom-margin strip OUTSIDE the grid, in muted body
    italics.

    Wins back ~17% of figure area vs the canonical `row_span=4` +
    caption-in-row-6 layout that earlier corpus exemplars used
    (todo 008 §A).

    Contract: **Goes under a title_block.** Call h.accent_stripe +
    h.title_block BEFORE. Image area starts at `_BODY_TOP`.

    Pass EXACTLY ONE image source:
        - `image_path` (str): a filesystem path. The pattern embeds
            via Keynote-safe normalization (RGBA→RGB JPEG ≤ 1920px).
        - `image_callable`: a callable that takes `slide, *, left,
            top, width, height` and adds the picture. Use this when
            embedding a paper figure or region from inside a code
            slide, e.g.
                p.figure_full(slide,
                    image_callable=lambda **kw: h.image_figure(slide, 3, **kw),
                    caption="Fig 3 · ...", palette=palette, fonts=fonts,
                    type_scale=type_scale, sw=sw, sh=sh)

    Content limits:
        caption: ≤ ~120 chars; 1 line preferred, 2 max.
    """
    if (image_path is None) == (image_callable is None):
        raise TypeError(
            "figure_full requires EXACTLY ONE of `image_path` or "
            "`image_callable`"
        )
    img_top = _BODY_TOP
    margin = _SIDE_MARGIN
    img_left = margin
    img_w = sw - 2 * margin
    cap_h = Inches(0.55) if caption else Inches(0.0)
    # Image owns the full grid + everything down to the caption strip.
    img_h = sh - img_top - cap_h - Inches(0.05)

    if image_path is not None:
        img_buf = _normalize_for_helper(image_path)
        pic = slide.shapes.add_picture(
            img_buf, img_left, img_top, width=img_w, height=img_h,
        )
        # Contain-fit (letterbox; never crop).
        pic_ar = (pic.width / pic.height) if pic.height else 1
        box_ar = (img_w / img_h) if img_h else 1
        if pic_ar > box_ar:
            new_w = img_w
            new_h = int(img_w / pic_ar)
        else:
            new_h = img_h
            new_w = int(img_h * pic_ar)
        pic.width = new_w
        pic.height = new_h
        pic.left = img_left + (img_w - new_w) // 2
        pic.top = img_top + (img_h - new_h) // 2
    else:
        # image_callable is responsible for embedding into the given box.
        image_callable(slide, left=img_left, top=img_top,
                       width=img_w, height=img_h)

    if caption:
        _emit_text(slide, caption,
                   left=margin, top=sh - Inches(0.5),
                   width=sw - 2 * margin, height=Pt(34),
                   size_pt=type_scale.get("caption", 12),
                   color=_muted(palette["foreground"]),
                   font_name=fonts.get("body"),
                   italic=True, align=PP_ALIGN.CENTER)


def _normalize_for_helper(path: str):
    """Local thin wrapper around deck_render._normalize_image_for_pptx
    — kept here to break a circular import (deck_render imports this
    module, and figure_full / title_and_image_grid need the
    normalization)."""
    from . import deck_render
    return deck_render._normalize_image_for_pptx(path)


# ─── Structural patterns (todo 006) ──────────────────────────────────────
#
# Where the 10 patterns above are organized by design *intent* (thesis,
# comparison, timeline, …), these are organized by content *shape* —
# the PowerPoint master-layout taxonomy. The agent picks the structural
# pattern first (do I have a body? images? how many?) and may then opt
# into a more design-loaded intent pattern instead.


def title_slide(slide, *, title: str, subtitle: str = "",
                eyebrow: str = "",
                palette, fonts, type_scale, sw, sh):
    """Deck opener / cover (PowerPoint base 1 — Title Slide). Centered
    eyebrow + large title + short accent rule + subtitle.

    Distinct from `chapter_divider` (a *mid-deck* section break) — this
    is the *opener* at slide 1 with author / venue / date metadata.

    Contract: **Owns the whole slide.** Do NOT call h.accent_stripe or
    h.title_block before; this pattern is the slide.

    Content limits:
        title:    ≤ ~50 chars at display size (wraps to 2 lines).
        subtitle: ≤ ~80 chars (1 line; e.g. "Author — Venue 2026").
        eyebrow:  ≤ ~30 chars; rendered upper-case in accent color.
    """
    title_pt = max(48, type_scale.get("cover_title", 40) + 8)
    sub_pt = type_scale.get("head", 26)
    eyebrow_pt = type_scale.get("caption", 12)

    fg = palette["foreground"]
    fg_muted = _muted(fg)
    accent = palette["accent"]
    margin_x = sw // 12
    box_w = sw - 2 * margin_x

    title_h = Pt(title_pt * 1.25)
    rule_gap = Pt(_h.SPACING_UNIT_PT * 2)
    rule_h = Pt(6)
    eyebrow_h = Pt(eyebrow_pt * 1.8) if eyebrow else 0
    sub_h = Pt(sub_pt * 1.8) if subtitle else 0
    block_h = (eyebrow_h + (Pt(4) if eyebrow else 0)
               + title_h
               + (rule_gap + rule_h if subtitle else 0)
               + (rule_gap + sub_h if subtitle else 0))
    block_top = (sh - block_h) // 2

    cursor = block_top
    if eyebrow:
        _emit_text(slide, eyebrow.upper(),
                   left=margin_x, top=cursor,
                   width=box_w, height=eyebrow_h,
                   size_pt=eyebrow_pt, color=accent,
                   font_name=fonts.get("body"), bold=True,
                   align=PP_ALIGN.CENTER)
        cursor += eyebrow_h + Pt(4)
    _emit_text(slide, title,
               left=margin_x, top=cursor,
               width=box_w, height=title_h,
               size_pt=title_pt, color=fg,
               font_name=fonts.get("display"), bold=True,
               align=PP_ALIGN.CENTER, line_spacing=1.05)
    cursor += title_h
    if subtitle:
        rule_w = max(Inches(1.2), sw // 8)
        _accent_rule(slide, left=(sw - rule_w) // 2,
                     top=cursor + rule_gap,
                     width=rule_w, height=rule_h, color=accent)
        _emit_text(slide, subtitle,
                   left=margin_x, top=cursor + rule_gap + rule_h + rule_gap,
                   width=box_w, height=sub_h,
                   size_pt=sub_pt, color=fg_muted,
                   font_name=fonts.get("body"),
                   align=PP_ALIGN.CENTER)


def title_and_image_grid(slide, *, title: str,
                          images: list[dict], cols: int = 2,
                          palette, fonts, type_scale, sw, sh):
    """N images in a `cols`-column grid (PowerPoint base 4 extended).
    1 image = full-half; 2 = side-by-side; 4 = 2×2; etc. Optional per-
    image caption strip below each tile. Layout effect (SKILL §5c):
    grid implies COMPARISON; for PROGRESSION use cols=N rows=1.

    Contract: **Goes under a title_block.** Pattern starts at _BODY_TOP.

    Content limits:
        images: 1–6 items, each {path (real filesystem path), caption?
            (≤ ~60 chars; 1 line)}.
        cols: 1, 2, 3, or 4 (default 2).
    """
    if not images:
        return
    cols = max(1, min(4, int(cols)))
    n = len(images)
    rows = (n + cols - 1) // cols

    side = _SIDE_MARGIN
    top = _BODY_TOP
    gap = Pt(_h.SPACING_UNIT_PT * 2)
    grid_w = sw - 2 * side - gap * (cols - 1)
    grid_h = sh - top - _BOTTOM_MARGIN - gap * (rows - 1)
    cell_w = grid_w // cols
    cell_h = grid_h // rows

    cap_pt = type_scale.get("caption", 12)
    fg_muted = _muted(palette["foreground"])

    for i, img in enumerate(images):
        r, c = i // cols, i % cols
        x = side + c * (cell_w + gap)
        y = top + r * (cell_h + gap)
        path = img.get("path")
        caption = (img.get("caption") or "").strip()
        cap_h = Pt(cap_pt * 1.8) if caption else 0
        img_h = cell_h - cap_h - (Pt(4) if caption else 0)
        # Embed image with Keynote-safe normalization
        if path:
            img_buf = _normalize_for_helper(path)
            pic = slide.shapes.add_picture(img_buf, x, y, width=cell_w, height=img_h)
            # Letterbox-fit (contain) — never crop a tile image.
            if pic.width != cell_w or pic.height != img_h:
                # contain math
                pic_ar = pic.width / pic.height if pic.height else 1
                box_ar = cell_w / img_h if img_h else 1
                if pic_ar > box_ar:
                    new_w = cell_w
                    new_h = int(cell_w / pic_ar)
                else:
                    new_h = img_h
                    new_w = int(img_h * pic_ar)
                pic.width = new_w
                pic.height = new_h
                pic.left = x + (cell_w - new_w) // 2
                pic.top = y + (img_h - new_h) // 2
        if caption:
            _emit_text(slide, caption,
                       left=x, top=y + img_h + Pt(4),
                       width=cell_w, height=cap_h,
                       size_pt=cap_pt, color=fg_muted,
                       font_name=fonts.get("body"),
                       italic=True, align=PP_ALIGN.CENTER)
