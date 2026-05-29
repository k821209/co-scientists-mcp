"""Deck rendering + PPTX export (Phase 3).

`render_slide` materializes a slide's image based on its render_mode:
  - paper-figure : copies the existing figure blob from the paper
  - ai-image     : substitutes concept placeholders in prompt, calls
                   the configured ImageGenerator
  - code-shape   : agent supplies a local PNG path (the MCP can't safely
                   sandbox arbitrary Python execution); we just register
                   the path
  - hybrid       : treated as code-shape for now

The rendered PNG is uploaded to Storage at
  papers/{slug}/decks/{deck_id}/slides/{slide_number}.png
and the slide doc's `image_blob_path` field is updated.

`export_deck_to_pptx` emits a .pptx via python-pptx:
  - image slides   : the rendered PNG, aspect-fitted and centered
  - `text` slides  : a NATIVE (editable) title + bullet layout, themed
                     from the deck concept's palette — not a picture
The slide size follows the deck's `aspect_ratio`. A sibling .pdf is
produced via LibreOffice (`soffice`) when available — python-pptx's
PPTX is occasionally rejected by Keynote, so the PDF is the portable
fallback. python-pptx ships in the base install.
"""
from __future__ import annotations

import os
import pathlib
import re
import subprocess
import tempfile

from ..backends.base import NotFound
from ..state import State
from ..util import now_iso
from . import decks as _decks
from . import figures as _figures


# ─── placeholder resolution ─────────────────────────────────────────────────


_KV_RE = re.compile(r"^\s*([a-zA-Z_]\w*)\s*:\s*(.+?)\s*$")


def _parse_concept(concept: str | None) -> dict[str, str]:
    """Extract {key: value} pairs from a freeform concept string.

    The skill writes the concept like:
        Palette:
          bg: #fafaf7  surface: #ffffff  text: #1a1a1a  accent: #b58900
        Typography:
          display: Inter Bold     body: Inter Regular

    We tokenize on whitespace+colon to harvest any `name: value` pairs the
    skill drops in. Best-effort; placeholders that don't resolve stay as
    `{name}` in the output prompt for visibility.
    """
    if not concept:
        return {}
    pairs: dict[str, str] = {}
    # Split on whitespace tokens like "accent: #b58900" — also handle
    # multi-word values when the line is one-per-key.
    for raw in concept.splitlines():
        line = raw.strip()
        if not line or line.endswith(":"):
            continue
        # Multiple pairs on a single line ("bg: #fff surface: #f00 text: #000")
        chunks = re.findall(r"([a-zA-Z_]\w*)\s*:\s*([^\s][^,\n]*?)(?=\s+[a-zA-Z_]\w*\s*:|$)", line)
        if chunks:
            for k, v in chunks:
                pairs[k.lower()] = v.strip()
            continue
        m = _KV_RE.match(line)
        if m:
            pairs[m.group(1).lower()] = m.group(2).strip()
    return pairs


def resolve_placeholders(text: str, concept: str | None) -> str:
    """Replace `{name}` tokens with values harvested from the concept.

    Unknown placeholders stay literal so the issue is visible (a slide
    with `{accent}` in its final prompt is a bug the user can see).
    """
    if not text:
        return text
    if "{" not in text:
        return text
    kv = _parse_concept(concept)

    def _sub(m: re.Match) -> str:
        key = m.group(1).lower()
        return kv.get(key, m.group(0))
    return re.sub(r"\{([a-zA-Z_][\w-]*)\}", _sub, text)


def _apply_image_style(prompt: str, deck: dict) -> str:
    """Prepend `deck.image_style` (if set) to an ai-image prompt so every
    region in the deck inherits a consistent visual treatment without
    the agent having to repeat the style hint per slide. No-op when
    image_style is empty / missing. Style goes FIRST so the rest of the
    prompt anchors the subject — gpt-image-2 weights early tokens.
    """
    style = (deck.get("image_style") or "").strip()
    if not style:
        return prompt
    body = (prompt or "").strip()
    if not body:
        return style
    return f"{style}. {body}"


# ─── slide rendering ────────────────────────────────────────────────────────

# Deck aspect ratio → (width_in, height_in) for the PPTX page.
_ASPECT_TO_SIZE = {
    "16:9": (13.333, 7.5),
    "16:10": (12.0, 7.5),
    "4:3": (10.0, 7.5),
}

# Standard aspect ratios an image generator accepts — used to match a
# region's box shape (the original co-scientist's todo 052: pass the
# region aspect to the generator, not the deck aspect).
_STD_ASPECTS = {"1:1": 1.0, "4:3": 4 / 3, "3:4": 3 / 4, "16:9": 16 / 9, "9:16": 9 / 16}


def _nearest_aspect(width: float, height: float) -> str:
    """Pick the standard aspect-ratio string closest to width:height."""
    if height <= 0:
        return "1:1"
    ratio = width / height
    return min(_STD_ASPECTS, key=lambda k: abs(_STD_ASPECTS[k] - ratio))


def _deck_aspect_size(deck: dict) -> tuple[float, float]:
    """The deck's slide dimensions in inches."""
    return _ASPECT_TO_SIZE.get(
        deck.get("aspect_ratio") or "16:9", _ASPECT_TO_SIZE["16:9"],
    )


def _slide_blob_path(state: State, slug: str, deck_id: str, slide_number: int) -> str:
    return state.project_path(
        "papers", slug, "decks", deck_id, "slides", f"{slide_number:04d}.png",
    )


def _region_blob_path(
    state: State, slug: str, deck_id: str, slide_number: int, region_id: str,
) -> str:
    return state.project_path(
        "papers", slug, "decks", deck_id, "slides",
        f"{slide_number:04d}", "regions", f"{region_id}.png",
    )


def _figure_png(state: State, slug: str, figure_number) -> bytes:
    """Fetch a paper figure's image bytes, or raise NotFound."""
    fig = _figures.get_figure(state, slug, int(figure_number))
    if not fig or not fig.get("blob_path"):
        raise NotFound(
            f"figure {figure_number!r} has no blob_path on paper {slug!r}"
        )
    blob = state.backend.get_blob(fig["blob_path"])
    if blob is None:
        raise NotFound(f"figure {figure_number!r} blob is empty")
    return blob


def _png_size(data: bytes) -> tuple[int, int] | None:
    """(width, height) in pixels from a PNG header, or None if not a PNG."""
    if (len(data) >= 24 and data[:8] == b"\x89PNG\r\n\x1a\n"
            and data[12:16] == b"IHDR"):
        import struct
        return struct.unpack(">II", data[16:24])
    return None


def _image_dims(data: bytes) -> dict:
    """`{image_width, image_height}` for a PNG, else empty — stored on the
    placeholder so layout can reconcile the box against the real image."""
    sz = _png_size(data)
    return {"image_width": sz[0], "image_height": sz[1]} if sz else {}


# Tokens that signal a `code` field is an in-place python-pptx snippet
# (renders at export time via _add_code_slide) rather than a `code-shape`
# external script (produces a PNG path the agent supplies via local_path).
# Used by `_infer_render_mode` to disambiguate (todo 010).
_PYPPTX_SNIPPET_SIGNALS = (
    "slide.shapes", "slide.shapes.add", "h.text", "h.title_block",
    "h.accent_stripe", "h.bullet_list", "h.card", "h.image_",
    "h.icon", "h.grid", "h.table", "h.deck_chrome",
    "p.title_slide", "p.title_and_body", "p.chapter_divider",
    "p.hero_with_trailing_evidence", "p.evidence_stack",
    "p.flow_pipeline", "p.metric_tile_row", "p.before_after_split",
    "p.contrast_pair", "p.quadrant_map", "p.numbered_milestone_arc",
    "p.zoom_in_callout", "p.figure_full", "p.gantt_chart",
    "p.title_two_content", "p.title_and_image_grid",
    "MSO_SHAPE.", "Inches(", "Pt(",
)


def _infer_render_mode(slide: dict) -> str:
    """Pick a render mode from which fields the agent populated (todo
    010 — defer the design decision to authoring time). Priority:

      regions[]                    → "hybrid"
      code has python-pptx signal  → "code"
      figure_number                → "paper-figure"
      image_blob_path              → "ai-image" (already-rendered image)
      prompt (no blob yet)         → "ai-image" (will be generated)
      code (no python-pptx signal) → "code-shape" (external PNG script)
      otherwise                    → "text"

    Explicit `render_mode` on the slide always wins — call this only
    when `slide.get("render_mode")` is None / missing.
    """
    if slide.get("regions"):
        return "hybrid"
    code = (slide.get("code") or "").strip()
    if code and any(s in code for s in _PYPPTX_SNIPPET_SIGNALS):
        return "code"
    if slide.get("figure_number") is not None:
        return "paper-figure"
    if slide.get("image_blob_path"):
        return "ai-image"
    if (slide.get("prompt") or "").strip():
        return "ai-image"
    if code:
        return "code-shape"
    return "text"


def _resolve_mode(slide: dict) -> str:
    """Explicit `render_mode` if set; otherwise infer from populated
    fields (todo 010)."""
    explicit = slide.get("render_mode")
    if explicit:
        return explicit
    return _infer_render_mode(slide)


def render_slide(
    state: State,
    slug: str,
    deck_id: str,
    slide_id: str,
    *,
    local_path: str | None = None,
) -> dict:
    """Materialize a slide's image(s) and upload to Storage.

    paper-figure / ai-image : the MCP renders it.
    code-shape              : the agent supplies `local_path` (a PNG it
                              produced locally — the MCP won't exec code).
    hybrid                  : renders every region it can; returns a
                              per-region summary (see render_region).
    text                    : nothing to render — native text at export.
    """
    deck = _decks.get_deck(state, slug, deck_id)
    slide = state.backend.get_doc(_decks._slide_path(state, slug, deck_id, slide_id))
    if slide is None:
        raise NotFound(f"slide not found: {slide_id!r}")
    mode = _resolve_mode(slide)
    slide_number = slide.get("slide_number") or 0

    if mode == "text":
        raise ValueError(
            "text slides carry no image — they render as native PPTX "
            "text on export; nothing to render here"
        )
    if mode == "code":
        # Code slide with regions = image-placeholder layout. Render
        # the regions same as hybrid. Code without regions has nothing
        # to materialize ahead of export.
        if slide.get("regions"):
            return _render_hybrid_slide(state, slug, deck_id, slide, deck)
        raise ValueError(
            "code slide has no regions — its `code` runs at PPTX export "
            "time and populates the slide natively. To pre-render image "
            "placeholders, call set_slide_regions on this slide first."
        )
    if mode == "hybrid":
        return _render_hybrid_slide(state, slug, deck_id, slide, deck)

    png: bytes
    if mode == "paper-figure":
        fig_num = slide.get("figure_number")
        if fig_num is None:
            raise ValueError("paper-figure slide has no figure_number")
        png = _figure_png(state, slug, fig_num)
    elif mode == "ai-image":
        prompt = resolve_placeholders(slide.get("prompt") or "", deck.get("concept"))
        prompt = _apply_image_style(prompt, deck)
        if not prompt.strip():
            raise ValueError("ai-image slide has no prompt")
        png = state.require_image_gen().generate(prompt=prompt, aspect_ratio="16:9")
    elif mode == "code-shape":
        if not local_path:
            raise ValueError(
                "code-shape slide requires local_path — produce the PNG "
                "yourself (matplotlib, python-pptx shapes, etc.) and pass it in."
            )
        png = pathlib.Path(local_path).expanduser().read_bytes()
    else:
        raise ValueError(f"unknown render_mode: {mode!r}")

    blob_path = _slide_blob_path(state, slug, deck_id, slide_number)
    state.backend.put_blob(blob_path, png)
    now = now_iso()
    state.backend.update_doc(
        _decks._slide_path(state, slug, deck_id, slide_id),
        {
            "image_blob_path": blob_path,
            "rendered_at": now,
            "status": "rendered",
            "updated_at": now,
            **_image_dims(png),
        },
    )
    return {
        "slide_id": slide_id,
        "slide_number": slide_number,
        "blob_path": blob_path,
        "size_bytes": len(png),
        "mode": mode,
    }


# ─── hybrid (multi-region) slides ───────────────────────────────────────────


def _update_region(
    state: State, slug: str, deck_id: str, slide_id: str,
    region_id: str, fields: dict,
) -> None:
    """Patch one region (matched by id) inside a slide's `regions` array."""
    path = _decks._slide_path(state, slug, deck_id, slide_id)
    slide = state.backend.get_doc(path)
    regions = list(slide.get("regions") or [])
    for r in regions:
        if r.get("id") == region_id:
            r.update(fields)
            break
    state.backend.update_doc(path, {"regions": regions, "updated_at": now_iso()})


def _render_one_region(
    state: State, slug: str, deck_id: str, slide: dict, deck: dict,
    region: dict, *, local_path: str | None = None,
) -> dict:
    """Render a single region's image, store it, stamp the region doc."""
    rmode = region.get("render_mode")
    region_id = region.get("id")
    slide_number = slide.get("slide_number") or 0

    if rmode == "paper-figure":
        if region.get("figure_number") is None:
            raise ValueError(f"region {region_id}: paper-figure needs figure_number")
        png = _figure_png(state, slug, region["figure_number"])
    elif rmode == "ai-image":
        prompt = resolve_placeholders(region.get("prompt") or "", deck.get("concept"))
        prompt = _apply_image_style(prompt, deck)
        if not prompt.strip():
            raise ValueError(f"region {region_id}: ai-image needs a prompt")
        w_in, h_in = _deck_aspect_size(deck)
        aspect = _nearest_aspect(
            (region.get("w") or 1.0) * w_in, (region.get("h") or 1.0) * h_in,
        )
        png = state.require_image_gen().generate(prompt=prompt, aspect_ratio=aspect)
    elif rmode == "code-shape":
        if not local_path:
            raise ValueError(
                f"region {region_id}: code-shape needs local_path — render the "
                "PNG yourself and pass it to render_region."
            )
        png = pathlib.Path(local_path).expanduser().read_bytes()
    else:
        raise ValueError(f"region {region_id}: unknown render_mode {rmode!r}")

    blob_path = _region_blob_path(state, slug, deck_id, slide_number, region_id)
    state.backend.put_blob(blob_path, png)
    _update_region(state, slug, deck_id, slide["id"], region_id, {
        "image_blob_path": blob_path,
        "rendered_at": now_iso(),
        **_image_dims(png),
    })
    return {
        "region_id": region_id,
        "mode": rmode,
        "blob_path": blob_path,
        "size_bytes": len(png),
    }


def _render_hybrid_slide(
    state: State, slug: str, deck_id: str, slide: dict, deck: dict,
) -> dict:
    """Render every auto-renderable region of a hybrid slide. code-shape
    regions land in `skipped[]` for the agent to do via render_region."""
    regions = slide.get("regions") or []
    if not regions:
        raise ValueError(
            "hybrid slide has no regions — call set_slide_regions first"
        )
    rendered: list[dict] = []
    skipped: list[dict] = []
    errors: list[dict] = []
    for r in regions:
        tag = {
            "slide_id": slide["id"], "slide_number": slide.get("slide_number"),
            "region_id": r.get("id"),
        }
        if r.get("render_mode") == "code-shape":
            skipped.append({
                **tag,
                "reason": "code-shape region — render_region with local_path",
            })
            continue
        try:
            rendered.append({**tag, **_render_one_region(
                state, slug, deck_id, slide, deck, r,
            )})
        except Exception as e:
            errors.append({**tag, "error": str(e)})
    return {
        "slide_id": slide["id"],
        "slide_number": slide.get("slide_number"),
        "mode": "hybrid",
        "rendered": rendered,
        "skipped": skipped,
        "errors": errors,
    }


def render_region(
    state: State,
    slug: str,
    deck_id: str,
    slide_id: str,
    region_id: str,
    *,
    local_path: str | None = None,
) -> dict:
    """Render one region of a hybrid slide. paper-figure / ai-image are
    done by the MCP; a code-shape region needs `local_path` — a PNG you
    produced locally."""
    slide = state.backend.get_doc(_decks._slide_path(state, slug, deck_id, slide_id))
    if slide is None:
        raise NotFound(f"slide not found: {slide_id!r}")
    deck = _decks.get_deck(state, slug, deck_id)
    region = next(
        (r for r in (slide.get("regions") or []) if r.get("id") == region_id),
        None,
    )
    if region is None:
        raise NotFound(f"region not found: {region_id!r} on slide {slide_id!r}")
    return _render_one_region(
        state, slug, deck_id, slide, deck, region, local_path=local_path,
    )


def _slide_is_rendered(s: dict) -> bool:
    """A slide is 'done' when its image(s) exist: text → always; code
    with regions → every region has an image; code without regions →
    always (materializes natively at export); hybrid → every region has
    an image; else → has image_blob_path."""
    mode = _resolve_mode(s)
    if mode == "text":
        return True
    if mode == "code":
        regions = s.get("regions") or []
        # No regions → nothing to pre-render. With regions → every
        # placeholder needs an image before export so h.image_region
        # resolves.
        return all(r.get("image_blob_path") for r in regions)
    if mode == "hybrid":
        regions = s.get("regions") or []
        return bool(regions) and all(r.get("image_blob_path") for r in regions)
    return bool(s.get("image_blob_path"))


def render_deck(state: State, slug: str, deck_id: str) -> dict:
    """Render every slide we can do automatically. code-shape slides and
    code-shape regions land in `skipped[]` for agent follow-up. When every
    non-text slide (and every region) is rendered, deck.status → 'rendered'."""
    slides = _decks.list_slides(state, slug, deck_id)
    results: dict[str, list] = {"rendered": [], "skipped": [], "errors": []}
    for s in slides:
        mode = _resolve_mode(s)
        tag = {"slide_id": s["id"], "slide_number": s.get("slide_number")}
        if mode == "text":
            results["skipped"].append({
                **tag, "reason": "text slide — native PPTX text on export",
            })
            continue
        if mode == "code":
            # Code slide may carry image placeholders (regions[]) the
            # snippet will pull in via h.image_region(). Render those
            # regions ahead of export. No regions → nothing to do.
            if s.get("regions"):
                try:
                    hr = _render_hybrid_slide(
                        state, slug, deck_id, s,
                        _decks.get_deck(state, slug, deck_id),
                    )
                    results["rendered"].extend(hr["rendered"])
                    results["skipped"].extend(hr["skipped"])
                    results["errors"].extend(hr["errors"])
                except Exception as e:
                    results["errors"].append({**tag, "error": str(e)})
            else:
                results["skipped"].append({
                    **tag, "reason": "code slide — runs at PPTX export time",
                })
            continue
        if mode == "code-shape":
            results["skipped"].append({
                **tag, "reason": "needs local PNG from agent",
            })
            continue
        if mode == "hybrid":
            try:
                hr = render_slide(state, slug, deck_id, s["id"])
                results["rendered"].extend(hr["rendered"])
                results["skipped"].extend(hr["skipped"])
                results["errors"].extend(hr["errors"])
            except Exception as e:
                results["errors"].append({**tag, "error": str(e)})
            continue
        try:
            results["rendered"].append(render_slide(state, slug, deck_id, s["id"]))
        except Exception as e:
            results["errors"].append({**tag, "error": str(e)})

    refreshed = _decks.list_slides(state, slug, deck_id)
    if all(_slide_is_rendered(s) for s in refreshed):
        _decks.update_deck(state, slug, deck_id, status="rendered")
    return results


# ─── PPTX export ─────────────────────────────────────────────────────────────


def _theme_colors(concept: str | None) -> dict[str, str]:
    """Harvest every documented Palette key as hex (todo 007 axis 4).
    Concept's `Palette:` block exposes 7 named colors; we propagate all
    of them, computing sensible defaults for any missing keys so the
    snippet doesn't have to hardcode hex literals when it wants e.g.
    `palette["muted"]`."""
    kv = _parse_concept(concept)
    accent = kv.get("accent") or "#2E7D32"
    background = kv.get("bg") or kv.get("background") or "#FFFFFF"
    foreground = kv.get("text") or kv.get("foreground") or "#1A1A1A"
    surface = kv.get("surface") or background
    # `muted` defaults to a 65/35 blend of foreground+background (a
    # legible secondary-text gray that adapts to dark/light themes).
    # `secondary` falls back to a darker accent shift; `highlight` to a
    # warmer accent shift — both rough heuristics that the agent can
    # override via the concept's Palette block when precision matters.
    muted = kv.get("muted") or _blend_hex(foreground, background, 0.45)
    secondary = kv.get("secondary") or _shift_hex(accent, -0.18)
    highlight = kv.get("highlight") or _shift_hex(accent, +0.18)
    return {
        "accent": accent,
        "background": background,
        "foreground": foreground,
        "surface": surface,
        "muted": muted,
        "secondary": secondary,
        "highlight": highlight,
    }


def _hex_to_tuple(hexstr: str) -> tuple[int, int, int]:
    raw = (hexstr or "").lstrip("#")
    if len(raw) != 6:
        raw = "000000"
    try:
        return (int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except ValueError:
        return (0, 0, 0)


def _tuple_to_hex(rgb: tuple[int, int, int]) -> str:
    r, g, b = (max(0, min(255, int(c))) for c in rgb)
    return f"#{r:02x}{g:02x}{b:02x}"


def _blend_hex(a: str, b: str, t: float) -> str:
    """Linear interpolation between two hex colors. t=0 → all a, t=1 → all b."""
    ar, ag, ab = _hex_to_tuple(a)
    br, bg, bb = _hex_to_tuple(b)
    return _tuple_to_hex((
        ar * (1 - t) + br * t,
        ag * (1 - t) + bg * t,
        ab * (1 - t) + bb * t,
    ))


def _shift_hex(hexstr: str, delta: float) -> str:
    """Shift each channel of a hex color by `delta` (-1..+1) toward
    white (delta > 0) or black (delta < 0). Cheap accent-variant
    generator for `secondary`/`highlight` defaults."""
    target = "#FFFFFF" if delta > 0 else "#000000"
    return _blend_hex(hexstr, target, abs(delta))


def _hex_to_rgb(value: str, fallback: str):
    """Parse '#rrggbb' into a python-pptx RGBColor, tolerating junk."""
    from pptx.dml.color import RGBColor  # type: ignore
    raw = (value or "").lstrip("#")
    if len(raw) != 6:
        raw = fallback.lstrip("#")
    try:
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except ValueError:
        fb = fallback.lstrip("#")
        return RGBColor(int(fb[0:2], 16), int(fb[2:4], 16), int(fb[4:6], 16))


def _rel_luminance(hexstr: str) -> float:
    """WCAG relative luminance of an '#rrggbb' color."""
    raw = (hexstr or "").lstrip("#")
    if len(raw) != 6:
        raw = "000000"

    def _chan(v: int) -> float:
        c = v / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    try:
        r, g, b = (int(raw[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return 0.0
    return 0.2126 * _chan(r) + 0.7152 * _chan(g) + 0.0722 * _chan(b)


def _contrast_ratio(hex1: str, hex2: str) -> float:
    """WCAG contrast ratio between two colors (1.0 … 21.0)."""
    lo, hi = sorted((_rel_luminance(hex1), _rel_luminance(hex2)))
    return (hi + 0.05) / (lo + 0.05)


def _ensure_contrast(fg_hex: str, bg_hex: str) -> str:
    """Guarantee legible text: if `fg` on `bg` fails WCAG AA (4.5:1),
    snap it to black or white — whichever reads better on `bg`."""
    if _contrast_ratio(fg_hex, bg_hex) >= 4.5:
        return fg_hex
    return ("#000000" if _contrast_ratio("#000000", bg_hex)
            >= _contrast_ratio("#FFFFFF", bg_hex) else "#FFFFFF")


def _normalize_image_for_pptx(img_path: str, target_width: int = 1920):
    """RGBA / palette PNGs and oversized images → RGB JPEG ≤ target_width px,
    returned as a BytesIO. Mirrors the original co-scientist's todo 032
    normalization: Keynote rejects python-pptx output when slides embed RGBA
    or huge PNGs, so every picture goes through this guard before
    `add_picture`. If PIL can't process the source (corrupt, exotic format,
    Pillow missing), returns the original path so embedding still works."""
    try:
        import io
        from PIL import Image  # type: ignore — python-pptx pulls this in
        img = Image.open(img_path)
        img.load()  # surface truncation errors here, not on .split()
        if img.mode in ("RGBA", "LA", "P"):
            bg = Image.new("RGB", img.size, "white")
            if img.mode == "RGBA":
                bg.paste(img, mask=img.split()[-1])
            else:
                rgba = img.convert("RGBA")
                bg.paste(rgba, mask=rgba.split()[-1])
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")
        if img.width > target_width:
            new_h = round(img.height * target_width / img.width)
            img = img.resize((target_width, new_h), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=92, optimize=True)
        buf.seek(0)
        return buf
    except Exception:
        # Last-ditch fallback — Keynote-safety is best-effort. Embedding the
        # original is still strictly better than failing the export.
        return img_path


def _place_picture(slide, img_path: str, *, left, top, box_w, box_h,
                   fit: str = "contain") -> None:
    """Place an image inside the box at (left, top, box_w, box_h).

    fit='contain' — scaled to fit entirely, centered, letterboxed; never
    crops (figures, charts, tables). fit='cover' — fills the box, the
    overflowing edges cropped (eyecatch, decorative, photos).

    Every image is normalized to Keynote-safe RGB JPEG ≤ 1920px (todo 032)
    before embedding."""
    pic = slide.shapes.add_picture(_normalize_image_for_pptx(img_path), left, top)
    iw, ih = pic.width, pic.height
    if fit == "cover" and iw > 0 and ih > 0:
        image_ar, box_ar = iw / ih, box_w / box_h
        if image_ar > box_ar:            # too wide → crop left/right
            crop = (1 - box_ar / image_ar) / 2
            pic.crop_left = pic.crop_right = crop
        elif image_ar < box_ar:          # too tall → crop top/bottom
            crop = (1 - image_ar / box_ar) / 2
            pic.crop_top = pic.crop_bottom = crop
        pic.left, pic.top = int(left), int(top)
        pic.width, pic.height = int(box_w), int(box_h)
    else:                                # contain
        scale = min(box_w / iw, box_h / ih)
        pic.width = int(iw * scale)
        pic.height = int(ih * scale)
        pic.left = int(left + (box_w - pic.width) / 2)
        pic.top = int(top + (box_h - pic.height) / 2)


# Typography for native text. Presentation-scale but tuned a notch down
# from "shouty" — readable from a room without dominating dense content.
# Defaults; the deck concept can override any via a `Type scale:` block.
_BODY_PT = 20
_HEAD_PT = 26
_TITLE_PT = 32
_LINE_SPACING = 1.22

_DEFAULT_TYPE_SCALE = {
    # ─── legacy keys (kept for back-compat with existing callers) ────
    "title": _TITLE_PT,        # content-slide title
    "head": _HEAD_PT,          # body markdown # heading
    "body": _BODY_PT,          # body bullets / paragraphs
    "hybrid_body": 18,         # narrower body box → smaller body type
    "hybrid_head": 22,
    "cover_title": 40,         # title slide hero (kept; new code: display_cover)
    "cover_subtitle": 20,
    "caption": 12,             # region captions
    "line_spacing": _LINE_SPACING,
    # ─── semantic role keys (todo 004 §E — canonical naming) ─────────
    # New code is encouraged to use these by *role* instead of by
    # arbitrary pixel size. Old keys above stay as informal aliases.
    "display_cover":    48,    # cover slide title (slide 1)
    "display_hero":     44,    # thesis / takeaway hero line
    "display_chapter":  56,    # mid-deck chapter divider
    "headline_section": 28,    # in-slide section header
    "title_slide":      _TITLE_PT,   # standard slide title
    "body_large":       22,    # generous body
    "body_standard":    _BODY_PT,    # standard body
    "body_small":       16,    # small body / dense lists
    "label_tag":        12,    # pills, tags
    "label_caption":    12,    # figure captions
    "scale_ratio":      1.25,  # perfect-fourth ratio (relationship constant)
}


def _theme_type_scale(concept: str | None) -> dict:
    """Per-deck typography sizes — title / head / body / cover_* / caption /
    line_spacing. The concept can override any with a 'Type scale:' block:

        Type scale:
          title: 30  head: 24  body: 18  line_spacing: 1.2

    Unspecified keys fall back to `_DEFAULT_TYPE_SCALE`. Mirrors the
    original co-scientist's todo 035 (theme JSON `type_scale`).
    """
    kv = _parse_concept(concept)
    out = dict(_DEFAULT_TYPE_SCALE)
    # Float-valued keys; everything else (point sizes) coerces to int.
    float_keys = {"line_spacing", "scale_ratio"}
    for k in out:
        if k in kv:
            try:
                out[k] = (
                    float(kv[k]) if k in float_keys else int(float(kv[k]))
                )
            except (ValueError, TypeError):
                pass  # malformed value → keep default
    return out

_WEIGHT_WORDS = {
    "thin", "extralight", "ultralight", "light", "regular", "book", "normal",
    "medium", "semibold", "demibold", "bold", "extrabold", "ultrabold",
    "black", "heavy", "italic", "oblique",
}


def _font_family(value: str | None) -> str | None:
    """'Inter Bold' → 'Inter' — drop trailing weight/style words so the
    name is a font *family* PowerPoint can resolve."""
    if not value or not value.strip():
        return None
    words = value.strip().split()
    while len(words) > 1 and words[-1].lower() in _WEIGHT_WORDS:
        words.pop()
    return " ".join(words) or None


def _theme_fonts(concept: str | None) -> dict[str, str | None]:
    """Display / body / mono font families from the concept's Typography
    block. None when unspecified — PowerPoint's default is then kept."""
    kv = _parse_concept(concept)
    body = _font_family(kv.get("body") or kv.get("body_font"))
    display = _font_family(
        kv.get("display") or kv.get("heading") or kv.get("title_font")
    )
    mono = _font_family(kv.get("mono") or kv.get("code_font"))
    return {"display": display or body, "body": body or display, "mono": mono}


def _render_simple_body(slide, body: str, *, box, fg, fonts, Pt,
                        body_pt: int = _BODY_PT,
                        line_spacing: float = _LINE_SPACING) -> None:
    """Plain-text body. Each non-empty line becomes one paragraph at
    `body_pt` (Korean-aware autofit may shrink further to fit the box),
    with any leading bullet/star marker stripped. **No markdown
    parsing** — see docs/todo/002."""
    from . import slide_render_helpers as _h
    left, top, w, h = box
    # Strip leading markers and gather lines so the autofit calc sees
    # the same text that ends up rendered.
    lines = [raw.strip().lstrip("-•*").lstrip()
             for raw in body.splitlines() if raw.strip()]
    if not lines:
        return
    # Autofit considers the entire body's wrap + line count, then sets
    # one font size for every paragraph in the textbox so the rhythm
    # stays even.
    actual_pt = _h.autofit_pt(
        "\n".join(lines),
        max_width_emu=w, max_height_emu=h,
        start_pt=body_pt, line_spacing=line_spacing,
        min_pt=max(12, body_pt - 6),
    )
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        from pptx.enum.text import MSO_AUTO_SIZE  # type: ignore
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
    except Exception:
        pass
    first = True
    for ln in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = ln
        run.font.size = Pt(actual_pt)
        run.font.color.rgb = fg
        if fonts.get("body"):
            run.font.name = fonts["body"]



def _add_slide_frame(slide, row, *, sw, sh, accent, fg, bg, fonts,
                     Inches, Pt, MSO_SHAPE,
                     type_scale: dict = _DEFAULT_TYPE_SCALE) -> None:
    """Themed background + accent stripe + native title — the shared shell
    of native text slides and hybrid (multi-region) slides."""
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
    except Exception:
        pass  # leave the default background if this build won't set it
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.16))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.shadow.inherit = False

    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.45), sw - Inches(1.4), Inches(1.0),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    # Auto-shrink the title when it would overflow the box (long titles).
    try:
        from pptx.enum.text import MSO_AUTO_SIZE  # type: ignore
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    run = p.add_run()
    run.text = row.get("title") or ""
    run.font.size = Pt(type_scale["title"])
    run.font.bold = True
    run.font.color.rgb = fg
    if fonts.get("display"):
        run.font.name = fonts["display"]

    # Short accent rule under the title — a visual anchor for the slide.
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.52), Inches(2.2), Pt(4),
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = accent
    rule.shadow.inherit = False


def _add_text_slide(slide, row, *, sw, sh, accent, fg, bg, fonts,
                    Inches, Pt, MSO_SHAPE,
                    type_scale: dict = _DEFAULT_TYPE_SCALE) -> None:
    """Themed title frame + plain-text body bullets. The body is
    rendered straight (no markdown parsing) — slides that need richer
    visual treatment should set a `code` field with a python-pptx
    render function (see docs/todo/002)."""
    _add_slide_frame(slide, row, sw=sw, sh=sh, accent=accent, fg=fg, bg=bg,
                     fonts=fonts, Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                     type_scale=type_scale)
    body = (row.get("body") or "").strip()
    if body:
        _render_simple_body(
            slide, body,
            box=(Inches(0.7), Inches(1.95),
                 sw - Inches(1.4), sh - Inches(2.5)),
            fg=fg, fonts=fonts, Pt=Pt,
            body_pt=type_scale["body"],
            line_spacing=type_scale["line_spacing"],
        )


def _add_title_slide(slide, row, *, sw, sh, accent, fg, bg, fonts,
                     Inches, Pt, MSO_SHAPE, PP_ALIGN, MSO_ANCHOR,
                     type_scale: dict = _DEFAULT_TYPE_SCALE) -> None:
    """A cover layout for role='title' slides — the title large and
    centered (vertically + horizontally), the body as a centered
    subtitle. Distinct from the top-anchored content-slide frame."""
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
    except Exception:
        pass
    box = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.5), sw - Inches(2.0), sh - Inches(1.0),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Auto-shrink for very long cover titles.
    try:
        from pptx.enum.text import MSO_AUTO_SIZE  # type: ignore
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.1
    run = p.add_run()
    run.text = row.get("title") or ""
    run.font.size = Pt(type_scale["cover_title"])
    run.font.bold = True
    run.font.color.rgb = fg
    if fonts.get("display"):
        run.font.name = fonts["display"]

    body = (row.get("body") or "").strip()
    first_sub = True
    for raw in body.splitlines():
        line = raw.strip().lstrip("#-*•> ").strip()
        if not line:
            continue
        sp = tf.add_paragraph()
        sp.alignment = PP_ALIGN.CENTER
        sp.line_spacing = 1.2
        if first_sub:
            sp.space_before = Pt(18)
            first_sub = False
        run = sp.add_run()
        run.text = line
        run.font.size = Pt(type_scale["cover_subtitle"])
        run.font.color.rgb = fg
        if fonts.get("body"):
            run.font.name = fonts["body"]


def _add_hybrid_slide(slide, row, state, tmpd, *, sw, sh, accent, fg, bg,
                      fonts, Inches, Pt, MSO_SHAPE,
                      type_scale: dict = _DEFAULT_TYPE_SCALE) -> int:
    """Themed title frame + native body bullets (left half, when body is
    set) + each region placed as its own picture. The body-on-the-left
    convention gives the "title + bullets + figure-on-the-right" layout
    natively; agents position image regions (typically on the right)
    accordingly. Returns the count of unrendered regions."""
    _add_slide_frame(slide, row, sw=sw, sh=sh, accent=accent, fg=fg, bg=bg,
                     fonts=fonts, Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                     type_scale=type_scale)
    body = (row.get("body") or "").strip()
    if body:
        # Half-width box wraps more, so smaller body type than a text slide.
        _render_simple_body(
            slide, body,
            box=(Inches(0.7), Inches(1.65),
                 int(sw / 2) - Inches(0.8), sh - Inches(1.95)),
            fg=fg, fonts=fonts, Pt=Pt,
            body_pt=type_scale["hybrid_body"],
            line_spacing=type_scale["line_spacing"],
        )
    unrendered = 0
    for r in row.get("regions") or []:
        left = int(r["x"] * sw)
        top = int(r["y"] * sh)
        box_w = int(r["w"] * sw)
        box_h = int(r["h"] * sh)
        caption = r.get("caption")
        cap_h = Inches(0.34) if caption else 0
        img_h = box_h - cap_h
        blob = (
            state.backend.get_blob(r["image_blob_path"])
            if r.get("image_blob_path") else None
        )
        if blob:
            tmp = pathlib.Path(tmpd) / f"region_{row['id']}_{r['id']}.png"
            tmp.write_bytes(blob)
            _place_picture(slide, str(tmp),
                           left=left, top=top, box_w=box_w, box_h=img_h,
                           fit=r.get("fit") or "contain")
        else:
            unrendered += 1
            box = slide.shapes.add_textbox(left, top, box_w, img_h)
            box.text_frame.word_wrap = True
            box.text_frame.text = (
                f"[region {r.get('id')}: {r.get('render_mode')} — not rendered]"
            )
        if caption:
            cap = slide.shapes.add_textbox(left, top + img_h, box_w, cap_h)
            cf = cap.text_frame
            cf.word_wrap = True
            run = cf.paragraphs[0].add_run()
            run.text = caption
            run.font.size = Pt(type_scale["caption"])
            run.font.color.rgb = fg
            if fonts.get("body"):
                run.font.name = fonts["body"]
    return unrendered


# ─── code slides: agent-authored python-pptx snippets ───────────────────


def _build_code_namespace(slide, row, state, slug, tmpd, *,
                          palette, fonts, type_scale, aspect, sw, sh):
    """Build the namespace handed to a code slide's `exec()`. Pre-binds
    the slide object, theme primitives, python-pptx imports, and the
    `h` helpers namespace — including image-loading closures that
    capture `state` so the snippet doesn't have to thread it through.

    The snippet runs as top-level statements (not a function). The
    convention: `slide`, `title`, `body`, `row`, `palette`, `fonts`,
    `type_scale`, `aspect`, `sw`, `sh`, plus python-pptx (`Pt`,
    `Inches`, `MSO_SHAPE`, `PP_ALIGN`, `MSO_ANCHOR`, `RGBColor`), plus
    `h` for the helpers. See docs/todo/002 + paper-deck SKILL §5a.
    """
    from pptx.util import Pt, Inches, Emu                  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE                 # type: ignore
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR        # type: ignore
    from pptx.dml.color import RGBColor                    # type: ignore
    from types import SimpleNamespace
    from . import slide_render_helpers as _h
    from . import slide_patterns as _p

    # Image helpers accept slide as the first positional arg for
    # consistency with every other helper / pattern (todo 007 axis 1).
    # Calls without slide still work via the closure capture — that
    # keeps the older `h.image_path("/path", ...)` form running.

    def _peel_slide(args, want_one_after: bool):
        """Resolve a varargs call shape: (target,) or (slide, target).
        Returns (target_slide, target). Falls back to the closure-
        captured `slide` when the caller omitted it."""
        if len(args) == (2 if want_one_after else 1):
            target_slide, *rest = args
            if rest:
                return target_slide, rest[0]
            return target_slide, None
        if len(args) == (1 if want_one_after else 0):
            return (slide, args[0]) if want_one_after else (slide, None)
        raise TypeError(
            "expected (slide, value) or (value,); got "
            f"{len(args)} positional args"
        )

    def _image_from_path(*args, left, top, width, height, fit="contain"):
        _, path = _peel_slide(args, want_one_after=True)
        _place_picture(slide, str(path),
                       left=left, top=top, box_w=width, box_h=height, fit=fit)

    def _image_from_region(*args, left, top, width, height, fit="contain"):
        _, region_id = _peel_slide(args, want_one_after=True)
        region = next(
            (r for r in (row.get("regions") or [])
             if r.get("id") == region_id), None,
        )
        if region is None:
            raise ValueError(f"region {region_id!r} not on this slide")
        bp = region.get("image_blob_path")
        if not bp:
            raise ValueError(
                f"region {region_id!r} has no image_blob_path — "
                "render it via render_region first"
            )
        blob = state.backend.get_blob(bp)
        if not blob:
            raise ValueError(f"region {region_id!r} blob is empty")
        tmp = pathlib.Path(tmpd) / f"region_{row['id']}_{region_id}.png"
        tmp.write_bytes(blob)
        _place_picture(slide, str(tmp),
                       left=left, top=top, box_w=width, box_h=height, fit=fit)

    def _image_from_figure(*args, left, top, width, height, fit="contain"):
        _, figure_number = _peel_slide(args, want_one_after=True)
        png = _figure_png(state, slug, int(figure_number))
        tmp = pathlib.Path(tmpd) / f"figure_{figure_number}_{row['id']}.png"
        tmp.write_bytes(png)
        _place_picture(slide, str(tmp),
                       left=left, top=top, box_w=width, box_h=height, fit=fit)

    # Combine static helpers + state-aware image helpers into one `h`.
    h = SimpleNamespace(
        accent_stripe=_h.accent_stripe,
        title_block=_h.title_block,
        bullet_list=_h.bullet_list,
        card=_h.card,
        card_grid=_h.card_grid,
        pull_quote=_h.pull_quote,
        image_path=_image_from_path,
        image_region=_image_from_region,
        image_figure=_image_from_figure,
        grid=_h.grid,                       # todo 004 §D — design grid
        SPACING_UNIT_PT=_h.SPACING_UNIT_PT,  # 8pt vertical rhythm
        icon=_h.icon,                       # todo 004 §C — iconography
        icon_names=_h.icon_names,           # list available icon names
        autofit_pt=_h.autofit_pt,           # Korean-aware text autofit
        estimate_text_width_pt=_h.estimate_text_width_pt,
        text=_h.text,                       # todo 007 §D — one-call textbox
        Cell=_h.Cell,                       # Grid.cell() return type
        deck_chrome=_h.deck_chrome,         # todo 009 B — eyebrow/footer/page
        table=_h.table,                     # todo 009 C — native pptx table
    )
    # Whole-slide patterns (todo 004 §B, 006) — bound as `p`.
    p = SimpleNamespace(
        # Intent-axis patterns (todo 004 §B)
        hero_with_trailing_evidence=_p.hero_with_trailing_evidence,
        chapter_divider=_p.chapter_divider,
        metric_tile_row=_p.metric_tile_row,
        evidence_stack=_p.evidence_stack,
        flow_pipeline=_p.flow_pipeline,
        before_after_split=_p.before_after_split,
        contrast_pair=_p.contrast_pair,
        quadrant_map=_p.quadrant_map,
        numbered_milestone_arc=_p.numbered_milestone_arc,
        zoom_in_callout=_p.zoom_in_callout,
        # Structural-axis patterns (todo 006 — PowerPoint master layouts)
        title_slide=_p.title_slide,
        title_and_body=_p.title_and_body,
        title_two_content=_p.title_two_content,
        title_and_image_grid=_p.title_and_image_grid,
        figure_full=_p.figure_full,         # todo 008 §A
        gantt_chart=_p.gantt_chart,         # todo 009 D
    )

    return {
        # Slide + content
        "slide": slide,
        "title": row.get("title") or "",
        "body": row.get("body") or "",
        "notes": row.get("notes") or "",
        "row": row,
        # Theme
        "palette": palette,
        "fonts": fonts,
        "type_scale": type_scale,
        "aspect": aspect,
        # Canvas
        "sw": sw, "sh": sh,
        # python-pptx imports the snippet will reach for
        "Pt": Pt, "Inches": Inches, "Emu": Emu,
        "MSO_SHAPE": MSO_SHAPE,
        "PP_ALIGN": PP_ALIGN,
        "MSO_ANCHOR": MSO_ANCHOR,
        "RGBColor": RGBColor,
        # Helpers (full + short alias)
        "h": h, "helpers": h,
        # Whole-slide patterns
        "p": p, "patterns": p,
    }


def _detect_text_overlaps(slide, *, min_overlap_ratio: float = 0.2) -> list[dict]:
    """Walk `slide.shapes` and find pairs of textbox shapes whose
    bounding boxes intersect by at least `min_overlap_ratio` of the
    SMALLER box's area (todo 016).

    Only textbox↔textbox pairs count: card border rectangles, accent
    rules, dividers, and other decorative shapes lack `text_frame` so
    they're filtered out automatically — that's intentional, since a
    card border rect always "overlaps" the title textbox inside it
    and we don't want to flag those.

    Returns a list of `{a_preview, b_preview, overlap_ratio}` dicts.
    Empty list = clean slide.
    """
    boxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        try:
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
        except AttributeError:
            continue
        if w <= 0 or h <= 0:
            continue
        preview = text.splitlines()[0][:40]
        boxes.append(((l, t, l + w, t + h), preview))

    out: list[dict] = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            (l1, t1, r1, b1), p1 = boxes[i]
            (l2, t2, r2, b2), p2 = boxes[j]
            ix1, iy1 = max(l1, l2), max(t1, t2)
            ix2, iy2 = min(r1, r2), min(b1, b2)
            if ix1 >= ix2 or iy1 >= iy2:
                continue
            inter = (ix2 - ix1) * (iy2 - iy1)
            a1 = (r1 - l1) * (b1 - t1)
            a2 = (r2 - l2) * (b2 - t2)
            min_area = min(a1, a2) or 1
            ratio = inter / min_area
            if ratio >= min_overlap_ratio:
                out.append({
                    "a_preview": p1,
                    "b_preview": p2,
                    "overlap_ratio": round(ratio, 2),
                })
    return out


def _add_code_slide(slide, row, state, slug, tmpd, *, sw, sh,
                    accent, fg, bg, fonts, palette_full, Inches, Pt,
                    MSO_SHAPE,
                    type_scale: dict = _DEFAULT_TYPE_SCALE) -> tuple[str | None, list[dict]]:
    """Execute the slide's `code` against a prepared namespace.

    Returns `(err, overlap_warnings)`. `err` is the error message
    string when execution raises (the caller then degrades to plain
    text), or None on success. `overlap_warnings` is the list of
    text-on-text overlap pairs found AFTER successful execution —
    empty when no overlap or when execution failed (no shapes to
    check).
    """
    code = (row.get("code") or "").strip()
    if not code:
        return "no code on slide", []
    # Full 7-key palette as RGBColor objects (todo 007 axis 4).
    palette = {
        "accent": accent, "background": bg, "foreground": fg,
        "surface": _hex_to_rgb(palette_full["surface"], "#FFFFFF"),
        "muted": _hex_to_rgb(palette_full["muted"], "#6C757D"),
        "secondary": _hex_to_rgb(palette_full["secondary"], "#2E7D32"),
        "highlight": _hex_to_rgb(palette_full["highlight"], "#FFC107"),
    }
    aspect = "16:9"  # only used as a hint inside the snippet
    ns = _build_code_namespace(
        slide, row, state, slug, tmpd,
        palette=palette, fonts=fonts, type_scale=type_scale,
        aspect=aspect, sw=sw, sh=sh,
    )
    try:
        exec(compile(code, f"<slide {row.get('id', '?')} code>", "exec"), ns)
    except Exception as e:  # noqa: BLE001 — surface to the caller
        return f"{type(e).__name__}: {e}", []
    warnings = _detect_text_overlaps(slide)
    return None, warnings


def _render_pdf_to_pngs(pdf_path: pathlib.Path, out_dir: pathlib.Path,
                        *, dpi: int = 150) -> list[pathlib.Path]:
    """Render each page of `pdf_path` to a PNG at out_dir/slide_{N:03d}.png
    via PyMuPDF (todo 004 §A). Per-slide PNGs let the agent run a vision
    critique loop — Read each PNG, score against a design rubric, rewrite
    weak slides' `code`, re-export. Best-effort: returns [] when PyMuPDF
    is missing or the PDF can't be opened, so the rest of the export
    still succeeds."""
    try:
        import pymupdf  # type: ignore
    except ImportError:
        return []
    try:
        doc = pymupdf.open(str(pdf_path))
    except Exception:
        return []
    mat = pymupdf.Matrix(dpi / 72, dpi / 72)
    pngs: list[pathlib.Path] = []
    try:
        for i in range(len(doc)):
            page = doc[i]
            pix = page.get_pixmap(matrix=mat, alpha=False)
            out = out_dir / f"slide_{i + 1:03d}.png"
            pix.save(str(out))
            pngs.append(out)
    except Exception:
        pass  # partial results returned (best-effort)
    finally:
        doc.close()
    return pngs


def _pdf_via_soffice(pptx_path: pathlib.Path) -> pathlib.Path | None:
    """Convert the .pptx to a sibling .pdf via LibreOffice. Returns the PDF
    path, or None if soffice/libreoffice is missing or the conversion fails.
    python-pptx output is occasionally rejected by Keynote — the PDF is the
    portable fallback."""
    for binary in ("soffice", "libreoffice"):
        try:
            proc = subprocess.run(
                [binary, "--headless", "--convert-to", "pdf",
                 "--outdir", str(pptx_path.parent), str(pptx_path)],
                capture_output=True, text=True, timeout=180,
            )
        except FileNotFoundError:
            continue  # try the next binary name
        except subprocess.TimeoutExpired:
            return None
        pdf = pptx_path.with_suffix(".pdf")
        return pdf if (proc.returncode == 0 and pdf.is_file()) else None
    return None


def export_deck_to_pptx(
    state: State,
    slug: str,
    deck_id: str,
    *,
    output_path: str,
) -> dict:
    """Build a .pptx — and, when LibreOffice is available, a sibling .pdf.

    Per slide render_mode:
      - image slides (paper-figure / ai-image / code-shape) embed the
        rendered PNG, aspect-fitted and centered.
      - `hybrid` slides get a themed title frame plus each region placed
        as its own picture at its x/y/w/h box.
      - `text` slides — and any slide still missing a render — become
        NATIVE editable text (title + bullets) themed from the deck
        concept's palette.
    The slide size follows the deck's `aspect_ratio`.
    """
    try:
        from pptx import Presentation                       # type: ignore
        from pptx.util import Inches, Pt                    # type: ignore
        from pptx.enum.shapes import MSO_SHAPE              # type: ignore
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR    # type: ignore
    except ImportError as e:
        raise RuntimeError(
            "python-pptx not installed — reinstall the package: "
            "pip install -e ~/co-scientists-mcp/apps/local-mcp"
        ) from e

    deck = _decks.get_deck(state, slug, deck_id)
    slides = _decks.list_slides(state, slug, deck_id)
    if not slides:
        raise ValueError(f"deck {deck_id!r} has no slides")

    aspect = deck.get("aspect_ratio") or "16:9"
    w_in, h_in = _ASPECT_TO_SIZE.get(aspect, _ASPECT_TO_SIZE["16:9"])
    colors = _theme_colors(deck.get("concept"))
    # Guarantee the body text is legible even if the concept's palette
    # pairs low-contrast colors.
    fg_hex = _ensure_contrast(colors["foreground"], colors["background"])
    accent = _hex_to_rgb(colors["accent"], "#2E7D32")
    fg = _hex_to_rgb(fg_hex, "#1A1A1A")
    bg = _hex_to_rgb(colors["background"], "#FFFFFF")
    fonts = _theme_fonts(deck.get("concept"))
    type_scale = _theme_type_scale(deck.get("concept"))

    prs = Presentation()
    sw, sh = Inches(w_in), Inches(h_in)
    prs.slide_width, prs.slide_height = sw, sh
    blank_layout = prs.slide_layouts[6]

    missing_renders: list[int] = []
    code_errors: list[dict] = []
    overlap_warnings: list[dict] = []
    image_slides = 0
    text_slides = 0
    hybrid_slides = 0
    code_slides = 0
    with tempfile.TemporaryDirectory() as tmpd:
        for s in slides:
            slide = prs.slides.add_slide(blank_layout)
            mode = _resolve_mode(s)
            if mode == "code":
                # Agent-authored python-pptx snippet builds the slide
                # natively (docs/todo/002). On exec failure we degrade
                # to plain text so the deck still exports — the error
                # surfaces in `code_errors[]` of the result.
                err, overlaps = _add_code_slide(
                    slide, s, state, slug, tmpd, sw=sw, sh=sh,
                    accent=accent, fg=fg, bg=bg, fonts=fonts,
                    palette_full=colors,
                    Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                    type_scale=type_scale,
                )
                if overlaps:
                    overlap_warnings.append({
                        "slide_number": s.get("slide_number"),
                        "slide_id": s.get("id"),
                        "pairs": overlaps,
                    })
                if err is None:
                    code_slides += 1
                else:
                    code_errors.append({
                        "slide_number": s.get("slide_number"),
                        "slide_id": s.get("id"),
                        "error": err,
                    })
                    # Wipe the (possibly partial) slide and rebuild as text.
                    # python-pptx doesn't expose a clean remove-slide API
                    # without rewriting the XML, so we re-add as a fresh
                    # slide at the end of the deck order. Simpler: just
                    # populate the existing slide via _add_text_slide;
                    # any partial shapes from the failed exec will still
                    # be there, but they don't crash the export.
                    _add_text_slide(
                        slide, s, sw=sw, sh=sh,
                        accent=accent, fg=fg, bg=bg, fonts=fonts,
                        Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                        type_scale=type_scale,
                    )
                    text_slides += 1
            elif mode == "hybrid":
                n_missing = _add_hybrid_slide(
                    slide, s, state, tmpd, sw=sw, sh=sh,
                    accent=accent, fg=fg, bg=bg, fonts=fonts,
                    Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                    type_scale=type_scale,
                )
                hybrid_slides += 1
                if n_missing:
                    missing_renders.append(s.get("slide_number") or 0)
            else:
                blob = (
                    state.backend.get_blob(s["image_blob_path"])
                    if s.get("image_blob_path") else None
                )
                if blob:
                    tmp = pathlib.Path(tmpd) / f"slide_{s['id']}.png"
                    tmp.write_bytes(blob)
                    # Full-bleed single image — always contain (a bare
                    # figure must never be cropped to fill the slide).
                    _place_picture(slide, str(tmp),
                                   left=0, top=0, box_w=sw, box_h=sh)
                    image_slides += 1
                elif (s.get("role") or "") == "title":
                    # Cover slide — centered hero layout, not the
                    # top-anchored content frame.
                    _add_title_slide(
                        slide, s, sw=sw, sh=sh, accent=accent, fg=fg, bg=bg,
                        fonts=fonts, Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                        PP_ALIGN=PP_ALIGN, MSO_ANCHOR=MSO_ANCHOR,
                        type_scale=type_scale,
                    )
                    text_slides += 1
                    if mode != "text":
                        missing_renders.append(s.get("slide_number") or 0)
                else:
                    # Native editable text slide — a `text` slide by design,
                    # or a not-yet-rendered slide degraded gracefully.
                    _add_text_slide(
                        slide, s, sw=sw, sh=sh, accent=accent, fg=fg, bg=bg,
                        fonts=fonts, Inches=Inches, Pt=Pt, MSO_SHAPE=MSO_SHAPE,
                        type_scale=type_scale,
                    )
                    text_slides += 1
                    if mode != "text":
                        missing_renders.append(s.get("slide_number") or 0)
            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s["notes"]

    out = pathlib.Path(output_path).expanduser()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    pdf_path = _pdf_via_soffice(out)

    # Upload + index each artifact so the dashboard's Presentation tab
    # lists them.
    now = now_iso()

    def _publish(path: pathlib.Path, fmt: str) -> str:
        blob_path = state.project_path(
            "papers", slug, "decks", deck_id, "exports", path.name,
        )
        data = path.read_bytes()
        state.backend.put_blob(blob_path, data)
        state.backend.set_doc(blob_path, {
            "filename": path.name,
            "format": fmt,
            "blob_path": blob_path,
            "size_bytes": len(data),
            "slide_count": len(slides),
            "missing_renders": missing_renders,
            "created_at": now,
        })
        return blob_path

    pptx_blob = _publish(out, "pptx")
    pdf_blob = _publish(pdf_path, "pdf") if pdf_path else None

    # Per-slide PNGs for the agent's vision critique loop (todo 004 §A).
    # Needs the PDF as source; skipped silently when soffice/PyMuPDF
    # missing — same posture as the PDF itself.
    slide_pngs: list[dict] = []
    if pdf_path is not None:
        pngs = _render_pdf_to_pngs(pdf_path, out.parent)
        for i, p in enumerate(pngs, start=1):
            blob_path = state.project_path(
                "papers", slug, "decks", deck_id, "exports", p.name,
            )
            state.backend.put_blob(blob_path, p.read_bytes())
            slide_pngs.append({
                "slide_number": i,
                "local_path": str(p),
                "blob_path": blob_path,
            })

    return {
        "deck_id": deck_id,
        "deck_title": deck.get("title"),
        "aspect_ratio": aspect,
        "slide_count": len(slides),
        "image_slides": image_slides,
        "text_slides": text_slides,
        "hybrid_slides": hybrid_slides,
        "code_slides": code_slides,
        "missing_renders": missing_renders,
        "code_errors": code_errors,
        "overlap_warnings": overlap_warnings,
        "local_path": str(out),
        "blob_path": pptx_blob,
        "pdf_local_path": str(pdf_path) if pdf_path else None,
        "pdf_blob_path": pdf_blob,
        "pdf_skipped": pdf_path is None,
        "size_bytes": os.path.getsize(out),
        "slide_pngs": slide_pngs,
        "slide_pngs_skipped": pdf_path is None or not slide_pngs,
    }
