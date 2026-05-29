"""Helpers exposed to a slide's `code` field at PPTX export time
(docs/todo/002).

A slide with `render_mode='code'` carries a `code` field that is a
python-pptx snippet. The exporter `exec`s that code in a namespace where
the slide object, theme primitives, python-pptx imports, and these
helpers are pre-bound — so the snippet can compose the slide natively
without dragging in boilerplate.

The static helpers in this module don't need the runtime state (Storage
/ Firestore). Image-loading helpers that DO need state are constructed
per-slide as closures in `deck_render._build_code_namespace` and added
to the same `h` namespace.

Design intent: keep these helpers thin. They cover the boilerplate the
agent would otherwise repeat per slide (accent stripe, themed title,
bullet column, card / card-grid), so per-slide code can focus on the
layout that distinguishes that slide.

Usage shape inside a slide's `code` field:

    h.accent_stripe(slide)
    h.title_block(slide, title)
    h.card_grid(slide, [
        {"title": "Memory", "body": "decisions stack across sessions"},
        {"title": "Hooks",  "body": "run on specific events"},
        {"title": "Slash",  "body": "reusable commands"},
        {"title": "Context","body": "what the LLM is looking at"},
    ], left=Inches(0.7), top=Inches(1.9),
       width=sw - Inches(1.4), height=sh - Inches(2.3),
       cols=2)
"""
from __future__ import annotations

from typing import NamedTuple

from pptx.util import Inches, Pt  # type: ignore
from pptx.enum.shapes import MSO_SHAPE  # type: ignore
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore


class Cell(NamedTuple):
    """Geometry of a grid cell. Tuple-iterates as `(left, top, width,
    height)` so existing code can unpack it; attribute-accessible as
    `.left / .top / .width / .height` for modern style (todo 007 axis
    3). Both forms are first-class."""
    left: int
    top: int
    width: int
    height: int


# Icon vocabulary (todo 004 §C). Each semantic name maps to either an
# MSO_SHAPE built-in (preferred — native PowerPoint shape, fully
# recolorable + editable + scalable) or a fallback Unicode glyph
# rendered in a textbox. The agent calls `h.icon(slide, name, ...)`
# with one of these names; unknown names raise.
_ICON_SHAPES = {
    # Direction / flow
    "arrow-right":   MSO_SHAPE.RIGHT_ARROW,
    "arrow-left":    MSO_SHAPE.LEFT_ARROW,
    "arrow-up":      MSO_SHAPE.UP_ARROW,
    "arrow-down":    MSO_SHAPE.DOWN_ARROW,
    "arrow-both":    MSO_SHAPE.LEFT_RIGHT_ARROW,
    "arrow-vert":    MSO_SHAPE.UP_DOWN_ARROW,
    "arrow-quad":    MSO_SHAPE.QUAD_ARROW,
    "arrow-curve":   MSO_SHAPE.CURVED_RIGHT_ARROW,
    "arrow-loop":    MSO_SHAPE.CIRCULAR_ARROW,
    "chevron":       MSO_SHAPE.CHEVRON,
    # Status / emphasis
    "lightning":     MSO_SHAPE.LIGHTNING_BOLT,
    "burst":         MSO_SHAPE.EXPLOSION1,
    "burst-big":     MSO_SHAPE.EXPLOSION2,
    "star":          MSO_SHAPE.STAR_5_POINT,
    "heart":         MSO_SHAPE.HEART,
    "plus":          MSO_SHAPE.MATH_PLUS,
    "minus":         MSO_SHAPE.MATH_MINUS,
    "warning":       MSO_SHAPE.ISOSCELES_TRIANGLE,
    # Data / storage
    "database":      MSO_SHAPE.CAN,
    "cloud":         MSO_SHAPE.CLOUD,
    "document":      MSO_SHAPE.FOLDED_CORNER,
    # Process / decision
    "decision":      MSO_SHAPE.DIAMOND,
    "molecule":      MSO_SHAPE.HEXAGON,
    "stop":          MSO_SHAPE.OCTAGON,
    "input":         MSO_SHAPE.PARALLELOGRAM,
    "gear":          MSO_SHAPE.GEAR_6,
    # Time / mode
    "sun":           MSO_SHAPE.SUN,
    "moon":          MSO_SHAPE.MOON,
    # Generic geometry
    "circle":        MSO_SHAPE.OVAL,
    "square":        MSO_SHAPE.RECTANGLE,
    "rounded":       MSO_SHAPE.ROUNDED_RECTANGLE,
    "ring":          MSO_SHAPE.DONUT,
    # Annotation / brackets
    "brace-left":    MSO_SHAPE.LEFT_BRACE,
    "brace-right":   MSO_SHAPE.RIGHT_BRACE,
    "brace-pair":    MSO_SHAPE.DOUBLE_BRACE,
}

# Unicode-glyph fallbacks for things python-pptx doesn't ship as a
# native shape (check, x, info, …). Rendered as text inside a textbox
# — still recolorable and scalable.
_ICON_GLYPHS = {
    "check":         "✓",   # ✓
    "x":             "✗",   # ✗
    "info":          "ⓘ",   # ⓘ
    "question":      "?",
    "exclaim":       "!",
    "dot":           "●",   # ●
    "open-dot":      "○",   # ○
    "diamond-glyph": "◆",   # ◆
    "asterisk":      "✱",   # ✱
    "skull":         "☠",   # ☠
    "atom":          "⚛",   # ⚛
    "leaf":          "⚘",   # ⚘
    "flask":         "⚚",   # ⚚
    "dna":           "\U0001F9EC",   # 🧬
    "microscope":    "\U0001F52C",   # 🔬
    "chart":         "\U0001F4C8",   # 📈
    "lock":          "\U0001F512",   # 🔒
}


def icon(slide, name: str, *, left, top, size, palette,
         color=None, fonts=None):
    """Place a named icon at `(left, top)` with `size × size` bounding
    box (todo 004 §C). The vocabulary maps semantic names to MSO_SHAPE
    built-ins (preferred — native, recolorable, editable, scalable) or
    a Unicode-glyph fallback rendered in a textbox.

    color: defaults to `palette["accent"]`. Pass an RGBColor to override.

    Available shape names: see _ICON_SHAPES keys (`arrow-*`, `chevron`,
    `lightning`, `burst`, `star`, `heart`, `plus`, `minus`, `warning`,
    `database`, `cloud`, `document`, `decision`, `molecule`, `stop`,
    `input`, `gear`, `sun`, `moon`, `circle`, `square`, `rounded`,
    `ring`, `brace-*`).

    Glyph fallback names (textbox-rendered): `check`, `x`, `info`,
    `question`, `exclaim`, `dot`, `open-dot`, `diamond-glyph`,
    `asterisk`, `skull`, `atom`, `leaf`, `flask`, `dna`, `microscope`,
    `chart`, `lock`.

    Raises ValueError for unknown names.
    """
    fill = color if color is not None else palette["accent"]
    shape_const = _ICON_SHAPES.get(name)
    if shape_const is not None:
        sh = slide.shapes.add_shape(shape_const, left, top, size, size)
        sh.line.fill.background()
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.shadow.inherit = False
        return sh
    glyph = _ICON_GLYPHS.get(name)
    if glyph is None:
        raise ValueError(
            f"unknown icon name: {name!r}. "
            f"See _ICON_SHAPES / _ICON_GLYPHS in slide_render_helpers.py "
            "for the full vocabulary."
        )
    tb = slide.shapes.add_textbox(left, top, size, size)
    tf = tb.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = glyph
    # Size text to ~75% of the box (Pt is 1/72 inch; EMU = 12700 Pt)
    pt_size = max(8, int(size / 12700 * 0.75))
    run.font.size = Pt(pt_size)
    run.font.color.rgb = fill
    if fonts and fonts.get("display"):
        run.font.name = fonts["display"]
    return tb


def icon_names() -> list[str]:
    """Sorted list of all icon names `h.icon()` accepts. Useful for the
    agent to consult before picking one."""
    return sorted(list(_ICON_SHAPES) + list(_ICON_GLYPHS))


# ─── text autofit (todo follow-up — box overlap fix) ─────────────────────


def _char_width_factor(cp: int) -> float:
    """Approximate width of a glyph at font_pt = 1.0, by codepoint
    range. Hangul / CJK / Kana count as ~1.0 × font_pt (wide); ASCII
    counts as ~0.55 × font_pt; punctuation / spaces around 0.3–0.4.
    LibreOffice / PowerPoint render at slightly different widths, so
    this is a conservative estimate used for autofit only."""
    if (0x1100 <= cp <= 0x11FF or 0x3130 <= cp <= 0x318F
            or 0xAC00 <= cp <= 0xD7AF):
        return 1.0   # Hangul Jamo / Compatibility Jamo / Syllables
    if (0x4E00 <= cp <= 0x9FFF or 0x3040 <= cp <= 0x309F
            or 0x30A0 <= cp <= 0x30FF):
        return 1.0   # CJK Unified / Hiragana / Katakana
    if cp >= 0x2000 and cp < 0x3000:
        return 0.4   # general punctuation block
    if cp < 128:
        if cp == 32:
            return 0.3
        ch = chr(cp)
        if ch.isalnum() or ch in "._/-+":
            return 0.55
        return 0.4
    return 0.6


def estimate_text_width_pt(text: str, font_pt: int) -> float:
    """Rough width-in-pt for `text` rendered at `font_pt`. Sums per-
    character width factors; Korean-aware so Korean-heavy text doesn't
    fool the autofit into thinking it'll fit."""
    return sum(_char_width_factor(ord(c)) * font_pt for c in text)


def measure_text_height_pt(text: str, *, max_width_emu, font_pt: int,
                            line_spacing: float = 1.22) -> float:
    """Estimated rendered height (in pt) of `text` wrapped to
    `max_width_emu` at `font_pt`. Used by content-driven sizing so a
    card or box can shrink to its content instead of stretching to fill
    a caller-provided rectangle. Korean-aware via `estimate_text_width_pt`.
    Empty / whitespace text returns 0.
    """
    if not text or not text.strip():
        return 0.0
    width_pt = max(1, max_width_emu / 12700)
    lines_raw = [ln for ln in text.splitlines() if ln.strip()]
    if not lines_raw:
        return 0.0
    total_lines = 0
    for ln in lines_raw:
        w = estimate_text_width_pt(ln, font_pt)
        wraps = max(1, int((w + width_pt - 1) // width_pt))
        total_lines += wraps
    return total_lines * font_pt * line_spacing


def autofit_pt(text: str, *, max_width_emu, max_height_emu, start_pt: int,
                line_spacing: float = 1.22, min_pt: int = 10) -> int:
    """Largest pt in [min_pt, start_pt] at which `text` fits inside an
    EMU-sized box, conservatively wrapping per estimated character
    width. Falls back to `min_pt` if even that doesn't fit (caller's
    job to handle over-min overflow).

    The autofit happens BEFORE python-pptx runs, so the rendered PNG
    (via soffice — which doesn't fully honor TEXT_TO_SHAPE auto-shrink)
    matches what the user will see. PowerPoint's own renderer will
    still auto-shrink further if needed.
    """
    if not text or not text.strip():
        return start_pt
    width_pt = max(1, max_width_emu / 12700)
    height_pt = max(1, max_height_emu / 12700)
    lines_raw = [ln for ln in text.splitlines() if ln.strip()]
    if not lines_raw:
        return start_pt

    def line_count_at(pt: int) -> int:
        total = 0
        for ln in lines_raw:
            w = estimate_text_width_pt(ln, pt)
            # Ceiling of w / width_pt, min 1 line
            wraps = max(1, int((w + width_pt - 1) // width_pt))
            total += wraps
        return total

    pt = start_pt
    while pt > min_pt:
        if line_count_at(pt) * pt * line_spacing <= height_pt:
            return pt
        pt -= 1
    return min_pt





# 8pt vertical rhythm — every vertical gap a snippet places should be an
# integer multiple of this (todo 004 §D). Pulled out as a constant so
# snippets can write `Pt(8 * 2)` rather than magic numbers.
SPACING_UNIT_PT = 8


class Grid:
    """12-column / N-row design grid for `code` slides (todo 004 §D).

    A `Grid` is pure coordinate math — it does NOT add any shapes. Use
    `grid.cell(col=, span=, row=, row_span=)` to get a `(left, top,
    width, height)` EMU tuple you can pass to the layout helpers
    (`h.card`, `h.bullet_list`, `h.image_*`, etc.).

    The grid leaves room at the top for `h.title_block` (default
    margin_top=Inches(1.8)) so titles and grid content don't fight.
    The 8pt vertical rhythm is enforced as the row gap.

    Replaces ad-hoc `Inches(0.7)` constants scattered across snippets
    with a single shared geometry — alignment becomes default, not
    something the agent has to think about.
    """

    def __init__(self, *, sw, sh, cols: int = 12, rows: int = 6,
                 gutter=None, margin_x=None, margin_top=None,
                 margin_bot=None, row_gap=None, row_h=None):
        if cols < 1:
            raise ValueError(f"cols must be >= 1, got {cols}")
        if rows < 1:
            raise ValueError(f"rows must be >= 1, got {rows}")
        self.sw = sw
        self.sh = sh
        self.cols = cols
        self.rows = rows
        self.gutter = gutter if gutter is not None else Pt(SPACING_UNIT_PT)
        self.margin_x = margin_x if margin_x is not None else Inches(0.6)
        self.margin_top = margin_top if margin_top is not None else Inches(1.8)
        self.margin_bot = margin_bot if margin_bot is not None else Inches(0.6)
        self.row_gap = row_gap if row_gap is not None else Pt(SPACING_UNIT_PT)
        # Column width: usable width minus gutters, split N ways
        usable_w = sw - 2 * self.margin_x - self.gutter * (cols - 1)
        self.col_w = int(usable_w // cols)
        # Row height: usable height minus row gaps, split N ways
        usable_h = sh - self.margin_top - self.margin_bot
        if row_h is not None:
            self.row_h = row_h
        else:
            self.row_h = int(
                (usable_h - self.row_gap * (rows - 1)) // rows
            )

    def cell(self, col: int, span: int = 1, row: int = 1,
             row_span: int = 1):
        """Return `(left, top, width, height)` in EMU for a grid cell.

        col / row are 1-indexed (1..cols, 1..rows). `span` extends the
        cell rightward; `row_span` extends it downward. Spanning
        absorbs the gutters / row gaps between covered cells (so a
        col_span=3 block is *wider* than 3 stand-alone cells).
        """
        if not 1 <= col <= self.cols:
            raise ValueError(f"col must be 1..{self.cols}, got {col}")
        if span < 1 or col + span - 1 > self.cols:
            raise ValueError(
                f"span {span} from col {col} exceeds {self.cols} cols",
            )
        if not 1 <= row <= self.rows:
            raise ValueError(f"row must be 1..{self.rows}, got {row}")
        if row_span < 1 or row + row_span - 1 > self.rows:
            raise ValueError(
                f"row_span {row_span} from row {row} exceeds {self.rows} rows",
            )
        left = self.margin_x + (col - 1) * (self.col_w + self.gutter)
        width = self.col_w * span + self.gutter * (span - 1)
        top = self.margin_top + (row - 1) * (self.row_h + self.row_gap)
        height = self.row_h * row_span + self.row_gap * (row_span - 1)
        return Cell(left, top, width, height)

    def row(self, row: int = 1, row_span: int = 1):
        """Shorthand for `cell(1, span=cols, row=row, row_span=row_span)`
        — the whole row, useful for headlines or pull-quotes that span
        all columns."""
        return self.cell(1, span=self.cols, row=row, row_span=row_span)


def grid(*, sw, sh, **kwargs) -> Grid:
    """Build a `Grid` sized to the current slide. The `slide` argument
    isn't needed — the grid is pure coordinate math. Pass `cols`,
    `rows`, `gutter`, `margin_x`, `margin_top`, `margin_bot`,
    `row_gap`, `row_h` to override defaults (12 / 6 / Pt(8) / Inches
    (0.6) / Inches(1.8) / Inches(0.6) / Pt(8) / computed).

    Typical usage in a `code` snippet:

        g = h.grid(sw=sw, sh=sh)
        left, top, w_, h_ = g.cell(col=1, span=6, row=1, row_span=3)
        h.bullet_list(slide, items, left=left, top=top,
                      width=w_, height=h_, ...)
    """
    return Grid(sw=sw, sh=sh, **kwargs)


def _autoshrink(tf) -> None:
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
    except Exception:
        pass


def accent_stripe(slide, *, palette, sw, height_in: float = 0.16):
    """Top accent stripe — the deck's signature horizontal bar."""
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(height_in),
    )
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = palette["accent"]
    stripe.shadow.inherit = False
    return stripe


def title_block(slide, text: str, *, palette, fonts, type_scale, sw, sh,
                left=None, top=None, width=None, height=None,
                cover: bool = False, accent_rule: bool = True):
    """A themed title textbox.

    cover=False (default): top-anchored content-slide title at
    Inches(0.7), Inches(0.45). A short accent rule sits beneath when
    accent_rule=True.

    cover=True: vertically + horizontally centered, no rule, larger
    type. Use for `role='title'` slides.
    """
    if cover:
        left = Inches(1.0) if left is None else left
        top = Inches(0.5) if top is None else top
        width = sw - Inches(2.0) if width is None else width
        height = sh - Inches(1.0) if height is None else height
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _autoshrink(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1
        run = p.add_run()
        run.text = text or ""
        run.font.size = Pt(type_scale.get("cover_title", 40))
        run.font.bold = True
        run.font.color.rgb = palette["foreground"]
        if fonts.get("display"):
            run.font.name = fonts["display"]
        return box

    left = Inches(0.7) if left is None else left
    top = Inches(0.45) if top is None else top
    width = sw - Inches(1.4) if width is None else width
    height = Inches(1.0) if height is None else height
    title_pt_start = type_scale.get("title", 32)
    # Autofit (todo 014 title-fix): without this, long titles silently
    # wrap to 2 lines under soffice (visible in p.13 "Harness 네 가지
    # primitive — 대화가 소프트웨어가 되는 / 방식"). min_pt=20 keeps
    # titles readable; if even 20pt overflows the box width, we accept
    # the wrap.
    title_pt = autofit_pt(
        text or "", max_width_emu=width, max_height_emu=height,
        start_pt=title_pt_start, line_spacing=1.05, min_pt=20,
    )
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    _autoshrink(tf)
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = palette["foreground"]
    if fonts.get("display"):
        run.font.name = fonts["display"]
    if accent_rule:
        # Position the rule with enough vertical gap that it reads as a
        # *floating accent below the title* rather than an underline of
        # the title's first words. The previous tight gap (Pt(8) below
        # measured height) made short titles look underlined; restore
        # the old "loose accent" position for 1-line titles, and only
        # push the rule lower when the title actually wraps. (todo 014
        # title-fix v2.)
        title_natural_pt = measure_text_height_pt(
            text or "", max_width_emu=width,
            font_pt=title_pt, line_spacing=1.05,
        )
        rule_top = max(
            top + height + Pt(2),                       # original floor
            top + Pt(int(title_natural_pt) + 20),       # wrap fallback
        )
        # Width: extend to the title's visible width so the rule never
        # appears to underline only a prefix of a long title. Falls
        # back to Inches(2.2) for short titles. Capped at the title
        # box width.
        title_visible_pt = estimate_text_width_pt(text or "", title_pt)
        rule_w = max(Inches(2.2), Pt(int(title_visible_pt) + 4))
        rule_w = min(rule_w, width)
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left + Inches(0.02),
            rule_top, rule_w, Pt(4),
        )
        rule.line.fill.background()
        rule.fill.solid()
        rule.fill.fore_color.rgb = palette["accent"]
        rule.shadow.inherit = False
    return box


def deck_chrome(slide, *, palette, fonts, type_scale, sw, sh,
                eyebrow: str = "",
                page_number: int = None, total: int = None,
                footer: str = ""):
    """Deck-level chrome (todo 009 B): one call per slide adds the
    semantic eyebrow label upper-left + footer label bottom-left +
    page-number bottom-right. Use AFTER `h.accent_stripe` +
    `h.title_block` so the eyebrow sits above the title and the
    footer/page sit below the body.

    Skip the args you don't need — `eyebrow=""` leaves the eyebrow
    off, `page_number=None` skips the page number, `footer=""` skips
    the footer.
    """
    cap_pt = type_scale.get("caption", 12)
    label_pt = type_scale.get("label_tag", cap_pt)
    fg_muted = palette.get("muted")
    accent = palette["accent"]

    # Eyebrow: above the title (todo 009 — semantic navigation label
    # like "HOW · 추진 방법" or "WHY · 다른 제안과의 결정적 차이").
    if eyebrow:
        eyb = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.22),
            sw - Inches(1.4), Pt(label_pt * 1.6),
        )
        tf = eyb.text_frame
        tf.word_wrap = False
        _autoshrink(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = eyebrow.upper()
        run.font.size = Pt(label_pt)
        run.font.bold = True
        run.font.color.rgb = accent
        if fonts.get("body"):
            run.font.name = fonts["body"]

    # Footer: bottom-left, muted body type (deck title / project name).
    if footer:
        ft = slide.shapes.add_textbox(
            Inches(0.5), sh - Inches(0.35),
            sw - Inches(2.5), Pt(cap_pt * 1.8),
        )
        tf = ft.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = footer
        run.font.size = Pt(cap_pt)
        if fg_muted is not None:
            run.font.color.rgb = fg_muted
        if fonts.get("body"):
            run.font.name = fonts["body"]

    # Page number: bottom-right (e.g. "5 / 13").
    if page_number is not None:
        pgs = (f"{page_number} / {total}" if total is not None
               else f"{page_number}")
        pn = slide.shapes.add_textbox(
            sw - Inches(1.5), sh - Inches(0.35),
            Inches(1.0), Pt(cap_pt * 1.8),
        )
        tf = pn.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = pgs
        run.font.size = Pt(cap_pt)
        if fg_muted is not None:
            run.font.color.rgb = fg_muted
        if fonts.get("body"):
            run.font.name = fonts["body"]


def table(slide, *, headers: list[str], rows: list[list[str]],
          left, top, width, height,
          palette, fonts, type_scale,
          first_col_emphasis: bool = False):
    """Native python-pptx `MSO_SHAPE_TYPE.TABLE` — a real, editable
    PowerPoint table (todo 009 C). For tabular data (personnel,
    equipment, timelines, parameter sweeps) the native table beats
    re-implementing rows with card_grid: cells nudge + resize properly
    in the editor, headers + body get distinct styling, and the
    structure is grep-able.

    Headers render in display font + accent color. Body cells in body
    font + foreground. `first_col_emphasis=True` styles the first
    column in display font (useful for personnel / timeline labels).

    Returns the table shape so the caller can post-tweak (column
    widths, individual cell merges, etc.) if needed.
    """
    cols = len(headers)
    n_rows = len(rows) + 1   # +1 for header row
    shape = slide.shapes.add_table(
        n_rows, cols, left, top, width, height,
    )
    tbl = shape.table

    head_pt = max(11, type_scale.get("label_tag", 12))
    body_pt = max(10, type_scale.get("body_small",
                                       type_scale.get("body", 20) - 4))
    fg = palette["foreground"]
    accent = palette["accent"]
    # Header row
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        tf = cell.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(header)
        run.font.size = Pt(head_pt)
        run.font.bold = True
        run.font.color.rgb = accent
        if fonts.get("display"):
            run.font.name = fonts["display"]
    # Body rows
    for r, row in enumerate(rows, start=1):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            value = row[c] if c < len(row) else ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(value)
            run.font.size = Pt(body_pt)
            run.font.color.rgb = fg
            font_name = (
                fonts.get("display") if (first_col_emphasis and c == 0)
                else fonts.get("body")
            )
            if font_name:
                run.font.name = font_name
            if first_col_emphasis and c == 0:
                run.font.bold = True
    return shape


def text(slide, content: str, *, left, top, width, height,
         palette, size_pt: int = 20, color=None, font_name=None,
         bold: bool = False, italic: bool = False, align=None,
         anchor=None, line_spacing: float = 1.22,
         autofit: bool = True, min_pt: int = 12, fonts=None):
    """One-call themed textbox (todo 007 §D — DX helper). Drops the
    5-line `add_textbox + text_frame + paragraph + run + font` ceremony
    to a single call. Color defaults to `palette["foreground"]`; pass
    `color=palette["muted"]` for a caption look. Autofit-shrinks to
    fit the box by default (Korean-aware).

    Returns the textbox shape for further tweaking if needed.
    """
    fill = color if color is not None else palette.get("foreground")
    actual_pt = (
        autofit_pt(content or "", max_width_emu=width,
                   max_height_emu=height, start_pt=size_pt,
                   line_spacing=line_spacing, min_pt=min_pt)
        if autofit else size_pt
    )
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    _autoshrink(tf)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = content or ""
    run.font.size = Pt(actual_pt)
    run.font.bold = bold
    run.font.italic = italic
    if fill is not None:
        run.font.color.rgb = fill
    name = font_name or (fonts or {}).get("body")
    if name:
        run.font.name = name
    return tb


def vstack(slide, lines, *, left, top, width, palette,
           fonts=None, gap_pt: int = 4):
    """Stack text lines vertically with auto-measured heights (todo 014
    vstack-fix). Eliminates the "two textboxes at the same (left, top)
    overlap" bug pattern that breaks bespoke compositions — pass a list
    of styled-line dicts and the helper places each one beneath the
    previous, computing per-line height via `measure_text_height_pt`.

    Each `lines` item is a dict with keys:
      text        (str, required; empty/missing → renders as a gap)
      size_pt     (int, default 14)
      color       (RGBColor or palette key like "muted"/"accent", default
                   palette["foreground"])
      bold        (bool, default False)
      italic      (bool, default False)
      font_name   (str or "display"/"body", default body)
      align       (PP_ALIGN, default LEFT)
      line_spacing(float, default 1.22)
      pad_top_pt  (int, extra gap before this line in pt, default 0)

    Returns the y of the bottom of the stack — pass `top=returned_y +
    Pt(...)` to chain the next region. Use this whenever you would
    otherwise call `h.text(...)` twice with hand-computed `top` values
    for layouts like "number eyebrow + label below" or "title +
    subtitle + body".
    """
    y = top
    fill_default = palette.get("foreground") if palette else None
    fonts = fonts or {}
    for item in lines:
        if not isinstance(item, dict):
            item = {"text": str(item)}
        pad_top = item.get("pad_top_pt", 0)
        if pad_top:
            y += Pt(pad_top)
        body = item.get("text", "") or ""
        if not body.strip():
            # Empty/missing text → emit a gap proportional to size.
            size_pt = item.get("size_pt", 14)
            y += Pt(int(size_pt * 0.8) + gap_pt)
            continue
        size_pt = item.get("size_pt", 14)
        item_line_spacing = item.get("line_spacing", 1.22)
        color = item.get("color")
        if isinstance(color, str) and palette:
            color = palette.get(color, color)
        if color is None:
            color = fill_default
        font_name = item.get("font_name")
        if font_name in ("display", "body"):
            font_name = fonts.get(font_name)
        natural_pt = measure_text_height_pt(
            body, max_width_emu=width, font_pt=size_pt,
            line_spacing=item_line_spacing,
        )
        line_h = Pt(int(natural_pt) + 4)
        text(slide, body, left=left, top=y, width=width, height=line_h,
             palette=palette, size_pt=size_pt, color=color,
             font_name=font_name, bold=item.get("bold", False),
             italic=item.get("italic", False),
             align=item.get("align"),
             line_spacing=item_line_spacing,
             fonts=fonts)
        y += line_h + Pt(gap_pt)
    return y


def callout(slide, *, left, top, width, fill,
            palette, items=None,
            headline: str = "", body: str = "",
            fonts=None, type_scale=None,
            pad_pt: int = 14, gap_pt: int = 6,
            min_height=None, text_color=None,
            border_color=None, border_pt: float = 0):
    """Filled callout box whose background rectangle auto-sizes to fit
    its content (todo 014 callout-fix). Fixes the failure where bespoke
    code drew a fixed-size rect first, then placed text inside that
    overflowed the bottom (visible in v3 p.15 right "Next step" dark
    box where the closing sentence got clipped).

    Two ways to pass content:
    - `items=[{text, size_pt?, color?, bold?, italic?, font_name?,
       align?, line_spacing?, pad_top_pt?}, …]` — same dict shape as
       `h.vstack`. Use this when you have ≥3 stacked elements or
       custom styling per line.
    - `headline=` and/or `body=` — convenience for the common
      "small bold headline + body paragraph" callout. Resolves to:
        items=[
          {text: headline, size_pt: 13, bold: True, color: text_color},
          {text: body,     size_pt: 14,            color: text_color},
        ]

    `fill` is the background RGBColor (e.g., `palette["foreground"]`
    for a dark callout; `palette["surface"]` for a light one).
    `text_color` defaults to `palette["surface"]` for dark fills and
    `palette["foreground"]` for light fills (auto-detected by
    relative luminance).

    Returns `{"box": rect, "height_used": emu}`.
    """
    if items is None:
        items = []
        if headline:
            items.append({
                "text": headline, "size_pt": 13, "bold": True,
                "color": text_color,
            })
        if body:
            items.append({
                "text": body, "size_pt": 14,
                "color": text_color,
            })

    pad = Pt(pad_pt)
    inner_w = width - 2 * pad
    line_spacing_default = (
        type_scale.get("line_spacing", 1.22) if type_scale else 1.22
    )

    # Pre-measure total content height in pt, mirroring vstack's
    # per-line height math so the box exactly contains the stack.
    content_pt = 0
    last_i = len(items) - 1
    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"text": str(item)}
        pad_top = item.get("pad_top_pt", 0)
        if pad_top:
            content_pt += pad_top
        body_text = item.get("text", "") or ""
        size_pt = item.get("size_pt", 14)
        item_ls = item.get("line_spacing", line_spacing_default)
        if not body_text.strip():
            content_pt += int(size_pt * 0.8) + gap_pt
            continue
        natural_pt = measure_text_height_pt(
            body_text, max_width_emu=inner_w, font_pt=size_pt,
            line_spacing=item_ls,
        )
        content_pt += int(natural_pt) + 4
        if i < last_i:
            content_pt += gap_pt

    needed_h = Pt(content_pt) + 2 * pad
    box_h = max(min_height, needed_h) if min_height else needed_h

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, box_h,
    )
    if border_pt > 0 and border_color is not None:
        bg.line.color.rgb = border_color
        bg.line.width = Pt(border_pt)
    else:
        bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.shadow.inherit = False

    # If text_color wasn't set, pick foreground or surface based on
    # the fill's luminance (so the body remains readable on either
    # dark or light callout backgrounds).
    if text_color is None and palette:
        r, g, b = fill[0], fill[1], fill[2]
        luma = (0.299 * r + 0.587 * g + 0.114 * b)
        text_color = (
            palette.get("surface", palette.get("background"))
            if luma < 128 else palette.get("foreground")
        )
        for item in items:
            if item.get("color") is None:
                item["color"] = text_color

    vstack(slide, items, left=left + pad, top=top + pad,
           width=inner_w, palette=palette, fonts=fonts, gap_pt=gap_pt)

    return {"box": bg, "height_used": box_h}


def bullet_list(slide, items, *, palette, fonts, type_scale,
                left, top, width, height, bullet: str = "•",
                pack: bool = True):
    """A vertical bulleted list. `items` is a list of strings; each
    becomes one paragraph at the deck's body type (autofit-shrunk if
    the items don't fit at the start size). Pass bullet="" for a
    plain paragraph list.

    When `pack=True` (default, todo 014 D-fix v2), the textbox shrinks
    to natural content height — caller reads `box.height` after the
    call for the actual emitted height (so a wrapper rectangle can be
    sized to match). Pass `pack=False` for fixed-rectangle behavior.
    """
    body_pt = type_scale.get("body", 20)
    line_spacing = type_scale.get("line_spacing", 1.22)
    rendered = [(f"{bullet} {it}" if bullet else str(it)) for it in items]
    actual_pt = autofit_pt(
        "\n".join(rendered),
        max_width_emu=width, max_height_emu=height,
        start_pt=body_pt, line_spacing=line_spacing,
        min_pt=max(12, body_pt - 6),
    )
    if pack and rendered:
        natural_pt = sum(
            measure_text_height_pt(line, max_width_emu=width,
                                   font_pt=actual_pt,
                                   line_spacing=line_spacing)
            for line in rendered
        ) + (len(rendered) - 1) * 4  # space_after between paragraphs (pt)
        natural_h = Pt(int(natural_pt) + 4)
        used_h = min(height, natural_h) if natural_h < height else height
    else:
        used_h = height
    box = slide.shapes.add_textbox(left, top, width, used_h)
    tf = box.text_frame
    tf.word_wrap = True
    _autoshrink(tf)
    first = True
    for line in rendered:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.line_spacing = line_spacing
        para.space_after = Pt(4)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(actual_pt)
        run.font.color.rgb = palette["foreground"]
        if fonts.get("body"):
            run.font.name = fonts["body"]
    return box


def card(slide, *, left, top, width, height, title: str, body: str,
         palette, fonts, type_scale,
         accent_top: bool = True, accent_height_pt: int = 4,
         pack: bool = True):
    """A single titled card: rectangle + optional top accent stripe +
    title (bold) + body.

    When `pack=True` (default, todo 014 D-fix), the card SHRINKS to fit
    its content — if `pad + title_h + body_natural_h + pad < height`,
    the card rectangle's height becomes content-driven. This kills the
    "title at top, body at top, big empty space at the bottom" failure
    where a tall card holds short content. Pass `pack=False` to keep a
    fixed-height card (e.g. uniform-grid alignment); `card_grid` uses
    this internally after measuring max content height across the grid.

    Returns dict(card, title_box, body_box, height_used) — `height_used`
    is the actual emitted height (== `height` when pack=False or
    content overflows; less when shrunk).
    """
    bg = palette.get("surface", palette["background"])
    border = palette["accent"]
    title_pt = max(14, type_scale.get("head", 26) - 4)
    body_pt = max(12, type_scale.get("body", 20) - 4)
    line_spacing = type_scale.get("line_spacing", 1.22)
    pad = Pt(10)
    inner_w = width - 2 * pad
    title_h_natural_pt = (
        measure_text_height_pt(title or "", max_width_emu=inner_w,
                               font_pt=title_pt, line_spacing=1.1)
    )
    title_h = Pt(max(int(title_pt * 1.6), int(title_h_natural_pt) + 4))
    body_h_natural_pt = (
        measure_text_height_pt(body or "", max_width_emu=inner_w,
                               font_pt=body_pt, line_spacing=line_spacing)
    )
    body_natural_h = Pt(int(body_h_natural_pt) + 4) if body else Pt(0)

    natural_card_h = pad + title_h + body_natural_h + pad
    card_h = (
        min(height, natural_card_h) if pack and natural_card_h < height
        else height
    )
    body_h = card_h - pad - title_h - pad

    card_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, card_h,
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = bg
    card_shape.line.color.rgb = border
    card_shape.line.width = Pt(0.75)
    card_shape.shadow.inherit = False

    if accent_top:
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Pt(accent_height_pt),
        )
        stripe.line.fill.background()
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = palette["accent"]
        stripe.shadow.inherit = False

    tb = slide.shapes.add_textbox(
        left + pad, top + pad, inner_w, title_h,
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title or ""
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = palette["foreground"]
    if fonts.get("display"):
        run.font.name = fonts["display"]

    bb = slide.shapes.add_textbox(
        left + pad, top + pad + title_h,
        inner_w, body_h,
    )
    bf = bb.text_frame
    bf.word_wrap = True
    _autoshrink(bf)
    if body:
        # Body autofit (todo 014 B-fix): soffice doesn't honor
        # MSO_AUTO_SIZE.TEXT_TO_SHAPE during PDF preview, so an
        # over-tall body would wrap silently past the card's bottom
        # edge. Run the Korean-aware autofit so the chosen pt actually
        # fits inside body_h before render.
        body_actual_pt = autofit_pt(
            body, max_width_emu=inner_w, max_height_emu=body_h,
            start_pt=body_pt, line_spacing=line_spacing,
            min_pt=max(12, body_pt - 4),
        )
        p = bf.paragraphs[0]
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = body
        run.font.size = Pt(body_actual_pt)
        run.font.color.rgb = palette["foreground"]
        if fonts.get("body"):
            run.font.name = fonts["body"]
    return {"card": card_shape, "title_box": tb, "body_box": bb,
            "height_used": card_h}


def card_grid(slide, items, *, left, top, width, height,
              palette, fonts, type_scale,
              cols: int = 2, gap_pt: int = 12):
    """Lay out `items` (list of dicts with `title` + `body`, or the
    canonical `tag` + `body` — both accepted, todo 007 axis 2) as a
    grid of `card()`s with `cols` columns. Rows are derived from
    `len(items)`.

    Cells share a uniform height for grid alignment. The shared height
    is min(requested_cell_h, max(natural card heights)) — i.e. the
    grid tightens to its content when there's room (todo 014 D-fix),
    eliminating the under-fill where short content sits at the top of
    a tall card. Caller can rely on the returned cards' `height_used`
    + the grid's `top` + emitted gaps to position downstream shapes.
    """
    n = len(items)
    if n == 0:
        return []
    rows = (n + cols - 1) // cols
    gap = Pt(gap_pt)
    cell_w = (width - gap * (cols - 1)) // cols
    requested_cell_h = (height - gap * (rows - 1)) // rows

    title_pt = max(14, type_scale.get("head", 26) - 4)
    body_pt = max(12, type_scale.get("body", 20) - 4)
    line_spacing = type_scale.get("line_spacing", 1.22)
    pad = Pt(10)
    inner_w = cell_w - 2 * pad

    def _resolve_text(item, primary, fallback):
        return ((item.get(primary) if isinstance(item, dict) else None)
                or (item.get(fallback) if isinstance(item, dict) else None)
                or "")

    natural_heights = []
    for item in items:
        title = _resolve_text(item, "title", "tag")
        body = _resolve_text(item, "body", "note")
        title_h_pt = measure_text_height_pt(
            title, max_width_emu=inner_w, font_pt=title_pt, line_spacing=1.1
        )
        title_h = max(int(title_pt * 1.6), int(title_h_pt) + 4)
        body_h = (
            int(measure_text_height_pt(
                body, max_width_emu=inner_w,
                font_pt=body_pt, line_spacing=line_spacing,
            )) + 4
        ) if body else 0
        natural_heights.append(Pt(title_h + body_h) + 2 * pad)

    cell_h = min(requested_cell_h, max(natural_heights)) if natural_heights \
        else requested_cell_h

    out = []
    for i, item in enumerate(items):
        r, c = i // cols, i % cols
        cl = left + c * (cell_w + gap)
        ct = top + r * (cell_h + gap)
        title = _resolve_text(item, "title", "tag")
        body = _resolve_text(item, "body", "note")
        out.append(card(
            slide, left=cl, top=ct, width=cell_w, height=cell_h,
            title=title, body=body,
            palette=palette, fonts=fonts, type_scale=type_scale,
            pack=False,
        ))
    return out


def pull_quote(slide, text: str, *, palette, fonts, type_scale,
               left, top, width, height,
               bar_pt: int = 4, gap_pt: int = 10,
               pack: bool = True):
    """An emphasis block: vertical accent bar + italic text. Use for a
    punchline / take-home line you want the audience to focus on.

    Multi-line input: each non-empty line becomes one paragraph.
    **A blank line in the input adds an explicit half-line gap after
    the previous paragraph** (todo 014 pull_quote polish) — use this
    when you want a visual break between two parts of the quote
    without using two separate `pull_quote` calls.

    When `pack=True` (default, todo 014 D-fix v2), the accent bar and
    textbox shrink to natural content height when text is shorter than
    the requested rectangle — otherwise the bar floats above empty
    space below the quote. Return gains `height_used` for downstream
    layout. Pass `pack=False` for fixed-rectangle behavior.
    """
    body_pt = type_scale.get("body", 20)
    line_spacing = type_scale.get("line_spacing", 1.22)
    quote_pt = max(body_pt, int(body_pt * 1.1))
    text_left = left + Pt(bar_pt) + Pt(gap_pt)
    text_width = width - Pt(bar_pt) - Pt(gap_pt)
    half_gap_pt = int(quote_pt * 0.6)
    # Measure natural height: account for blank-line gaps in the input
    # so packing reserves the right space for them.
    raw_lines = (text or "").splitlines()
    non_empty = [ln for ln in raw_lines if ln.strip()]
    blank_count = max(0, sum(1 for ln in raw_lines if not ln.strip()) - 0)
    natural_pt_text = measure_text_height_pt(
        "\n".join(non_empty), max_width_emu=text_width, font_pt=quote_pt,
        line_spacing=line_spacing,
    )
    natural_pt_total = natural_pt_text + blank_count * half_gap_pt
    natural_h = Pt(int(natural_pt_total) + 4) if non_empty else Pt(0)
    used_h = (
        min(height, natural_h) if pack and natural_h < height and natural_h > 0
        else height
    )
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(bar_pt), used_h,
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["accent"]
    bar.shadow.inherit = False
    tb = slide.shapes.add_textbox(text_left, top, text_width, used_h)
    tf = tb.text_frame
    tf.word_wrap = True
    _autoshrink(tf)
    first = True
    pending_gap = False
    last_para = None
    for line in raw_lines:
        if not line.strip():
            # Blank line in input → add half-line gap after the
            # previous paragraph (rather than silently skipping).
            if last_para is not None:
                last_para.space_after = Pt(half_gap_pt)
                pending_gap = False  # already applied to last_para
            else:
                pending_gap = True  # apply to the next paragraph as space_before
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        if pending_gap:
            p.space_before = Pt(half_gap_pt)
            pending_gap = False
        run = p.add_run()
        run.text = line.strip()
        run.font.size = Pt(quote_pt)
        run.font.italic = True
        run.font.color.rgb = palette["foreground"]
        if fonts.get("body"):
            run.font.name = fonts["body"]
        last_para = p
    return {"bar": bar, "text_box": tb, "height_used": used_h}
